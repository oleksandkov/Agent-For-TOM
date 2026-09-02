#!/usr/bin/env python3
"""
Build Your Own Claude Code — a minimal AI coding agent.

Based on the architecture popularized by the "Build Your Own Claude Code" video
(Devtools Tech): a while-loop that calls an LLM, dispatches tool calls, and
feeds results back until the model returns plain text.

Features:
  - Agent loop with tool use
  - Tools: read_file, write_file, edit_file, list_files, run_command, search_code
  - Project context re-injection (AGENTS.md / AGENT.md)
  - Three-layer memory system (~/.tomas/memory/)
  - Risk-tiered permission system (low / medium / high)
  - Auto-compaction when the context window fills up

Env vars:
  ANTHROPIC_API_KEY   - required API key
  ANTHROPIC_BASE_URL  - optional, point to any Anthropic-compatible endpoint
  AGENT_MODEL         - optional, model name (e.g. claude-sonnet-5)
  AGENT_AUTO_APPROVE  - optional, "1" to auto-approve low-risk tools
"""

from __future__ import annotations

import os
import re
import sys
import time
import json
import base64
import shutil
import tempfile
import threading
import subprocess
import urllib.request
import urllib.error
import html as html_module
from pathlib import Path
import dataclasses
from typing import Any, Callable, Optional

from core import budget as core_budget
from core import context as core_context
from core import debug_log
from core import features as core_features
from core.console import CONSOLE
from core.context import ContextWindow

# ── Windows console setup ──
# The UI prints box-drawing and symbol glyphs (▌ ✧ ⚙ ⚡ ↳). On a console whose
# codepage is not UTF-8 (cp1251, cp1252, cp437 — common on non-English Windows)
# print() raises UnicodeEncodeError and kills the agent. Force UTF-8 with a
# replacement fallback, and enable VT100 so ANSI colours work on legacy conhost.
if sys.platform == "win32":
    for _stream in (sys.stdout, sys.stderr):
        try:
            _stream.reconfigure(encoding="utf-8", errors="replace")
        except Exception:
            pass
    try:
        os.system("")  # enables ANSI escape processing in legacy consoles
    except Exception:
        pass

# Playwright for browser-based fetching (JavaScript rendering).
# Whether it is installed is answered from the import system's metadata
# instead of by importing it: the package costs ~105 ms to load, and the CLI
# menus only ever need the boolean. The tools that drive a browser import it
# themselves, at the point they are about to use it.
import importlib.util

PLAYWRIGHT_AVAILABLE = importlib.util.find_spec("playwright") is not None
async_playwright = None  # bound on demand by the browser tools

# Load variables from .env into os.environ if present.
try:
    from dotenv import load_dotenv
    load_dotenv()
except ImportError:
    # python-dotenv not installed — rely on real env vars instead.
    pass

# NOTE: `anthropic` is deliberately *not* imported here. It is the single
# most expensive import in the project — 1.13 s of the CLI's 1.66 s cold
# start — and the menus never touch it: they import TOOLS, RISK_LEVELS and
# build_system_prompt from this module and go straight to drawing. It is
# imported where it is used, in `_get_client` (and already lazily in
# `core/loop.py`), so only a session that actually talks to a model pays it.

# MCP and skills support
from mcp_manager import MCPManager
from skills_manager import (build_skills_section, build_triggered_skills,
                            triggered_tool_allowlist,
                            discover_skills, cmd_skill_list, cmd_skill_run,
                            match_skills)

# Self-improving system
import self_improve

# Self-notes system
import self_notes

# Session management
from session_manager import (
    save_session, list_sessions, load_session,
    continue_session, delete_session, get_latest_session,
)

# Instructions management
import instructions_manager
from instructions_manager import (
    build_instructions_section, get_global_instructions,
)

# Learning system (facts, retrieval, reflection)
import learning

# Agent core (headless, event-emitting) + the terminal front end for it
from core import loop as core_loop
from core.loop import run_turn
from core.permissions import ApprovalStore
from core.state import AgentState
from adapters.terminal import TerminalAdapter, Thinking

# ---------------------------------------------------------------------------
# Configuration
# ---------------------------------------------------------------------------

MODEL = os.environ.get("AGENT_MODEL")
PROJECT_DIR = Path(os.environ.get("AGENT_PROJECT_DIR", os.getcwd())).resolve()

# Where /export writes. The project directory is the right default — an export
# the user cannot find has not been exported — but it is a module global rather
# than a literal so a test sweep that calls every command can redirect it, the
# same seam `session_manager.SESSION_DIR` and `self_notes.NOTES_DIR` provide.
# Without it, `test_command_surface` left a conversation dump in the repo root
# on every run.
EXPORT_DIR = PROJECT_DIR
MEMORY_DIR = Path.home() / ".tomas" / "memory"

#: The four places a remembered thing can go. Named here rather than beside
#: `route_memory` because the `save_memory` tool schema offers them as an
#: enum, and TOOLS is built long before the routing code is reached.
STORE_INSTRUCTION, STORE_RULE, STORE_FACT, STORE_NOTE = (
    "instruction", "rule", "fact", "note")
MEMORY_STORES = (STORE_INSTRUCTION, STORE_RULE, STORE_FACT, STORE_NOTE)

# Where throwaway helper scripts belong. The sandbox allows writes only under
# PROJECT_DIR, so "put it in a temp directory outside the project" — which
# BASE_PROMPT used to say — asked for something that could not work: the model
# tried ~/.tomas/tmp, was refused, tried %TEMP%, was refused again, and only
# then wrote into the repo root it was being told to keep clean. One named
# location inside the sandbox resolves the contradiction.
SCRATCH_DIR = PROJECT_DIR / "_scratch"
def _env_int(name: str, default: int, minimum: int = 1) -> int:
    """An int from the environment, ignoring anything unusable."""
    try:
        value = int(os.environ.get(name, "").strip())
    except (TypeError, ValueError):
        return default
    return value if value >= minimum else default


# Output-token ceiling per turn. Overridable because the right value depends on
# the model: a reasoning model spends most of this budget thinking before it
# writes anything, so a limit that is generous for a chat model truncates it to
# nothing. Raise it when a long document or a reasoning model gets cut off.
MAX_TOKENS = _env_int("AGENT_MAX_TOKENS", 8192, minimum=256)
# When automatic compaction fires, as a fraction of CONTEXT_WINDOW — or None
# for "never, only /compact". Read from the persisted budget settings so the
# choice survives a restart and is made in one place (the Context Budget page
# and /budget both write it); `AGENT_COMPACT_FIT` overrides it for a single
# run, the way every other budget knob here can be overridden.
def _compaction_threshold() -> Optional[float]:
    raw = os.environ.get("AGENT_COMPACT_FIT", "").strip()
    if raw:
        try:
            return core_context.fit_fraction_from_percent(int(raw))
        except ValueError:
            pass
    try:
        import core.budget as _budget
        return _budget.compaction_fit_fraction(_budget.load_settings())
    except Exception:
        return core_context.DEFAULT_FIT_FRACTION


COMPACTION_THRESHOLD: Optional[float] = _compaction_threshold()
# Cost ceiling on top of the fit rule. Past roughly this much conversation the
# cost of re-reading it every turn outweighs what the extra history is worth.
#
# It used to be a flat 120,000 combined with `min(...)`, which meant a model
# with a 1,000,000-token window compacted as though it had 120,000 — the window
# the user picked the model *for* was discarded silently, the same shape of bug
# as the old `min(MAX_TOKENS, max_output_tokens)` clamp. The limit now scales
# with the real window (see core.context.default_cost_limit); this reproduces
# the old number exactly for every window up to 480,000.
#
#   unset  → derived from the window
#   0      → cost rule off, only the fit rule applies
#   N      → exactly N, as before
COMPACTION_COST_LIMIT: Optional[int] = None
_compact_at_raw = os.environ.get("AGENT_COMPACT_AT", "").strip()
if _compact_at_raw:
    try:
        _compact_at = int(_compact_at_raw)
        COMPACTION_COST_LIMIT = max(0, _compact_at)
    except ValueError:
        pass
#: Retained because it is the floor the derived limit never goes below.
COMPACTION_CEILING = core_context.COST_FLOOR
DEFAULT_CONTEXT_WINDOW = 200_000  # fallback if API doesn't report context window (standard Claude tier)
CONTEXT_WINDOW = DEFAULT_CONTEXT_WINDOW  # will be updated dynamically at startup

# ── Characters per token, by what is being counted ──
# These decide when compaction fires, so being wrong in either direction costs
# something real: too low and a long session overflows the window before the
# summariser runs, too high and it compacts history that still fitted.
#
# Prose runs about 4 chars/token. JSON runs far denser — punctuation, quotes
# and short keys all tokenise separately — measured at ~3.5 across the 64 tools
# of three real MCP servers (32,226 chars, ~9,200 tokens).
#
# Both were wrong before: messages used //3 (over-counting prose by ~30%) while
# tool schemas used //6 (under-counting JSON by ~40%, or ~5,400 tokens at a
# 128-tool ceiling). The two errors pulled in opposite directions, so the total
# looked plausible while neither half was.
CHARS_PER_TOKEN_PROSE = 4
CHARS_PER_TOKEN_JSON = 3.5

# Known model context windows — the *last* fallback, consulted only when the
# endpoint could not be probed and the fetched catalogue is not on disk.
#
# **A hand-maintained table of numbers rots, and this one had.** Checked
# against the live catalogue on 2026-08-13, seven of eighteen entries were
# wrong, in both directions: `deepseek-v4-flash-free` said 1,000,000 and serves
# 200,000 (the *paid* model's window, copied onto the free one), while
# `claude-opus-5` said 200,000 and serves 1,000,000. A model told it has five
# times the window it has will not compact until far too late; one told it has
# a fifth will compact constantly for nothing.
#
# The entries below are corrected, but correcting them is not the fix — they
# will rot again. `resolve_context_window` now asks `zen_catalog` first, which
# is fetched rather than typed. Prefer adding nothing here; if a model needs a
# number, it needs it in the catalogue.
MODEL_CONTEXT_MAP: dict[str, int] = {
    # Zen models
    "deepseek-v4-flash-free": 200_000,     # was 1_000_000 — the paid model's
    "deepseek-v4-flash": 1_000_000,
    "big-pickle": 200_000,                 # was 128_000
    "nemotron-3.5-lightning-free": 262_144,
    "nemotron-3-ultra-free": 1_000_000,
    "laguna-s-2.1-free": 256_000,
    "mimo-v2.5-free": 200_000,
    "hy3-free": 190_000,
    # Anthropic
    "claude-fable-5": 1_000_000,           # was 200_000
    "claude-opus-5": 1_000_000,            # was 200_000
    "claude-sonnet-5": 1_000_000,          # was 200_000
    "claude-haiku-4-5-20251001": 200_000,
    "claude-haiku-4-5": 200_000,
    "claude-sonnet-4-5": 1_000_000,        # was 200_000
    "claude-sonnet-4": 1_000_000,          # was 200_000
    "claude-3-5-sonnet-20241022": 200_000,
    "claude-3-haiku-20240307": 200_000,
    "claude-opus-4-5": 200_000,
    "claude-opus-4": 200_000,
    # OpenRouter / common
    "gpt-4o-mini": 128_000,
    "gpt-4o": 128_000,
    "gpt-4": 128_000,
    "gpt-3.5-turbo": 16_385,
    # DeepSeek
    "deepseek-chat": 128_000,
    "deepseek-reasoner": 128_000,
}

# Will be updated at startup if we can fetch from API
_current_context_window: int = DEFAULT_CONTEXT_WINDOW
#: Where `_current_context_window` came from — see core.context.WINDOW_SOURCES.
#: Displayed alongside the number so "measured 200,000" is distinguishable
#: from "knew nothing and assumed 200,000".
_current_context_source: str = "default"
AUTO_APPROVE_LOW = os.environ.get("AGENT_AUTO_APPROVE", "1") == "1"
YOLO_MODE = False  # when True, all tools are auto-approved without any prompt
# Bypass mode: yolo, plus the turn is never stopped to ask whether to keep
# going. The two are separate questions and yolo only ever answered the first.
# Observed: a session that approved every tool still halted at the 40-call
# checkpoint and was saved incomplete, mid-task, after 56 calls. Bounded by
# AgentState.max_auto_continuations so "do not ask me" cannot become "bill me
# without limit".
BYPASS_MODE = False

#: Mode order used by Tab cycling and by every place that names the modes.
#: One list, so a new mode cannot be added to the switcher and forgotten in the
#: badge, the status line and the help.
MODE_CYCLE = ("auto", "default", "yolo", "bypass")
ALL_MODES = ("auto", "default", "strict", "yolo", "bypass")


def current_mode_name() -> str:
    """The single answer to "what mode am I in?".

    Derived from the flags rather than stored beside them: several call sites
    used to recompute this inline, so a new mode has to be taught to each of
    them or the banner and the status line start disagreeing.
    """
    if BYPASS_MODE:
        return "bypass"
    if YOLO_MODE:
        return "yolo"
    return "auto" if AUTO_APPROVE_LOW else "default"


def mode_color(mode: Optional[str] = None) -> str:
    """The colour a mode is always shown in. Red means "nothing will stop"."""
    mode = mode or current_mode_name()
    if mode in ("yolo", "bypass"):
        return RED
    if mode == "auto":
        return GREEN
    return YELLOW


#: One line naming every key that switches modes, so the four places that
#: print it cannot drift apart when a mode is added.
MODE_KEYS_HINT = ("F5 auto · F6 default · F7 strict · F8 yolo · F9 bypass "
                  "· Tab cycles")


def set_mode(name: str) -> None:
    """Apply a mode by name. The one place the flags are written together."""
    global AUTO_APPROVE_LOW, YOLO_MODE, BYPASS_MODE
    YOLO_MODE = False
    BYPASS_MODE = False
    if name == "auto":
        AUTO_APPROVE_LOW = True
    elif name == "default":
        AUTO_APPROVE_LOW = False
    elif name == "strict":
        AUTO_APPROVE_LOW = False
        APPROVALS.clear()
        for key in list(RISK_LEVELS.keys()):
            if key not in BUILTIN_TOOL_NAMES:
                RISK_LEVELS[key] = "high"
    elif name == "yolo":
        AUTO_APPROVE_LOW = True
        YOLO_MODE = True
    elif name == "bypass":
        AUTO_APPROVE_LOW = True
        YOLO_MODE = True
        BYPASS_MODE = True

# ── Session token tracking ──
# Per-session, not per-process. These used to accumulate for the life of the
# interpreter, so two sessions run back to back reported byte-identical usage
# and a session that did no work still claimed 1.6M input tokens.
_session_tokens = {"input": 0, "output": 0, "calls": 0, "cached_input": 0,
                   # Measurement, not billing: what the streamed calls cost
                   # before being discarded and re-issued non-streamed, and how
                   # many of those re-issues bought nothing.
                   "duplicate_input": 0, "duplicate_calls": 0,
                   "would_have_served": 0, "stream_malformed_tool_args": 0}
_last_turn_usage = {"input": 0, "output": 0, "cached_input": 0}
# How the last turn ended, for the `advanced_diagnostics` line. The core
# already records all three on `AgentState` — they reached the session file
# and nothing else, so the answer to "why did that turn stop there?" existed
# on disk while the user watching it happen had no way to ask.
_last_turn_diag = {"stop_reason": "", "error": "", "tool_calls": 0}

# ── Session telemetry (P6-8) ──
# Per-turn wall clock and per-tool-call outcome, so a saved session can say
# which call was slow and which one failed.
_turn_timings: list[float] = []
_tool_log: list[dict] = []
_session_started_at: float = time.time()
# Turns that produced no assistant reply (e.g. retries exhausted on a 429).
_failed_turns: list[dict] = []
# Every compaction, with the arithmetic that triggered it. Compaction used to
# print a line to stdout and nothing else, so a saved session could not
# distinguish "compaction never ran" from "compaction ran and did not help" —
# which is exactly the question the V3 token analysis could not answer.
_context_log: list[dict] = []
# Turns that finished cleanly and did nothing: no tool call, and a reply too
# short to be an answer. `mimo-v2.5-free` replied "My Lord." — eight
# characters, zero tools — to four consecutive "read this file" instructions
# and every one was counted a success, because the only test applied was
# `reply.strip()` being non-empty.
#
# Kept separate from _failed_turns on purpose: this is a heuristic, and a
# heuristic must not be able to mark a good session incomplete. It reports,
# it does not judge.
_low_content_turns: list[dict] = []

#: A reply shorter than this, with no tool call, did no work. Deliberately
#: well under anything real: the shortest genuine reply in the V4 sweep was
#: 19 output tokens, while the empty ones were 4.
LOW_CONTENT_REPLY_CHARS = 24

# The current turn's Esc-interrupt signal, if the adapter driving it exposes
# one (see adapters.terminal.TerminalAdapter.esc_interrupt). Set by
# build_state() at the start of every turn. handle_run_command polls it to
# kill a running subprocess immediately rather than only at the next loop
# checkpoint, which for a shell command can be the whole timeout away.
# _call_mcp_tool_interruptibly polls the same signal for the same reason —
# an MCP call has no cancellation point of its own.
_CURRENT_INTERRUPT: Optional[threading.Event] = None


def reset_session_state() -> None:
    """Start a fresh session's accounting. Called when a session begins or
    when /clear discards the conversation."""
    global _session_started_at, _synthetic_replies
    # /clear discards the transcript, prefill included, so the allowance for
    # scaffolding turns has to go with it — otherwise the every-3rd-reply
    # cadence would be off by one in the other direction for the rest of the
    # session.
    _synthetic_replies = 0
    # Every key, including ones added since this dict was created. Naming them
    # here meant a counter added later survived the session boundary — the
    # exact bug TestSessionTokenIsolation exists to catch.
    _session_tokens.update(dict.fromkeys(_session_tokens, 0))
    _last_turn_usage.update(dict.fromkeys(_last_turn_usage, 0))
    _last_turn_diag.update({"stop_reason": "", "error": "", "tool_calls": 0})
    _ollama_window_checked.clear()
    _turn_timings.clear()
    _tool_log.clear()
    _failed_turns.clear()
    _context_log.clear()
    _low_content_turns.clear()
    _LAST_TOOL_SELECTION.clear()   # a new session starts with no sticky set
    _session_started_at = time.time()


def _record_compaction(strategy: str, before: int, after: int,
                       plan, error: Optional[str] = None) -> None:
    """Record one compaction so the session file can show it happened."""
    entry = {
        "turn": len(_turn_timings) + 1,
        "strategy": strategy,
        "before_tokens": before,
        "after_tokens": after,
        "reclaimed_tokens": max(0, before - after),
        "trigger": getattr(plan, "trigger", 0),
        "reason": getattr(plan, "reason", ""),
        "window": getattr(plan, "window", 0),
    }
    if error:
        entry["error"] = error
    _context_log.append(entry)


def session_telemetry() -> dict:
    """Telemetry block for the session file."""
    total = round(sum(_turn_timings), 2)
    n = len(_turn_timings)
    return {
        "turn_metrics": {
            "total_duration_sec": total,
            "avg_turn_sec": round(total / n, 2) if n else 0.0,
            "turn_timings": [round(t, 2) for t in _turn_timings],
        },
        "tool_log": list(_tool_log),
        "failed_turns": list(_failed_turns),
        "context_events": list(_context_log),
        "low_content_turns": list(_low_content_turns),
    }

# ── Session continuation ──
# Set by agent_cli.py before calling main() to continue a previous session.
CONTINUE_SESSION_ID: Optional[str] = None

# ── Client factory: supports ANTHROPIC_EXTRA_HEADERS env var (JSON) ──
_client_instance = None

def _get_client():
    """Return a cached client for the active provider.

    An endpoint that speaks OpenAI wire format gets the in-process adapter
    (no daemon, real incremental streaming); everything else gets the
    Anthropic SDK. Both present the same surface to core/loop.py.
    """
    global _client_instance
    # Read current env values each time — update_dotenv in the CLI sets os.environ
    key = os.environ.get("ANTHROPIC_API_KEY", "")
    base = os.environ.get("ANTHROPIC_BASE_URL", "")
    extra_hdr = os.environ.get("ANTHROPIC_EXTRA_HEADERS", "")
    # Build a simple cache key so we know when env changed
    cache_key = f"{key}::{base}::{extra_hdr}"
    if _client_instance is not None:
        prev_key = getattr(_client_instance, "_cache_key", None)
        if prev_key == cache_key:
            return _client_instance

    client = None
    try:
        import openai_adapter
        if openai_adapter.should_use_adapter():
            client = openai_adapter.build_from_active()
    except Exception:
        client = None

    if client is None:
        import anthropic  # deferred: see the note at the top of this module
        headers = json.loads(extra_hdr) if extra_hdr else None
        client = anthropic.Anthropic(
            api_key=key or None,
            base_url=base or None,
            default_headers=headers,
        )
    _client_instance = client
    _client_instance._cache_key = cache_key  # type: ignore[attr-defined]
    return _client_instance

def reinit_client():
    """Force the client to be re-created on next use (called after provider change)."""
    global _client_instance
    _client_instance = None


def _ensure_zen_proxy():
    """Start the Zen proxy daemon — only when explicitly asked for.

    Translation now runs in-process (openai_adapter), so the daemon is no
    longer on the agent's critical path. It stays available because pointing
    *other* tools at Zen is a real use for it: set TOMAS_ZEN_PROXY=1.
    """
    if os.environ.get("TOMAS_ZEN_PROXY", "") != "1":
        return
    base_url = os.environ.get("ANTHROPIC_BASE_URL", "")
    if "127.0.0.1:6446" in base_url or "localhost:6446" in base_url:
        try:
            from zen_proxy import check_status, start_proxy
            if not check_status(6446):
                print(f'  {DIM}Starting Zen proxy on port 6446...{RESET}')
                start_proxy(6446, daemon=True)
                import time
                time.sleep(0.5)
                if check_status(6446, use_cache=False):  # verifies our own start
                    print(f'  {GREEN}✓{RESET}  Zen proxy is running')
                else:
                    print(f'  {YELLOW}⚠{RESET}  Zen proxy may not have started')
            # else: proxy already running, nothing to do
        except ImportError:
            print(f'  {RED}✗{RESET}  zen_proxy module not found — cannot start proxy')
        except Exception as exc:
            print(f'  {RED}✗{RESET}  Failed to start Zen proxy: {exc}')


def _probe_models_endpoint(model: str) -> int:
    """Ask the endpoint for this model's window. 0 when it does not say.

    Split out of `_fetch_model_context_window` so the *probe* and the
    *fallbacks* are separable: `resolve_context_window` needs to know which of
    the two answered, and a function that silently substitutes a default
    cannot report that.
    """
    base_url = os.environ.get("ANTHROPIC_BASE_URL", "")
    if not base_url:
        return 0
    try:
        models_url = base_url.rstrip("/") + "/v1/models"
        req = urllib.request.Request(models_url, method="GET")
        # Use the same API key if set
        key = os.environ.get("ANTHROPIC_API_KEY", "")
        if key:
            req.add_header("x-api-key", key)
            req.add_header("anthropic-version", "2023-06-01")
        with urllib.request.urlopen(req, timeout=5) as resp:
            data = json.loads(resp.read().decode())
            # OpenAI-style: data.data[].context_window
            # Anthropic-style: data.data[].context_window
            models_list = data if isinstance(data, list) else data.get("data", [])
            for m in models_list:
                if m.get("id") == model:
                    cw = m.get("context_window") or m.get("context") or 0
                    if cw:
                        return int(cw)
            # Check for a top-level context_window field (dict response)
            if isinstance(data, dict):
                cw = data.get("context_window") or data.get("context") or 0
                if cw:
                    return int(cw)
    except Exception:
        pass
    return 0


def _probed_capability_window(model: str) -> int:
    """The window a completed provider probe measured, or 0.

    `provider_manager` probes the same `/v1/models` data independently and
    stores it on `Capabilities`. Consulting it here is what stops the two from
    being separate truths: `/status` used to read one and `/provider` the
    other, so the same model could be reported as 1,000,000 in one command and
    200,000 in the next.
    """
    try:
        import provider_manager
        active = provider_manager.get_active()
        if active is None or active.model != model:
            return 0
        caps = active.capabilities
        if caps and caps.probed and caps.context_window:
            return int(caps.context_window)
    except Exception:
        pass
    return 0


#: Resolved windows, keyed by model. Probing costs an HTTP round trip, so it
#: happens once per model per session rather than on every status line.
_context_window_cache: dict[str, ContextWindow] = {}


def _catalog_window(model: str) -> int:
    """The window `zen_catalog` has on disk for `model`, or 0.

    Sits between a probe and `MODEL_CONTEXT_MAP` in `resolve_context_window`:
    it is fetched rather than typed, so it does not rot the way the table did,
    but it is still not a measurement of *this* endpoint, so a probe outranks
    it.

    **Cache only — never the network.** This runs while a turn is being sized
    and while the status line is drawn, and `catalog()` will otherwise spend up
    to eight seconds on an availability request. A cold cache simply falls
    through to the table, which is the behaviour that exists today; the cache is
    filled by the model picker and by any earlier catalogue read.

    `context_known` matters: `zen_catalog` reports an assumed 128,000 for
    models `models.dev` has never described, and adopting an assumption as a
    measured window is the exact failure this whole change is about.
    """
    if not model:
        return 0
    try:
        import zen_catalog

        entry = zen_catalog.catalog(allow_network=False).get(model)
    except Exception:
        return 0
    return int(entry.context) if entry and entry.context_known else 0


def resolve_context_window(model: Optional[str] = None,
                           refresh: bool = False) -> ContextWindow:
    """The context window for `model`, with the provenance of the number.

    Precedence runs from most to least trustworthy — an explicit override, a
    measurement, a known value, then a guess. Previously the order depended on
    whether `ANTHROPIC_BASE_URL` happened to be set, so an environment
    variable decided which source won rather than how reliable it was.
    """
    model = model or _get_model() or ""
    if not refresh and model in _context_window_cache:
        return _context_window_cache[model]

    override = _env_int("AGENT_CONTEXT_WINDOW", 0, minimum=1)
    if override:
        window = ContextWindow(override, "override", model)
    else:
        probed = _probed_capability_window(model) or _probe_models_endpoint(model)
        catalogued = 0 if probed else _catalog_window(model)
        if probed:
            window = ContextWindow(probed, "probed", model)
        elif catalogued:
            window = ContextWindow(catalogued, "catalog", model)
        elif model in MODEL_CONTEXT_MAP:
            window = ContextWindow(MODEL_CONTEXT_MAP[model], "known", model)
        else:
            window = ContextWindow(DEFAULT_CONTEXT_WINDOW, "default", model)

    _context_window_cache[model] = window
    return window


def context_window_divergence(model: Optional[str] = None) -> Optional[tuple]:
    """`(known, probed)` when the table and the endpoint disagree, else None.

    A stale `MODEL_CONTEXT_MAP` is worth knowing about: it is a hardcoded
    table and it will rot. Silence would just mean trusting whichever source
    happened to answer first.
    """
    model = model or _get_model() or ""
    known = MODEL_CONTEXT_MAP.get(model)
    probed = _probed_capability_window(model)
    if known and probed and known != probed:
        return (known, probed)
    return None


def _fetch_model_context_window() -> int:
    """The active model's window as a plain int.

    Kept as the historical entry point; the provenance-carrying answer is
    `resolve_context_window()`.
    """
    return resolve_context_window().tokens


def _refresh_context_window() -> int:
    """Re-fetch and apply the context window for the active model.

    Called after a mid-session provider/model switch (see the `/model` and
    `/provider` slash commands) so compaction math (`CONTEXT_WINDOW`) and the
    status line (`_current_context_window`) reflect the new model rather than
    the one the session started with.
    """
    global _current_context_window, CONTEXT_WINDOW, _current_context_source
    # A model switch invalidates the cached answer: the point of refreshing is
    # to stop reporting the window of the model the session started with.
    window = resolve_context_window(refresh=True)
    _current_context_window = window.tokens
    _current_context_source = window.source
    CONTEXT_WINDOW = _current_context_window
    return _current_context_window


# MCP manager (initialized at startup when main() is called)
mcp_manager: Optional[MCPManager] = None
COMBINED_TOOLS: list[dict] = []
# Every tool that exists after name resolution — built-ins plus every
# connected MCP tool. COMBINED_TOOLS is the subset that fits the budget;
# select_tools() picks that subset per turn from this list.
ALL_TOOLS: list[dict] = []
TOOL_TOKENS: int = 0  # estimated token count for tool definitions
#: One-shot latch for the "overhead leaves no room" warning — a configuration
#: problem deserves one clear message, not one per turn.
_warned_overhead: bool = False
MCP_TOOL_NAME_MAP: dict[str, str] = {}  # renamed_name -> original_name for conflicting MCP tools


def resolve_mcp_tool_conflicts(mcp_tools: list[dict], builtin_names: Optional[set] = None) -> tuple[list[dict], dict[str, str], int]:
    """
    Resolve MCP/built-in tool name collisions by prefixing conflicting MCP
    tools with 'mcp_' (e.g. read_file -> mcp_read_file).

    Returns (resolved_mcp_tools, name_map, renames) where:
      - resolved_mcp_tools: MCP tools with collisions renamed (originals untouched)
      - name_map:          {renamed_name: original_name} for renamed tools
      - renames:           number of tools that were renamed
    """
    if builtin_names is None:
        builtin_names = {t["name"] for t in TOOLS}
    resolved: list[dict] = []
    name_map: dict[str, str] = {}
    renames = 0
    for t in mcp_tools:
        original = t["name"]
        if original in builtin_names:
            new_name = f"mcp_{original}"
            renamed = dict(t)
            renamed["name"] = new_name
            renamed["description"] = f"[MCP: {original}] {renamed.get('description', '')}"
            resolved.append(renamed)
            name_map[new_name] = original
            renames += 1
        else:
            resolved.append(t)
    return resolved, name_map, renames


def apply_tool_cap(mcp_tools: list[dict], max_allowed: int = 128) -> tuple[list[dict], int]:
    """
    Merge built-in TOOLS with MCP tools and truncate to max_allowed tools.
    Built-in tools are always kept; excess MCP tools are dropped in order.

    This is the no-context fallback. When there is a user message to select
    against, `select_tools` picks by relevance instead — see P4-8. Kept
    because startup has no context yet and something must be sent.

    Returns (combined_tools, dropped).
    """
    n_builtin = len(TOOLS)
    keep = max(0, max_allowed - n_builtin)
    dropped = max(0, len(mcp_tools) - keep)
    if dropped:
        return TOOLS + mcp_tools[:keep], dropped
    return TOOLS + mcp_tools, 0


def tool_relevance(tool: dict, query_keywords: set) -> float:
    """Score one tool against the words the user just used.

    Same shape as `learning.retrieval.score_fact` — one scoring idea, two
    uses. Name matches count double: a tool called `take_screenshot` is a
    better answer to "screenshot" than one that merely mentions screenshots
    in its description.
    """
    if not query_keywords:
        return 0.0
    from learning.text import extract_keywords
    name = (tool.get("name") or "").replace("_", " ").replace("-", " ")
    desc = tool.get("description") or ""
    name_words = set(extract_keywords(name, max_keywords=8))
    desc_words = set(extract_keywords(desc, max_keywords=20))
    score = 2.0 * len(query_keywords & name_words) + len(query_keywords & desc_words)
    # Cheap substring credit for compound names retrieval would not split.
    lowered = f"{tool.get('name','')} {desc}".lower()
    score += 0.5 * sum(1 for kw in query_keywords if len(kw) > 4 and kw in lowered)
    return score


def tool_name_matches(tool: dict, query_keywords: set) -> int:
    """How many of the user's words appear in the tool's *name*.

    Separated from `tool_relevance` because the two are used for different
    decisions. A description hit is weak evidence — every browser tool mentions
    "time" in a timeout note — while a name hit means the request said what the
    tool is called. `SERVER_CORE_QUOTA` commits eight slots to one server, so
    it asks the stronger question.

    Measured: "remember that I prefer Ukrainian" matched `fill_form` at 2.0
    purely through its description, which was enough to make chrome-devtools
    the turn's primary server and hand it eight slots for a request about
    remembering a preference.
    """
    if not query_keywords:
        return 0
    from learning.text import extract_keywords
    name = (tool.get("name") or "").replace("_", " ").replace("-", " ")
    return len(query_keywords & set(extract_keywords(name, max_keywords=8)))


def _server_of(name: str) -> Optional[str]:
    """Which MCP server owns this tool. None for built-ins or when no manager
    is connected — callers must treat None as "ungrouped", not as an error."""
    if mcp_manager:
        try:
            return mcp_manager.get_server_for_tool(MCP_TOOL_NAME_MAP.get(name, name))
        except Exception:
            return None
    return None


def tool_simplicity(name: str) -> float:
    """How fundamental a tool name looks, from its segment count alone.

    Within one server the building blocks are the plainly-named tools —
    `add_heading`, `add_table` — while the specialised variants pile on
    qualifiers: `set_table_cell_shading`, `apply_table_alternating_rows`,
    `delete_footnote_robust`. A task that engages a server needs the former;
    it reaches for the latter by name, when it already knows it wants them.
    Segment count is the only signal here that does not require knowing what
    any particular server does.
    """
    segments = [s for s in re.split(r"[_\-]+", name or "") if s]
    return 1.0 / max(1, len(segments) - 1)


# How many tools the turn's primary server is guaranteed, before everything
# else competes on score. Sized to cover a working set (create + the handful
# of verbs that put content into what was created), not a whole server.
SERVER_CORE_QUOTA = 8

# A quota slot is only worth spending on a tool the turn plausibly needs: one
# the request actually scored, or one plainly enough named to be a building
# block (at most three segments). Without this the quota keeps filling from
# whatever is left — `set_table_cell_alternating_shading` and friends — after
# the real working set runs out, and spends budget other servers could use.
CORE_NAME_SIMPLICITY = 0.5

# How many tools carried over from last turn may keep their slot without
# scoring on *this* message.
#
# Follow-ups are why this is not zero: "yes, do that" and "now the same for the
# other file" score nothing at all, and dropping the working set on them would
# make every second turn re-fetch its tools. Bounded because the alternative is
# accumulation — each new topic adds its tools and none of the old ones ever
# leave, so after a few pivots the payload is full again and "selected by
# relevance" has quietly become "everything, eventually".
STICKY_CARRY_OVER = 8

# A score is only meaningful next to the best score on the same message.
#
# `tool_relevance` gives 2.0 per name-word hit and 1.0 per description-word
# hit, so a single incidental word in a description earns 1.0 and, on a
# fixed threshold, a slot. Measured against the live pool:
#
#   "take a screenshot of example.com"   6.5  browser_take_screenshot
#                                        0.5  resolve-library-id, query-docs
#   "what time is it in Tokyo"           3.0  get_current_time, convert_time
#                                        1.0  bulk_fetch, browser_wait_for
#                                             (their descriptions say "timeout")
#
# In both cases the gap between the tools the turn is about and the ones that
# merely share a word is large and obvious — but only relative to that turn.
# A fixed cut-off tuned for the first would discard everything in the second.
RELEVANCE_FLOOR_SHARE = 0.35

# …and never below this, or a message whose best match is itself incidental
# fills the payload with a dozen equally incidental ones.
MIN_RELEVANCE = 1.0


def relevance_floor(scores) -> float:
    """The score a tool must reach on this message to earn a slot on merit."""
    best = max(scores, default=0.0)
    if best <= 0:
        return MIN_RELEVANCE
    return max(MIN_RELEVANCE, best * RELEVANCE_FLOOR_SHARE)


# The MCP tools sent last turn. Used only as a tie-breaker, so it can never
# override relevance — see `rank` inside select_tools.
_LAST_TOOL_SELECTION: set[str] = set()


def _remember_tool_selection(names) -> None:
    _LAST_TOOL_SELECTION.clear()
    _LAST_TOOL_SELECTION.update(names)


def select_tools(all_tools: list[dict], context: str, budget: int,
                 server_of: Optional[Callable[[str], Optional[str]]] = None,
                 allowlist: Optional[set] = None,
                 ) -> tuple[list[dict], list[dict]]:
    """Fit the tool list to the budget by relevance to the current turn.

    Built-ins are always kept — they are what the agent is. The rest of the
    budget goes to the MCP tools most relevant to what the user is actually
    doing, rather than to whichever server connected first, which is what
    list-order truncation really selected on.

    Relevance alone is not enough, because it scores every tool in isolation.
    Asked to write a methodichka as .docx and .pdf, it handed the model
    `create_document` (3.5) and `convert_to_pdf` (3.0) while every tool that
    puts content *into* a document — `add_heading`, `add_paragraph`,
    `add_table` — scored exactly 0.0 and was withheld: nothing in those names
    or descriptions repeats the user's words. The model could open a document
    and convert one, but could not fill one, so it fell back to `write_file`
    and produced a .docx that was not a real .docx at all.

    So the server that best matches the turn is picked first, by summed
    relevance across its tools, and is guaranteed a `SERVER_CORE_QUOTA`
    working set — its highest-scoring tools, then its plainest-named ones.
    The rest of the budget is filled globally by score, as before.

    Returns (selected, withheld).
    """
    # A triggered skill may declare the tools it needs (`tools:` in its
    # frontmatter). Relevance scoring cannot reach this conclusion on its own:
    # it ranks tools against the *user's words*, and "зроби схожий файл" scores
    # the whole word-docs and pdf servers highly — ~13k tokens of schema per
    # turn — for a skill whose entire job is running one subprocess. The skill
    # knows; nothing else does. Built-ins are never withheld, so an allowlist
    # narrows the MCP payload and cannot strand the agent.
    if allowlist:
        kept, dropped = [], []
        for tool in all_tools:
            (kept if tool["name"] in allowlist else dropped).append(tool)
        all_tools = kept
        budget = max(budget, len(kept))
    else:
        dropped = []

    builtin_names = {t["name"] for t in TOOLS}
    builtins = [t for t in all_tools if t["name"] in builtin_names]
    mcp = [t for t in all_tools if t["name"] not in builtin_names]
    remaining = max(0, budget - len(builtins))

    if len(mcp) <= remaining:
        return builtins + mcp, dropped
    if remaining == 0:
        return builtins, mcp + dropped

    from learning.text import extract_keywords
    query_keywords = set(extract_keywords(context or "", max_keywords=15))
    if not query_keywords:
        return builtins + mcp[:remaining], mcp[remaining:] + dropped

    resolve = server_of or _server_of
    scores = [tool_relevance(t, query_keywords) for t in mcp]
    servers = [resolve(t.get("name", "")) for t in mcp]

    # Reuse last turn's set verbatim when it still covers this turn.
    #
    # Prefix caching is binary: one changed byte anywhere in the tool block and
    # the whole cached prefix is gone, so a set that is *mostly* stable is
    # worth almost nothing. Holding it byte-identical is what turns the system
    # prompt, the tools and the entire history into a cache hit.
    #
    # "Still covers this turn" means the single most relevant tool for the
    # message is already loaded. If the user pivots to a capability that is not
    # in hand, that tool outranks everything present and the set is re-picked —
    # which is the case adaptivity actually exists for.
    if _LAST_TOOL_SELECTION and any(s > 0 for s in scores):
        best = max(range(len(mcp)), key=lambda i: (scores[i], -i))
        if mcp[best].get("name", "") in _LAST_TOOL_SELECTION:
            held = [t for t in mcp if t.get("name", "") in _LAST_TOOL_SELECTION]
            if len(held) == min(remaining, len(_LAST_TOOL_SELECTION)):
                kept = sorted(held, key=lambda t: t.get("name", ""))
                names = {t.get("name", "") for t in kept}
                return (builtins + kept,
                        [t for t in mcp if t.get("name", "") not in names] + dropped)

    # Rank servers by summed relevance, not by their single best tool: one
    # incidental keyword hit should not outrank a server the whole request is
    # about. `memory` matched "create" twice; `word-docs` matched across
    # create/convert/copy and wins on the total.
    server_totals: dict[str, float] = {}
    for srv, score in zip(servers, scores):
        if srv is not None:
            server_totals[srv] = server_totals.get(srv, 0.0) + score

    # Sort key shared by both phases: score, then whether the tool was already
    # loaded last turn, then plainest name, then the original order so ties
    # never depend on dict iteration.
    #
    # The stickiness term is what makes this affordable to send. Tools are
    # serialised *before* the messages, and every provider that caches a prompt
    # — DeepSeek's automatic caching, Anthropic's cache_control, OpenAI's
    # prefix cache — matches on an exact byte prefix. Selecting purely by score
    # swapped 38 of 64 tool slots between consecutive turns, so the cached
    # prefix was invalidated at the first block and the system prompt, the
    # tools and the whole conversation were re-processed every single message.
    #
    # It is only a tie-breaker, so relevance is unaffected: a tool the turn
    # actually scores still beats a stale one. What it stabilises is the long
    # tail of zero-scoring filler, which is where all the churn was.
    sticky = _LAST_TOOL_SELECTION

    def rank(i: int) -> tuple:
        name = mcp[i].get("name", "")
        return (-scores[i], 0 if name in sticky else 1,
                -tool_simplicity(name), i)

    # A server only becomes "the turn's server" if the request named one of its
    # tools. Summed description hits used to be enough, and eight slots is far
    # too large a commitment to make on that evidence — see `tool_name_matches`.
    named_servers = {srv for srv, tool in zip(servers, mcp)
                     if srv is not None and tool_name_matches(tool, query_keywords)}

    taken: set[int] = set()
    primary = max(server_totals, key=lambda s: (server_totals[s], s), default=None)
    if primary is not None and primary in named_servers:
        core = sorted(
            (i for i, s in enumerate(servers)
             if s == primary
             and (scores[i] > 0
                  or tool_simplicity(mcp[i].get("name", "")) >= CORE_NAME_SIMPLICITY)),
            key=rank)
        taken.update(core[:min(SERVER_CORE_QUOTA, remaining)])

    # An empty slot is cheaper than a wrong tool.
    #
    # This loop used to fill the budget to the brim from whatever was left,
    # which meant that once the tools the turn actually scored ran out, the
    # rest of the payload was decided by name length and list order. Measured
    # against the live 257-tool pool at a 36-slot budget:
    #
    #   "what time is it in Tokyo"      21 of 36 slots went to the pdf server
    #   "hello"                         36 tools, none of them scoring at all
    #
    # At ~125 tokens per schema that is ~4,500 tokens per turn spent on tools
    # chosen by accident, and the model reads every one of them. So a slot is
    # now spent only on a tool that earned it: one this message scored, or one
    # already in hand from the last (bounded by STICKY_CARRY_OVER). The
    # primary server's core quota above is unaffected — that is the deliberate
    # exception, and the reason `add_heading` still ships with `create_document`
    # despite scoring zero.
    #
    # Under-filling is safe because it is *visible*: `withheld_tools_notice`
    # names what is not loaded, so a gap the model needs to close is one it can
    # see, which was never true of a gap hidden behind 36 irrelevant schemas.
    floor = relevance_floor(scores)
    # Carry the last set over only when *this* message has no direction of its
    # own. "yes, do that" scores nothing and needs the working set; "extract
    # the tables from this PDF" scores plenty and does not need the browser
    # tools left over from two turns ago. Carrying unconditionally is how the
    # payload silently refilled: 8 stale slots on every turn, whatever it was
    # about.
    carry_allowed = STICKY_CARRY_OVER if max(scores, default=0.0) <= 0 else 0
    carried = 0
    for i in sorted(range(len(mcp)), key=rank):
        if len(taken) >= remaining:
            break
        if i in taken:
            continue
        if scores[i] >= floor:
            taken.add(i)
        elif (carried < carry_allowed
              and mcp[i].get("name", "") in sticky):
            taken.add(i)
            carried += 1

    # Serialise by name, not by score. Selection is already decided by this
    # point and the model does not care about the order — but the byte stream
    # does: an unchanged set ordered by a score that shifts every turn produces
    # different bytes each time, and misses the prefix cache for no reason.
    chosen = sorted((mcp[i] for i in taken), key=lambda t: t.get("name", ""))
    withheld = [mcp[i] for i in sorted(set(range(len(mcp))) - taken, key=rank)]
    _remember_tool_selection(t.get("name", "") for t in chosen)
    return builtins + chosen, withheld + dropped


def withheld_tools_notice(withheld: list[dict]) -> str:
    """One line telling the model what it does not currently have.

    A silent capability gap is unrecoverable — the model cannot ask for a
    tool it has no idea exists. Naming the servers makes it recoverable.
    """
    if not withheld:
        return ""
    servers = []
    for t in withheld:
        origin = _tool_origin(t.get("name", ""))
        if origin.startswith("MCP: "):
            srv = origin[5:]
            if srv not in servers:
                servers.append(srv)
    where = f" from {', '.join(sorted(servers)[:6])}" if servers else ""
    return (f"\n\n{len(withheld)} additional tool(s){where} are connected but not "
            f"loaded for this turn. If you need one, say which capability you "
            f"need and the user can re-ask; do not assume it is unavailable.")


# How much of a tool's prose survives into the payload. Measured across the
# 64 tools of three real MCP servers: the mean schema is 503 chars (~125
# tokens), but the distribution has a long tail — `sequentialthinking` alone is
# 4,056 chars. At a 128-tool ceiling the block costs ~16,100 tokens a turn,
# four times the entire system prompt, and it is re-sent every single turn.
#
# These ceilings are set above the mean on purpose: the point is to clip the
# outliers, not to compress every tool. A description short enough to be
# ambiguous costs a wrong tool call, which is far dearer than the tokens saved.
MAX_TOOL_DESC_CHARS = _env_int("TOMAS_MAX_TOOL_DESC", 600, minimum=120)
MAX_TOOL_PARAM_DESC_CHARS = _env_int("TOMAS_MAX_PARAM_DESC", 220, minimum=40)

# Schema keys that describe the document rather than the call. Models do not
# read them and validators upstream do not require them.
_DROPPABLE_SCHEMA_KEYS = ("$schema", "title", "examples", "additionalProperties",
                          "$id", "definitions", "$defs")


def _clip(text: str, limit: int) -> str:
    """Trim prose to `limit`, preferring a sentence then a word boundary."""
    if not isinstance(text, str) or len(text) <= limit:
        return text
    cut = text[:limit]
    stop = max(cut.rfind(". "), cut.rfind("\n"))
    if stop > limit * 0.6:
        return cut[:stop + 1].rstrip()
    space = cut.rfind(" ")
    return (cut[:space] if space > limit * 0.6 else cut).rstrip() + "…"


def _compact_schema(node, depth: int = 0):
    """Recursively drop documentation-only keys and clip property prose.

    Everything that decides whether a call is well-formed — `type`,
    `properties`, `required`, `enum`, `items`, `default` — is preserved
    untouched. Only prose and metadata are reduced.
    """
    if isinstance(node, list):
        return [_compact_schema(v, depth + 1) for v in node]
    if not isinstance(node, dict):
        return node
    out = {}
    for key, value in node.items():
        if key in _DROPPABLE_SCHEMA_KEYS:
            continue
        if key == "description":
            # The top level of the schema is the tool's own blurb and is
            # clipped by the caller; nested ones are per-property.
            out[key] = _clip(value, MAX_TOOL_PARAM_DESC_CHARS) if depth else value
        else:
            out[key] = _compact_schema(value, depth + 1)
    return out


def compact_tool_schemas(tools: list[dict]) -> list[dict]:
    """Shrink the tool payload without changing what any tool accepts.

    Pure: returns new dicts and never mutates `ALL_TOOLS`, because the same
    tool objects are re-selected every turn and clipping them in place would
    compound turn after turn until the descriptions were gone.
    """
    compacted = []
    for tool in tools:
        if not isinstance(tool, dict):
            compacted.append(tool)
            continue
        new = dict(tool)
        if isinstance(new.get("description"), str):
            new["description"] = _clip(new["description"], MAX_TOOL_DESC_CHARS)
        schema = new.get("input_schema")
        if isinstance(schema, dict):
            new["input_schema"] = _compact_schema(schema)
        compacted.append(new)
    return compacted


def _max_tokens_was_set() -> bool:
    """True when the user named an output budget rather than taking the default."""
    return bool(os.environ.get("AGENT_MAX_TOKENS", "").strip())


#: Default sampling temperature. Low because this agent's output is documents,
#: code and tool arguments, none of which benefit from creative sampling — and
#: because the provider default it replaces is 1.0, which is what shook a small
#: free model apart mid-document (see `AgentState.temperature`). Not zero: some
#: endpoints treat 0 as "unset", and a model that repeats itself verbatim when
#: a tool call fails is its own failure mode.
DEFAULT_TEMPERATURE = 0.3

#: Wall-clock ceiling for one turn, seconds. `AGENT_MAX_TURN_SECONDS=0` removes
#: it.
#:
#: 20 minutes, from measurement, not taste. Across the labwork sweep the
#: slowest turn that *finished usefully* was 463 s; the two pathological ones
#: were 1,986 s (89 tool calls on one instruction) and a 75-minute stall that
#: produced nothing and had to be killed by hand. Raised from 900s after a
#: third pathological run: `document-style-match`'s Step 1 (open-ended
#: "analyze the sample") let a model re-measure margins and audit unrelated
#: files' conventions for the full 900s and hit the ceiling having never
#: reached Step 3 — the deadline was not generous, Step 1 had no exit
#: condition. That gap is now closed in the skill itself (a stop rule once
#: the four contract numbers are in hand) and the deadline now spends one
#: bounded call trying to salvage something instead of returning nothing
#: (`core.loop._deadline_wind_down`) — this bump is headroom for that salvage
#: call and for legitimately larger documents, not a substitute for either
#: fix. It still sits well below both pathological failures above.
MAX_TURN_SECONDS = _env_int("AGENT_MAX_TURN_SECONDS", 1200, minimum=0)


def effective_temperature() -> Optional[float]:
    """The sampling temperature for this turn, or None to send nothing.

    `AGENT_TEMPERATURE=` (empty) or `off` opts out entirely and restores the
    old behaviour of accepting whatever the endpoint does — kept reachable
    because a provider that rejects the parameter must remain usable, and
    because "the agent stopped sending a field it used to send" needs to be
    something a user can undo without editing source.
    """
    raw = os.environ.get("AGENT_TEMPERATURE")
    if raw is None:
        return DEFAULT_TEMPERATURE
    raw = raw.strip().lower()
    if raw in ("", "off", "none", "provider", "default"):
        return None
    try:
        return max(0.0, min(2.0, float(raw)))
    except ValueError:
        return DEFAULT_TEMPERATURE


def effective_max_tokens(caps) -> int:
    """This turn's output budget, honouring an explicit setting.

    `Capabilities.max_output_tokens` is documented as "the optimistic case;
    probing only ever narrows them" — but *nothing probes this field*, so its
    8192 default was being applied as a hard ceiling that no measurement ever
    justified. `min(MAX_TOKENS, caps.max_output_tokens)` therefore silently
    discarded AGENT_MAX_TOKENS: setting it to 32000 still ran at 8192, and the
    truncation message told the user to raise a variable that could not work.

    Measured against the endpoint that hit this: deepseek-v4-flash-free accepts
    max_tokens up to at least 65,536. The cap cost that provider 8x its real
    budget, which is what truncated a lab-manual build mid-write.

    So an explicit setting always wins here. Note this turns on the *field*
    never being probed, not on `caps.probed` — that flag goes true once any
    probe has run, while `max_output_tokens` stays untouched, so testing it
    would reinstate exactly the clamp this removes. If real probing is ever
    added for this field, make a measured value win instead.
    """
    if _max_tokens_was_set():
        return MAX_TOKENS
    # Otherwise the context budget decides, because the right reserve depends
    # on the window and 8192 was chosen when every window was assumed to be
    # 200,000. `output_reserve()` returns MAX_TOKENS if the budget is
    # unavailable, so the old behaviour is still the fallback.
    reserve = output_reserve()
    ceiling = getattr(caps, "max_output_tokens", 0) or reserve
    return min(reserve, ceiling) or reserve


def _claim_local_model() -> None:
    """Tell `ollama_runtime` this session is about to load a local model.

    Silent and best-effort: every failure here costs a model staying resident
    for the five minutes it would have stayed resident before this existed.
    """
    try:
        import provider_manager
        import ollama_runtime
        active = provider_manager.get_active()
        if active is None or active.type != "ollama" or not active.model:
            return
        ollama_runtime.claim(active.model,
                             provider_manager.ollama_native_root(active),
                             provider_manager.ollama_request_headers(active))
    except Exception:
        pass


def _release_local_models() -> None:
    """Give back the VRAM this session's local models are holding.

    Ollama keeps a model resident for five minutes after the last request, so
    a session that has ended goes on occupying the card for no reason —
    measured, three models and 8.5 GB left over from earlier runs. Unloading
    is one call; the care is in *not* unloading a model another live TOMAS is
    still using, which `ollama_runtime` decides by whether that session
    still holds its lock.

    Runs after the session is saved and after MCP is disconnected, because it
    is the one piece of cleanup whose failure genuinely does not matter: the
    five-minute timer is the fallback, and it is the behaviour that shipped.
    """
    try:
        import ollama_runtime
        results = ollama_runtime.unload_session_models()
    except Exception:
        return
    for model, freed, note in results:
        if freed:
            print(f'  {DIM}⏏  Unloaded {model} from memory{RESET}')
        elif note:
            print(f'  {DIM}⏏  Kept {model} loaded — {note}{RESET}')


#: Models whose real served window has already been checked this session, so
#: the reading costs one HTTP GET per model rather than one per turn.
_ollama_window_checked: set = set()
# Reset by `reset_session_state`, like every other per-session counter: a
# `/clear` starts a new session and the notice is worth showing again.


def _verify_ollama_window() -> None:
    """Compare the window TOMAS is budgeting against with what Ollama served.

    Every estimate made before the first call is reasoning about a server this
    process cannot see the configuration of. A request through the OpenAI shim
    allocates the runner itself, so `/api/ps` immediately afterwards reports
    the number that request was actually given — see
    `provider_manager.ollama_shim_window`. This is that reading, taken once
    per model per session, on a turn that has already happened.

    It reports and corrects; it never raises and never retries. When the two
    disagree the served number wins, because the disagreement is not a
    difference of opinion: the smaller one is what the model was handed, and
    everything past it was discarded before inference began.
    """
    global CONTEXT_WINDOW
    try:
        import provider_manager
        active = provider_manager.get_active()
        if active is None or active.type != "ollama" or not active.model:
            return
        if active.model in _ollama_window_checked:
            return
        served = provider_manager.ollama_shim_window(active)
        if not served:
            return                      # not loaded: no reading, not "no window"
        _ollama_window_checked.add(active.model)
        believed = CONTEXT_WINDOW
        if served >= believed:
            return
        # Persist it, so the next session budgets correctly from turn one
        # instead of re-learning this after the first prompt is truncated.
        try:
            active.capabilities.context_window = served
            provider_manager.persist_capabilities(active)
        except Exception:
            pass
        _context_window_cache.pop(active.model, None)
        CONTEXT_WINDOW = served
        print(f'\n  {RED}⚠{RESET}  {BOLD}Ollama served {served:,} tokens of '
              f'context, not the {believed:,} this model advertises.{RESET}')
        print(f'     {DIM}Anything past {served:,} is dropped before the model '
              f'sees it — no error, no warning. Budgeting at {served:,} from '
              f'here.{RESET}')
        print(f'     {DIM}Raise it by starting Ollama with '
              f'{RESET}{CYAN}OLLAMA_CONTEXT_LENGTH{RESET}{DIM} set, or give the '
              f'model a Modelfile with{RESET} {CYAN}PARAMETER num_ctx{RESET}'
              f'{DIM}.{RESET}')
    except Exception:
        pass


def _active_capabilities():
    """Probed capabilities of the active provider, or optimistic defaults.

    Never probes on this path — capabilities are data read from disk, and a
    turn must not pay network latency to find out what it can do.
    """
    try:
        import provider_manager
        return provider_manager.capabilities_for_active()
    except Exception:
        from types import SimpleNamespace
        return SimpleNamespace(
            streaming=True, tool_use=True, parallel_tool_calls=True,
            system_prompt=True, prompt_caching=False, vision=False,
            context_window=DEFAULT_CONTEXT_WINDOW, max_tools=128,
            max_output_tokens=MAX_TOKENS, probed_at=0.0)


def degrade_capability(field: str, reason: str = "") -> None:
    """Record that the provider cannot do something, and remember it.

    This is the whole of 'degrade, never fail': a capability the provider
    lacks costs the user a feature, never the session. The first discovery is
    paid once, then persisted so the next session starts already knowing.
    """
    try:
        import provider_manager
        provider = provider_manager.get_active()
        if provider is None or not getattr(provider.capabilities, field, False):
            return
        setattr(provider.capabilities, field, False)
        provider_manager.persist_capabilities(provider)
        note = f" ({reason})" if reason else ""
        print(f'  {YELLOW}⚠{RESET}  {provider.name} cannot {field.replace("_", " ")}'
              f'{note} — falling back{DIM}, remembered for next time{RESET}')
    except Exception:
        pass


def provider_tool_ceiling() -> int:
    """How many tools this endpoint will accept.

    Read from the active provider's probed capabilities. It used to be
    inferred from whether the model name contained the substring "free" —
    which cut a model called `my-free-model` to a quarter of its budget for
    its name, and gave a self-hosted endpoint the cloud default it could not
    honour.

    This is the *endpoint's* limit. What the window can afford is a separate
    and usually smaller question — see `tool_ceiling`.
    """
    try:
        import provider_manager
        return max(1, int(provider_manager.capabilities_for_active().max_tools))
    except Exception:
        return 128


# ══════════════════════════════════════════════════════════════════════
#  Context budget — what may occupy the window
# ══════════════════════════════════════════════════════════════════════
#
# The endpoint's tool ceiling answers "how many tools will this API accept?".
# It was being used to answer "how many tools should we send?", which is a
# different question with a different answer: 64 tools is fine for a 200,000
# token window and catastrophic for a 32,768 one, where measurement put the
# tool block at 18,079 tokens — 61% of everything, before the user typed. The
# policy lives in `core.budget`; this section is the wiring.

_budget_settings: Optional[Any] = None
_budget_cache: Optional[tuple] = None

# Bumped once, at the single point in init_mcp() where ALL_TOOLS is finalized
# — never on the turn path. Used to cache _avg_tool_tokens()'s json.dumps pass
# over the whole pool, which used to re-run unconditionally on every call,
# including the ones made just to check whether active_budget()'s own cache
# still applied.
_tool_pool_version = 0
_avg_tool_tokens_cache: Optional[tuple[int, int]] = None


def budget_settings(refresh: bool = False):
    """The user's stored budget choices, read once per session."""
    global _budget_settings
    if _budget_settings is None or refresh:
        _budget_settings = core_budget.load_settings()
    return _budget_settings


def save_budget_settings(settings) -> None:
    """Persist new choices and make them take effect immediately.

    The prompt cache is invalidated because two of the toggles (instructions,
    skills catalogue) live in the *stable* half of the system prompt, which is
    memoised on a filesystem fingerprint that no setting change would move.
    Without this a user could switch the catalogue off and watch the token
    count refuse to budge until they restarted.
    """
    global _budget_settings, _budget_cache
    core_budget.save_settings(settings)
    _budget_settings = settings
    _budget_cache = None
    invalidate_prompt_cache()
    _refresh_tool_tokens()


_features: Optional[Any] = None


def features(refresh: bool = False):
    """The user's feature switches, read once per session.

    Same shape as `budget_settings` deliberately: both are user choices
    persisted under `~/.tomas/`, both are read on the turn path, and one
    caching pattern for the two is one thing to get right rather than two.
    """
    global _features
    if _features is None or refresh:
        _features = core_features.load()
        # The debug recorder holds whole request payloads, so it stays off
        # until something asks for it — see core/debug_log.py. Applied on load
        # rather than only on toggle, or a switch left on in the settings file
        # would not start recording until the user toggled it again.
        _apply_debug_setting(_features.enabled("debug_view"))
    return _features


def save_features(new_features) -> None:
    """Persist the switches and make them take effect immediately."""
    global _features
    core_features.save(new_features)
    _features = new_features
    _apply_debug_setting(new_features.enabled("debug_view"))


def _apply_debug_setting(on: bool) -> None:
    """Turn recording on or off, and give it somewhere to stream to.

    The live file is per-process, not per-session: two TOMAS windows tailing
    one file would interleave their traffic into something neither user could
    read. It is created when recording starts rather than at import, so a
    session that never enables debug leaves nothing behind.
    """
    debug_log.set_enabled(on)
    if not on:
        debug_log.set_live_file(None)
        return
    if debug_log.live_file():
        return
    try:
        directory = Path(tempfile.gettempdir()) / "tomas-debug"
        directory.mkdir(parents=True, exist_ok=True)
        path = directory / f"session-{os.getpid()}.log"
        path.write_text(
            f"TOMAS debug log — session started "
            f"{time.strftime('%Y-%m-%d %H:%M:%S')}\n"
            f"Every model request and response is appended below as it "
            f"happens.\n", encoding="utf-8")
        debug_log.set_live_file(str(path))
    except OSError:
        # No live file is a degraded debug view, not a broken session:
        # /debug still reads the in-memory ring buffer.
        debug_log.set_live_file(None)


#: Assistant turns at the head of the transcript that are scaffolding, not
#: conversation. `prefill_messages` seeds a synthetic `user → assistant`
#: exchange, and its acknowledgement is not a reply to anybody — counting it
#: shifted the every-3rd-reply cadence by one, so the cap fired on the user's
#: *second* question. Set when the prefill is applied, cleared by
#: `reset_session_state` (after /clear the prefill is no longer in `messages`,
#: so subtracting for it would move the error the other way).
_synthetic_replies = 0


def _replies_so_far(messages: list) -> int:
    """How many replies the *user* has actually been given.

    Counted from the transcript rather than kept in a counter, so it survives
    /clear (which resets the conversation, and should reset the cadence with
    it) and a continued session (which should carry it on) without either case
    needing its own line of bookkeeping — minus the synthetic turns above,
    which are in the transcript but were never an answer to a question.
    """
    total = sum(1 for m in messages
                if isinstance(m, dict) and m.get("role") == "assistant")
    return max(0, total - _synthetic_replies)


def _avg_tool_tokens() -> int:
    """What one tool in the *current* pool actually costs.

    Passing a measured figure to `core.budget.resolve` is what lets a share
    mean what it says: 20% of the window is 20% whether the connected servers
    publish terse schemas or verbose ones.

    Cached on `_tool_pool_version` rather than recomputed every call: this
    used to `json.dumps` the entire tool pool unconditionally before
    `active_budget()` even checked *its own* cache, at least twice per turn,
    even though the pool only changes on MCP connect/reconnect.
    """
    global _avg_tool_tokens_cache
    if (_avg_tool_tokens_cache is not None
            and _avg_tool_tokens_cache[0] == _tool_pool_version):
        return _avg_tool_tokens_cache[1]
    pool = ALL_TOOLS or COMBINED_TOOLS or TOOLS
    value = 0 if not pool else max(1, estimate_tool_tokens(pool) // len(pool))
    _avg_tool_tokens_cache = (_tool_pool_version, value)
    return value


def active_budget():
    """The resolved budget for the current model and settings.

    Cached on everything that can change it, because this is read on the hot
    path — once per turn by the prompt builder and again by the tool selector.
    """
    global _budget_cache
    settings = budget_settings()
    window = CONTEXT_WINDOW or DEFAULT_CONTEXT_WINDOW
    ceiling = provider_tool_ceiling()
    per_tool = _avg_tool_tokens()
    key = (window, ceiling, per_tool, json.dumps(settings.to_dict(), sort_keys=True))
    if _budget_cache is not None and _budget_cache[0] == key:
        return _budget_cache[1]
    resolved = core_budget.resolve(settings, window, ceiling, per_tool)
    _budget_cache = (key, resolved)
    return resolved


def tool_ceiling() -> int:
    """How many tools this turn may send — the window's answer, not the API's."""
    try:
        return max(1, active_budget().tool_ceiling)
    except Exception:
        return provider_tool_ceiling()


def output_reserve() -> int:
    """Tokens held back for the reply.

    An explicit `AGENT_MAX_TOKENS` still wins, for the reason documented on
    `effective_max_tokens`: it is the escape hatch a reasoning model needs, and
    a budget share that silently overrode it would reinstate exactly the clamp
    that was removed. Absent that, the budget decides — 8,192 was 25% of a
    32,768-token window and nobody chose it.
    """
    if _max_tokens_was_set():
        return MAX_TOKENS
    try:
        return active_budget().output_reserve
    except Exception:
        return MAX_TOKENS


def enabled_tools(pool: Optional[list] = None) -> list:
    """The tool pool with anything the user switched off removed.

    Filtered here, at the pool, rather than at selection: a disabled tool left
    in `ALL_TOOLS` is still counted by `estimate_tool_tokens`, so the saving
    promised on the budget screen would appear nowhere in the numbers — and
    the tool would still reach the model on any turn selection ranked it well.
    """
    source = pool if pool is not None else (ALL_TOOLS or COMBINED_TOOLS or TOOLS)
    try:
        return core_budget.filter_tools(source, active_budget(), _server_of)
    except Exception:
        return list(source)


def _cap_pool(pool: list, ceiling: int) -> list:
    """The tools a turn would send from `pool`, given a ceiling.

    Not `apply_tool_cap`: that one takes *MCP* tools and prepends the built-ins
    itself, so handing it a pool that already contains them counts every
    built-in twice — which is how the budget screen first reported "18 of 12
    tools sent".
    """
    builtin_names = {t.get("name") for t in TOOLS}
    mcp_only = [t for t in pool if t.get("name") not in builtin_names]
    kept_builtins = [t for t in pool if t.get("name") in builtin_names]
    room = max(0, ceiling - len(kept_builtins))
    return kept_builtins + mcp_only[:room]


def _refresh_tool_tokens() -> int:
    """Recompute the tool-block estimate the compaction budget reserves.

    `TOOL_TOKENS` was measured once at startup from every discovered tool, so
    it neither shrank when the ceiling did nor when the user disabled a server
    — the compaction budget went on reserving for tools that were no longer
    being sent.
    """
    global TOOL_TOKENS
    try:
        sent = _cap_pool(enabled_tools(), tool_ceiling())
        TOOL_TOKENS = estimate_tool_tokens(compact_tool_schemas(sent))
    except Exception:
        pass
    return TOOL_TOKENS


def render_budget(breakdown, width: int = 22) -> list[str]:
    """The budget as lines of text. Shared by /budget and the TUI page.

    One renderer, because the two views drifting apart is how a user ends up
    being told two different numbers for the same thing — the bug that
    `_probed_capability_window` was written to stop between /status and
    /provider.
    """
    total = max(1, breakdown.window)
    lines = []
    for line in breakdown.lines:
        shown = line.tokens if line.enabled else 0
        filled = int(round(width * shown / total))
        bar = "█" * filled + "░" * (width - filled)
        if not line.enabled:
            state = f'{DIM}off{RESET}'
            bar = f'{DIM}{bar}{RESET}'
        elif line.protected:
            state = f'{GREEN}on{RESET}'
        else:
            state = 'on ' if line.toggleable else f'{DIM}—{RESET}  '
        pct = shown / total * 100
        lines.append(f'    {line.label:<22}{shown:>7,} {bar} {pct:>4.1f}%  {state}')
    filled = int(round(width * breakdown.output_reserve / total))
    lines.append(f'    {"Output reserve":<22}{breakdown.output_reserve:>7,} '
                 f'{"█" * filled + "░" * (width - filled)} '
                 f'{breakdown.output_reserve / total * 100:>4.1f}%')
    lines.append(f'    {DIM}{"─" * 58}{RESET}')
    room = breakdown.conversation_room
    colour = GREEN if breakdown.fits else RED
    lines.append(f'    {"Left for conversation":<22}{colour}{room:>7,}{RESET}'
                 f'  {room / total * 100:>4.1f}%')
    if breakdown.tools_sent and breakdown.tools_sent > tool_ceiling():
        # Not a miscount: `apply_tool_cap` keeps every built-in whatever the
        # ceiling says, because read/write/edit/run are what the agent is. A
        # ceiling below their number is therefore honoured for MCP tools only,
        # and saying so beats letting the two numbers look broken.
        lines.append(f'    {DIM}Ceiling is below the {len(TOOLS)} built-in tools, '
                     f'which are always kept.{RESET}')
    if not breakdown.fits:
        lines.append('')
        lines.append(f'  {RED}⚠{RESET}  The fixed overhead does not leave room for a '
                     f'reply. Compaction cannot')
        lines.append(f'     help — it only shrinks the conversation. Lower the tool '
                     f'ceiling or')
        lines.append(f'     switch to the economy profile.')
    return lines


def _budget_command(arg: str) -> str:
    """`/budget` — show the breakdown, or change one thing about it."""
    settings = budget_settings()
    words = arg.split()

    if words:
        head = words[0].lower()
        try:
            if head in ("auto", *core_budget.PRESETS):
                settings = dataclasses.replace(settings, profile=head)
            elif head == "tools" and len(words) > 1:
                value = None if words[1] == "auto" else max(0, int(words[1]))
                settings = dataclasses.replace(settings, tool_ceiling=value)
            elif head == "output" and len(words) > 1:
                value = None if words[1] == "auto" else max(1, int(words[1]))
                settings = dataclasses.replace(settings, output_reserve=value)
            elif head in ("on", "off") and len(words) > 1:
                key = words[1].replace("-", "_")
                if key not in core_budget.SECTION_KEYS:
                    return (f'  {RED}Unknown section {words[1]!r}.{RESET} '
                            f'{DIM}One of: {", ".join(core_budget.SECTION_KEYS)}{RESET}')
                sections = dict(settings.sections or {})
                sections[key] = (head == "on")
                settings = dataclasses.replace(settings, sections=sections)
            else:
                return (f'  {RED}Unrecognised: {arg!r}{RESET}\n'
                        f'  {DIM}/budget [auto|economy|balanced|full] · '
                        f'/budget tools <N|auto> · /budget output <N|auto> ·\n'
                        f'  /budget on|off <section>{RESET}')
        except ValueError:
            return f'  {RED}Expected a number: {arg!r}{RESET}'
        save_budget_settings(settings)

    breakdown = budget_breakdown()
    budget = active_budget()
    auto = " (auto)" if settings.profile not in core_budget.PRESETS else ""
    out = [
        f'  {BOLD}Context budget{RESET} {DIM}·{RESET} {budget.profile.label}{auto} '
        f'{DIM}·{RESET} {breakdown.window:,} token window',
        f'  {DIM}{budget.profile.summary}{RESET}',
        f'  {DIM}{"─" * 58}{RESET}',
    ]
    out += render_budget(breakdown)
    out += [
        '',
        f'  {DIM}Tools:{RESET} {breakdown.tools_sent} of {breakdown.tools_available} '
        f'sent per turn'
        + (f' {DIM}(ceiling set by hand){RESET}' if budget.tool_ceiling_is_manual
           else f' {DIM}(ceiling derived from the window){RESET}'),
    ]
    if budget.disabled_tools or budget.disabled_servers:
        out.append(f'  {DIM}Disabled:{RESET} '
                   f'{len(budget.disabled_tools)} tool(s), '
                   f'{len(budget.disabled_servers)} server(s)')
    out.append(f'  {DIM}Change with{RESET} {CYAN}/budget economy{RESET}{DIM},{RESET} '
               f'{CYAN}/budget tools 12{RESET}{DIM},{RESET} '
               f'{CYAN}/budget off skills_catalogue{RESET}')
    return '\n'.join(out)


def budget_breakdown(user_message: str = ""):
    """Measure where this turn's window actually goes.

    Every number is taken from the same functions the turn itself uses — the
    prompt builder, the tool selector, the compaction planner — rather than
    re-derived. A budget screen that estimated independently would drift from
    what is really sent, and then be worse than no screen at all.
    """
    window = CONTEXT_WINDOW or DEFAULT_CONTEXT_WINDOW
    budget = active_budget()
    prose = CHARS_PER_TOKEN_PROSE

    pool = enabled_tools()
    sent = _cap_pool(pool, tool_ceiling())
    tool_tokens = estimate_tool_tokens(compact_tool_schemas(sent))

    def _tok(text: str) -> int:
        return int(len(text or "") / prose)

    sections: dict = {}
    if budget.allows("instructions"):
        sections["instructions"] = _tok(build_instructions_section(PROJECT_DIR))
    if budget.allows("skills_catalogue"):
        sections["skills_catalogue"] = _tok(
            build_skills_section(max_chars=MAX_SKILLS_CHARS))
    for key, produce in (
        ("standing_rules", lambda: learning.directives_for_prompt()),
        ("learned_facts", lambda: learning.recall(user_message, k=5)),
        ("triggered_skills", lambda: build_triggered_skills(
            user_message, MAX_TRIGGERED_SKILL_CHARS)),
    ):
        if not budget.allows(key):
            continue
        try:
            sections[key] = _tok(produce())
        except Exception:
            sections[key] = 0

    plan = core_context.compaction_plan(
        used_tokens=0, window_tokens=window,
        reserve_tokens=0, fit_fraction=COMPACTION_THRESHOLD,
        cost_limit=COMPACTION_COST_LIMIT)

    return core_budget.build_breakdown(
        window, budget,
        base_tokens=_tok(BASE_PROMPT),
        environment_tokens=_tok(_environment_section()),
        section_tokens=sections,
        tool_tokens=tool_tokens,
        tools_sent=len(sent),
        tools_available=len(ALL_TOOLS or COMBINED_TOOLS or TOOLS),
        compaction_trigger=plan.trigger)


# ── ANSI color constants for the chat UI ──
RESET = '\033[0m'
BOLD = '\033[1m'
DIM = '\033[2m'
GREEN = '\033[92m'
CYAN = '\033[96m'
YELLOW = '\033[93m'
BLUE = '\033[94m'
MAGENTA = '\033[95m'
RED = '\033[91m'
GRAY = '\033[90m'
BOLD_OFF = '\033[22m'

# ---------------------------------------------------------------------------
# Tool definitions (schema sent to the model)
# ---------------------------------------------------------------------------

#: What `tab_act` accepts, spelled out here rather than imported from
#: `core.browser.ACTIONS`.
#:
#: The import would be the better dependency and costs 72 ms of `asyncio` at
#: module load — paid by every `agent_cli` menu, on every start, for a feature
#: most sessions never touch. agent.py already defers `asyncio` twice for the
#: same reason (see handle_fetch_url_with_browser), so the browser handlers
#: import their engine when they run. The duplication is the price, and
#: `test_browser_control.py` asserts the two lists agree.
#: What `doc_edit` accepts. Duplicated from `core.office.EDIT_ACTIONS` for
#: the same reason as ACTIONS_FOR_SCHEMA below: importing it would pull
#: pywin32 in at module load, and `agent_cli`'s menus never touch Word.
#: `test_office_control.py` asserts the two lists agree.
EDIT_ACTIONS_FOR_SCHEMA = ("replace", "insert_after", "insert_before",
                           "delete", "style", "find_replace",
                           "insert_equation", "equation")

ACTIONS_FOR_SCHEMA = ("click", "double_click", "type", "press", "select",
                      "hover", "check", "uncheck", "clear", "scroll")

TOOLS: list[dict] = [
    {
        "name": "read_file",
        "description": "Read a file from the filesystem. Returns contents with line numbers. Also extracts text from .pdf, .docx, .pptx, and .xlsx files -- no separate tool needed for those.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Absolute or project-relative path"},
                "offset": {"type": "integer", "description": "Line number to start reading from (1-indexed)"},
                "limit": {"type": "integer", "description": "Max lines to read. Defaults to 2000."},
            },
            "required": ["file_path"],
        },
    },
    {
        "name": "write_file",
        "description": "Write text content to a file. Creates the file (and parent dirs) if needed. Text formats only -- it cannot produce .docx/.pdf/.xlsx/.pptx, which are structured containers; build those with the word-docs MCP tools or a library.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Absolute or project-relative path"},
                "content": {"type": "string", "description": "Full file content to write"},
            },
            "required": ["file_path", "content"],
        },
    },
    {
        "name": "edit_file",
        "description": "Replace a string in a file. old_string must appear exactly once, unless replace_all is true.",
        "input_schema": {
            "type": "object",
            "properties": {
                "file_path": {"type": "string", "description": "Absolute or project-relative path"},
                "old_string": {"type": "string", "description": "Exact text to find"},
                "new_string": {"type": "string", "description": "Replacement text"},
                "replace_all": {"type": "boolean", "description": "Replace every occurrence instead of requiring a unique match. Use this for mechanical substitutions rather than issuing one call per site."},
            },
            "required": ["file_path", "old_string", "new_string"],
        },
    },
    {
        "name": "list_files",
        "description": "List files in a directory (non-recursive).",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Directory to list. Defaults to project root."},
            },
            "required": [],
        },
    },
    {
        "name": "run_command",
        "description": "Execute a shell command. Returns '[exit N — ok|FAILED]' followed by stdout and any stderr, so you never need to append '2>&1' or infer success from the text. For a process that stays running (a dev server, a watcher), launch it detached with Windows `start /b ...` (or a trailing `&`) rather than waiting on it directly — the call returns immediately once it's launched, and its output is not captured. Redirect its output to a log file, then use check_progress on that file instead of polling with more run_command calls.",
        "input_schema": {
            "type": "object",
            "properties": {
                "command": {"type": "string", "description": "Shell command to run"},
                "timeout": {"type": "integer", "description": "Timeout in seconds. Default 120."},
            },
            "required": ["command"],
        },
    },
    {
        "name": "check_progress",
        "description": "Check whether a background job is still making progress, without polling by hand. Point it at the log file a backgrounded command writes to (or a growing output directory like node_modules) — it reports the size, how much it changed since your last check on this same path, and flags a likely stall (e.g. from a full disk) instead of you re-running dir/type/tasklist in a loop.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Log file or directory to check for growth since the last call on this same path."},
            },
            "required": ["path"],
        },
    },
    {
        "name": "search_code",
        "description": "Search for a regex pattern across files. Returns the true match count; page through large result sets with offset.",
        "input_schema": {
            "type": "object",
            "properties": {
                "pattern": {"type": "string", "description": "Regex pattern"},
                "path": {"type": "string", "description": "Directory or single file to search. Defaults to project root."},
                "file_glob": {"type": "string", "description": "File pattern filter, e.g. '*.py'. Ignored when path is a file."},
                "offset": {"type": "integer", "description": "Skip this many matches. Default 0."},
            },
            "required": ["pattern"],
        },
    },
    {
        "name": "save_memory",
        "description": (
            "Persist something the user asked to be remembered, for future "
            "sessions. It is routed to the right store automatically — leave "
            "`store` unset unless you have a reason. The result says which "
            "store was chosen and why."),
        "input_schema": {
            "type": "object",
            "properties": {
                "key": {"type": "string", "description": "Short identifier (kebab-case)"},
                "description": {"type": "string", "description": "One-line summary for the index"},
                "content": {"type": "string", "description": "Full memory content"},
                "store": {
                    "type": "string",
                    "enum": list(MEMORY_STORES),
                    "description": (
                        "Optional override. instruction = permanent identity "
                        "(cached, free per turn); rule = unconditional, "
                        "injected every turn, max 10; fact = conditional "
                        "preference, retrieved by relevance; note = long "
                        "reference material, not in the prompt."),
                },
            },
            "required": ["key", "description", "content"],
        },
    },
    {
        "name": "fetch_url",
        "description": "Fetch a URL (HTTP/HTTPS). HTML is returned as readable text with scripts and styles removed. Retries once without TLS verification if the certificate cannot be validated, and labels the result. PDFs are not decoded — it says how to read them instead.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "URL to fetch"},
                "timeout": {"type": "integer", "description": "Timeout in seconds. Default 15."},
                "max_size": {"type": "integer", "description": "Max bytes to read. Default 200000; the response is clipped, not refused."},
                "raw": {"type": "boolean", "description": "Return HTML markup instead of extracted text. Default false — only set this when you need the tags themselves."},
            },
            "required": ["url"],
        },
    },
    {
        "name": "fetch_url_with_browser",
        "description": "Fetch content from a URL using a headless browser (Playwright). Supports JavaScript rendering, SPAs, and dynamic content. Returns the page content as text.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "URL to fetch"},
                "wait_for": {"type": "string", "description": "Selector to wait for before extracting content (optional)"},
                "timeout": {"type": "integer", "description": "Timeout in seconds. Default 30."},
                "max_size": {"type": "integer", "description": "Max response size in bytes. Default 50000000."},
                "screenshot": {"type": "boolean", "description": "Whether to take a screenshot. Default false."},
            },
            "required": ["url"],
        },
    },
    {
        "name": "search_web",
        "description": "Search the internet for current information on any topic. Returns a list of search results with titles, snippets, and URLs. Use this when you need up-to-date information from the web (news, weather, facts, etc.).",
        "input_schema": {
            "type": "object",
            "properties": {
                "query": {"type": "string", "description": "The search query (e.g. 'weather in Vinnytsia Ukraine' or 'latest AI news 2026')"},
                "max_results": {"type": "integer", "description": "Max results to return. Default 5."},
            },
            "required": ["query"],
        },
    },
    # ── Browser control ──
    # These attach to the browser the user already has open (see
    # core/browser.py); they are not the headless fetcher above. The
    # descriptions say so, because a model that picks `tab_read` to read
    # an arbitrary URL will get whatever tab happens to be in front.
    {
        "name": "tab_list",
        "description": "List the tabs of the user's OWN running browser and choose which one to work in. This attaches to a real browser with the user's logins -- it is not the headless fetcher. Call with no arguments to list tabs and see which one is attached. Use start_browser=true if no debuggable browser is running yet.",
        "input_schema": {
            "type": "object",
            "properties": {
                "select": {"type": "integer", "description": "Attach to this tab number (from a previous listing)."},
                "new_tab": {"type": "string", "description": "Open a new tab at this URL and attach to it. Prefer this over navigating the user's current tab away from what they were doing."},
                "start_browser": {"type": "boolean", "description": "Start a debuggable browser if none is running. It opens with its own profile, so the user signs in once."},
            },
            "required": [],
        },
    },
    {
        "name": "tab_snapshot",
        "description": "List the visible interactive elements of the attached tab as numbered refs ([e7] button \"Send\"). Call this before acting: tab_act addresses elements by ref, and refs come from here. Re-snapshot after any navigation. Covers the main frame only -- controls inside an iframe are not listed.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "tab_read",
        "description": "Read the visible text of the attached tab, or of one element by ref. Use this to understand page content; use tab_snapshot to find things to click.",
        "input_schema": {
            "type": "object",
            "properties": {
                "ref": {"type": "string", "description": "Read only this element (from tab_snapshot). Omit to read the whole page."},
                "max_chars": {"type": "integer", "description": "Clip at this many characters. Default 20000."},
            },
            "required": [],
        },
    },
    {
        "name": "tab_navigate",
        "description": "Point the attached tab at a URL, or move through its history. This changes what the user is looking at -- to avoid disturbing their current tab, open a new one with tab_list new_tab instead.",
        "input_schema": {
            "type": "object",
            "properties": {
                "url": {"type": "string", "description": "URL to open in the attached tab."},
                "action": {"type": "string", "enum": ["back", "forward", "reload"], "description": "History action, instead of a url."},
            },
            "required": [],
        },
    },
    {
        "name": "tab_act",
        "description": "Act in the attached tab as the user would: click, type, press a key, choose from a dropdown, hover, tick a box, scroll. Elements are addressed by ref from tab_snapshot. This runs in the user's real, logged-in browser, so an action here can send, buy, or delete for real -- act only on what was asked for.",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": list(ACTIONS_FOR_SCHEMA), "description": "What to do."},
                "ref": {"type": "string", "description": "Target element, e.g. 'e7', from tab_snapshot. Required for everything except a bare press or scroll."},
                "text": {"type": "string", "description": "Text for action=type; pixels for action=scroll."},
                "key": {"type": "string", "description": "Key for action=press, e.g. 'Enter', 'Escape', 'Control+A'."},
                "option": {"type": "string", "description": "Option label or value for action=select."},
                "submit": {"type": "boolean", "description": "Press Enter after typing. Default false."},
                "clear_first": {"type": "boolean", "description": "Replace the field's contents rather than appending. Default true."},
            },
            "required": ["action"],
        },
    },
    {
        "name": "tab_screenshot",
        "description": "Save a PNG of the attached tab and return its path. Read the file afterwards to look at it. Use when the layout matters and the text outline does not answer the question.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Where to save it. Defaults to a timestamped file under ~/.tomas/screenshots/."},
                "full_page": {"type": "boolean", "description": "Capture the whole scrollable page rather than the viewport. Default false."},
            },
            "required": [],
        },
    },
    # ── Live Office documents ──
    # These attach to the Word the user already has open (see core/office.py).
    # The descriptions say so, because a model that reaches for these to build
    # a document from scratch will edit whatever the user is looking at.
    {
        "name": "doc_list",
        "description": "List the documents open in the user's OWN running Word and choose which to work in. This edits the live document in their visible window -- it is not a file tool, and it is not for building a document from scratch (use write_file or the word-docs MCP server for a closed file). Call with no arguments to list.",
        "input_schema": {
            "type": "object",
            "properties": {
                "select": {"type": "integer", "description": "Attach to this document number from a previous listing."},
                "new_document": {"type": "boolean", "description": "Create a new empty document and attach to it."},
                "start_app": {"type": "boolean", "description": "Start Word if it is not running."},
            },
            "required": [],
        },
    },
    {
        "name": "doc_outline",
        "description": "List the paragraphs of the attached document as numbered refs with their styles ([p12] Heading 1 \"Introduction\"). Call before editing by ref. Refs are void after any edit that adds or removes a paragraph, and after the user types -- doc_edit refuses rather than guessing.",
        "input_schema": {"type": "object", "properties": {}, "required": []},
    },
    {
        "name": "doc_read",
        "description": "Read the text of the attached document, or of one paragraph by ref.",
        "input_schema": {
            "type": "object",
            "properties": {
                "ref": {"type": "string", "description": "Read only this paragraph, e.g. 'p3'. Omit for the whole document."},
                "max_chars": {"type": "integer", "description": "Clip at this many characters. Default 20000."},
            },
            "required": [],
        },
    },
    {
        "name": "doc_find",
        "description": "Find paragraphs containing text and return their refs. Preferred over reading the whole outline when you know what you are looking for, and the refs it returns are fresh.",
        "input_schema": {
            "type": "object",
            "properties": {
                "text": {"type": "string", "description": "Literal text to look for (case-insensitive)."},
                "max_hits": {"type": "integer", "description": "Stop after this many. Default 20."},
            },
            "required": ["text"],
        },
    },
    {
        "name": "doc_edit",
        "description": "Edit the document in the user's Word window. Prefer action=find_replace, which needs no ref. Equations are real Word objects, not text: rewrite one with action=equation ref=eqN, add one with action=insert_equation. Write maths in linear form -- 'E = \u222b_0^T P(t)dt' or 'x^2+y^2=z^2'; LaTeX-style \\int, \\sum, \\le are translated using Word's own table, but \\frac{a}{b} is not -- write a/b. Every edit is undoable with Ctrl+Z. Editing does NOT save -- use doc_save.",
        "input_schema": {
            "type": "object",
            "properties": {
                "action": {"type": "string", "enum": list(EDIT_ACTIONS_FOR_SCHEMA), "description": "What to do."},
                "ref": {"type": "string", "description": "Target: a paragraph like 'p3', or an equation like 'eq1' for action=equation. Required for everything except find_replace and a bare insert_equation."},
                "text": {"type": "string", "description": "New text (replace/insert_after/insert_before), or the replacement (find_replace)."},
                "find": {"type": "string", "description": "Text to search for, with action=find_replace."},
                "style": {"type": "string", "description": "Style name for action=style, e.g. 'Heading 1'."},
            },
            "required": ["action"],
        },
    },
    {
        "name": "doc_save",
        "description": "Save the attached document. Separate from doc_edit on purpose: saving is the one action Ctrl+Z cannot undo, so it is never implicit.",
        "input_schema": {
            "type": "object",
            "properties": {
                "path": {"type": "string", "description": "Save as this path instead of saving in place. Required if the document has never been saved."},
            },
            "required": [],
        },
    },
    {
        "name": "read_mcp_resource",
        "description": (
            "Read a resource exposed by a connected MCP server (a file, "
            "database row, document, or similar). Call with no arguments to "
            "list what is available. Resources are read-only data the server "
            "publishes; they are not tools."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "uri": {"type": "string", "description": "Resource URI. Omit to list all available resources."},
                "server": {"type": "string", "description": "MCP server name. Only needed when two servers publish the same URI."},
            },
        },
    },
    {
        "name": "ask_user_question",
        "description": (
            "Ask the user one or more multiple-choice questions and block "
            "until they answer, through an interactive arrow-key picker — the "
            "same mechanism Claude Code uses. Use this when a task is "
            "genuinely ambiguous (more than one reasonable interpretation, a "
            "choice with real consequences, missing information only the "
            "user has) instead of guessing and possibly doing the wrong "
            "thing. Do not use it for questions you could answer yourself by "
            "reading the code. Each question gets 2-4 short, mutually "
            "exclusive options (the user can always type a custom answer "
            "instead) plus an optional multiSelect flag for 'choose any that "
            "apply'. Returns the user's answer(s) for every question."
        ),
        "input_schema": {
            "type": "object",
            "properties": {
                "questions": {
                    "type": "array",
                    "minItems": 1,
                    "maxItems": 4,
                    "items": {
                        "type": "object",
                        "properties": {
                            "question": {"type": "string", "description": "The complete question to ask, clear and specific."},
                            "header": {"type": "string", "description": "Very short label for this question (max ~12 chars), e.g. 'Auth method'."},
                            "multiSelect": {"type": "boolean", "description": "True if the user may pick more than one option. Default false."},
                            "options": {
                                "type": "array",
                                "minItems": 2,
                                "maxItems": 4,
                                "items": {
                                    "type": "object",
                                    "properties": {
                                        "label": {"type": "string", "description": "Short (1-5 word) display text for this choice."},
                                        "description": {"type": "string", "description": "Optional one-line explanation of what picking this means."},
                                    },
                                    "required": ["label"],
                                },
                            },
                        },
                        "required": ["question", "options"],
                    },
                },
            },
            "required": ["questions"],
        },
    },
]

# ---------------------------------------------------------------------------
# Risk tiers for the permission system
# ---------------------------------------------------------------------------

RISK_LEVELS: dict[str, str] = {
    "read_file": "low",
    "list_files": "low",
    "search_code": "low",
    "edit_file": "medium",
    "write_file": "medium",
    "save_memory": "low",
    "run_command": "high",
    "check_progress": "low",
    "fetch_url": "low",
    "fetch_url_with_browser": "medium",
    "search_web": "low",
    # Browser control acts in the user's *own* signed-in browser, which is a
    # different proposition from the headless fetcher above. Reading is cheap
    # and reversible; `tab_act` is "high" for the same reason
    # `run_command` is — a click there can send the email, place the order or
    # delete the repo, and it does so as the user, with their session. There
    # is no sandbox to fall back on, so the permission prompt is the control.
    # `tab_list` is refined by params in `risk_for`: listing reads,
    # starting a browser launches a process.
    "tab_list": "low",
    "tab_snapshot": "low",
    "tab_read": "low",
    "tab_navigate": "medium",
    "tab_act": "high",
    "tab_screenshot": "low",
    # Live Office documents. Editing is "medium" where `tab_act` is "high",
    # and the difference is measured rather than assumed: Phase 0 showed one
    # Ctrl+Z reverses one COM edit, so a wrong edit is recoverable by the user
    # with a keystroke they already know. A browser click is not. `doc_save`
    # is "medium" for the opposite reason -- it is the one action Undo cannot
    # reverse. `doc_list` is refined by params in `risk_for`.
    "doc_list": "low",
    "doc_outline": "low",
    "doc_read": "low",
    "doc_find": "low",
    "doc_edit": "medium",
    "doc_save": "medium",
    "read_mcp_resource": "low",
    # "none", not "low": asking the user a question has no side effects at
    # all, and gating it behind "approve this tool call?" would put a
    # permission prompt in front of the very interaction that *is* the
    # human-in-the-loop control. See AgentState.needs_permission.
    "ask_user_question": "none",
}

#: Built-in tools that may run concurrently with one another.
#:
#: An allowlist, not a risk tier, because the two questions are different: a
#: call can be perfectly safe to *approve* and still unsafe to run twice at
#: once. Everything here reads and returns; nothing here shares a handle.
#:
#: Deliberately excluded, each for its own reason:
#:   run_command   — classified `low` for read-only commands, but it routes
#:                   multi-line payloads through a temp dir and reads the
#:                   module-level interrupt; two at once share both.
#:   search_web    — drives a headless browser. Two launches is two Chromes.
#:   read_mcp_resource — one stdio pipe per server carrying JSON-RPC; two
#:                   concurrent calls would interleave frames on it.
#:   tab_*        — one attached page and one ref map, shared by all six.
#:                   Even the reads: a snapshot taken while another call is
#:                   navigating would number elements on a page that no
#:                   longer exists, and the refs are what a later click
#:                   trusts.
#:   doc_*         — one COM thread, one attached document, one outline
#:                   fingerprint. Two calls at once would interleave on the
#:                   same Word instance, and the fingerprint guard would be
#:                   comparing against an edit still in flight.
#:   save_memory   — writes.
#:   ask_user_question — is a conversation, and those are sequential.
PARALLEL_SAFE_TOOLS = frozenset({
    "read_file", "list_files", "search_code", "fetch_url",
})

#: Escape hatch. Set TOMAS_PARALLEL_TOOLS=0 to go back to strictly sequential
#: execution without editing anything.
PARALLEL_TOOLS_ENABLED = os.environ.get("TOMAS_PARALLEL_TOOLS", "1") != "0"


def parallel_safe(name: str, params: Optional[dict] = None) -> bool:
    """May this call overlap with the others in its batch?

    Takes the params for the same reason `risk_for` does — so a future tool
    whose safety depends on its arguments has somewhere to say so — and
    ignores them for the current allowlist, which is safe by tool identity.
    """
    return PARALLEL_TOOLS_ENABLED and name in PARALLEL_SAFE_TOOLS


#: The tools whose risk tier strict mode leaves alone — they are this agent's
#: own, and their tiers are deliberate. Everything else (i.e. MCP) is reset to
#: "high". Previously spelled out as a literal set in two places, which drifted:
#: `read_mcp_resource` was in RISK_LEVELS but in neither copy.
BUILTIN_TOOL_NAMES = frozenset(RISK_LEVELS)

# Commands that only read. `run_command` is the most-used tool in the corpus
# (72 of 209 calls) and every one of them blocked on a prompt, because
# `git status` and `rm -rf` shared a single tier.
READONLY_CMD = re.compile(
    r'^\s*(git\s+(status|log|diff|show|branch)\b'
    r'|(?:[\w.\\/:-]*[\\/])?python(?:\.exe)?\s+-m\s+(?:unittest|pytest)\b'
    r'|pytest\b|dir\b|ls\b|type\b|cat\b|echo\b|where\b|which\b|findstr\b'
    r'|pip\s+(list|show|freeze)\b)',
    re.I,
)

# Any of these turns a command into something a single-token classifier cannot
# reason about. The corpus contains `del x && type y` — a delete wearing a
# read's clothes — and `dir /b n* & echo --- & type n`, where cmd.exe treats a
# lone `&` as a separator. Chaining is not classified; it is high.
_CMD_SEPARATORS = ("&&", "||", ";", "|", ">", "<", "&", "`", "$(")


def risk_for(name: str, params: Optional[dict] = None) -> str:
    """Risk tier for a specific call, not just a tool name."""
    if name == "run_command":
        cmd = (params or {}).get("command", "")
        if any(sep in cmd for sep in _CMD_SEPARATORS):
            return "high"
        if re.search(r'\b(del|rm|rmdir|erase|move|ren|format|curl|wget)\b', cmd, re.I):
            return "high"
        return "low" if READONLY_CMD.match(cmd) else "high"
    if name == "doc_list":
        # Listing and switching read. Starting Word, or adding a document to
        # the user's session, is a change to what is on their screen.
        params = params or {}
        if params.get("start_app") or params.get("new_document"):
            return "medium"
        return "low"
    if name == "tab_list":
        # Listing tabs and switching between them reads. Starting a browser
        # spawns a detached process, which is not something to auto-approve
        # on the strength of the tool's name.
        params = params or {}
        if params.get("start_browser"):
            return "medium"
        if params.get("new_tab"):
            return "medium"
        return "low"
    return RISK_LEVELS.get(name, "high")

# Patterns that are always blocked from run_command
BLOCKED_PATTERNS = ["rm -rf /", "mkfs", "> /dev/sd", "dd if=/dev/zero", ":(){:|:&};:"]

# Ceiling on a single read_file result. The line limit alone is not a size
# limit — one 2000-line read put 45 KB into the context in a single result.
MAX_READ_FILE_CHARS = 20000

# ---------------------------------------------------------------------------
# Path helpers
# ---------------------------------------------------------------------------

def _resolve(path: str) -> Path:
    # ~/.tomas is documented as readable, but Path() treats "~" as an ordinary
    # directory name: "~/.tomas/memory/MEMORY.md" resolved to
    # PROJECT_DIR/~/.tomas/... and came back "file not found", which the model
    # only recovered from by guessing the absolute path.
    # Only the home form is expanded. expanduser() maps "~weird" to
    # C:\Users\weird without complaint, so a file genuinely named "~weird"
    # must keep resolving inside the project.
    s = _repair_drive(str(path))
    p = Path(s)
    if s[:1] == "~" and (len(s) == 1 or s[1] in "/\\"):
        p = p.expanduser()
    if not p.is_absolute():
        p = PROJECT_DIR / p
    return p.resolve()


def _repair_drive(s: str) -> str:
    """Put back a drive letter the path lost on its way here.

    A path pasted or dragged into the prompt can arrive a character short:
    measured, three sessions were handed `:\\Github\\Agent-For-TOM\\...` with
    the leading `C` gone, and each spent between three and seven tool calls
    hunting for a file that was exactly where it said it was. `:` cannot
    begin a relative path on Windows — it is not a legal filename character —
    so a string starting with it is unambiguously a mangled absolute path and
    repairing it cannot shadow a real file. A leading separator with no drive
    is the same case one character further along.

    Only ever *offered*: the repaired path is returned as a string and still
    has to exist. Nothing here decides a file is present.
    """
    if not s:
        return s
    drive = PROJECT_DIR.drive           # 'C:' on Windows, '' elsewhere
    if not drive:
        return s
    if s[0] == ":" and len(s) > 1 and s[1] in "/\\":
        # Unconditional: `:` is not a legal filename character, so this string
        # cannot name anything relative. Repairing even when the result does
        # not exist is what makes the *error* readable — left alone it was
        # joined to the project directory and reported as
        # `C:\...\Agent-For-TOM\:\Github\...`, which names no file anyone typed.
        return drive[0] + s
    if s[0] in "/\\" and not s.startswith("//") and not s.startswith("\\\\"):
        candidate = drive + s
        if Path(candidate).exists():
            return candidate
    return s


#: How many same-named files a not-found error may list. Enough to disambiguate
#: a real duplicate, few enough that a common name does not fill the reply.
_NOT_FOUND_SUGGESTIONS = 5


def _not_found_error(path: Path) -> str:
    """"Not found" plus where it actually is, when that is answerable.

    A bare "file not found" is a dead end the model can only leave by
    guessing: measured across three sessions handed a slightly wrong sample
    path, the recoveries cost 3, 4 and 7 tool calls — `dir /s /b`, `list_files`
    on three directories, a `search_code`, and in one case two more reads of
    a path that had already failed once. The name is almost always right and
    only the directory is wrong, so one bounded search by basename ends it in
    the same call that reported the problem.
    """
    suggestions: list[Path] = []
    name = path.name
    if name and not any(c in name for c in "*?"):
        try:
            for found in PROJECT_DIR.rglob(name):
                if found.is_file():
                    suggestions.append(found)
                    if len(suggestions) >= _NOT_FOUND_SUGGESTIONS:
                        break
        except OSError:
            pass
    if not suggestions:
        return f"Error: file not found: {path}"
    if len(suggestions) == 1:
        return (f"Error: file not found: {path}\n"
                f"There is one file named {name} in the project: "
                f"{suggestions[0]}")
    listed = "\n".join(f"  {s}" for s in suggestions)
    return (f"Error: file not found: {path}\n"
            f"Files named {name} in the project:\n{listed}")

# All user state lives under ~/.tomas (sessions, notes, memory, learned
# skills). Locking the agent out of it does not make anything safer — it just
# pushes the same read through `run_command`, which is a higher risk tier.
TOMAS_HOME = (Path.home() / ".tomas").resolve()


def _within(p: Path, root: Path) -> bool:
    return p == root or root in p.parents


def _safe(p: Path, write: bool = False) -> bool:
    """Project dir is read-write. ~/.tomas is read-only: it is written through
    the typed APIs (save_memory, self_notes, session_manager) that own each
    file's schema, never by hand."""
    if _within(p, PROJECT_DIR):
        return True
    return not write and _within(p, TOMAS_HOME)


def _outside_project_error(path: Path, write: bool = False) -> str:
    """Say which rule was hit *and where to go instead*.

    Naming only the rule sent the model looking for a second forbidden
    location: told ~/.tomas was read-only it tried the system temp directory,
    was refused again, and spent two tool calls before landing in the project.
    An error that names the writable scratch path ends that search in one.
    """
    if write and _within(path, TOMAS_HOME):
        return (f"Error: {path} is under ~/.tomas, which is read-only. "
                f"For a throwaway helper script use {SCRATCH_DIR}; for durable "
                f"notes use save_memory or the self_notes API.")
    if write:
        return (f"Error: path outside project: {path}. Writes are only allowed "
                f"under {PROJECT_DIR}. For a throwaway helper script use "
                f"{SCRATCH_DIR} — the system temp directory is not writable "
                f"from here.")
    return (f"Error: path outside project: {path}. "
            f"Readable roots: {PROJECT_DIR} (read-write), {TOMAS_HOME} (read-only).")

# ---------------------------------------------------------------------------
# Tool handlers
# ---------------------------------------------------------------------------

# Opening these in text mode (the default read_file path) just decodes their
# raw binary bytes as UTF-8 with replacement characters -- the PDF header and
# compressed font streams come back as "%PDF-1.5" and garbage, not the
# document's actual text. Each extractor turns the format into plain text
# read_file can line-number and paginate like any other file.

def _extract_pdf_text(path: Path) -> str:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    pages = [f"--- page {i} ---\n{page.extract_text() or ''}"
             for i, page in enumerate(reader.pages, 1)]
    return "\n".join(pages)


def _extract_docx_text(path: Path) -> str:
    import docx
    document = docx.Document(str(path))
    parts = [p.text for p in document.paragraphs]
    for table in document.tables:
        for row in table.rows:
            parts.append(" | ".join(cell.text for cell in row.cells))
    return "\n".join(parts)


def _extract_pptx_text(path: Path) -> str:
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- slide {i} ---")
        for shape in slide.shapes:
            if shape.has_text_frame:
                parts.append(shape.text_frame.text)
    return "\n".join(parts)


def _extract_xlsx_text(path: Path) -> str:
    import openpyxl
    wb = openpyxl.load_workbook(str(path), data_only=True)
    parts: list[str] = []
    for ws in wb.worksheets:
        parts.append(f"--- sheet: {ws.title} ---")
        for row in ws.iter_rows(values_only=True):
            parts.append("\t".join("" if c is None else str(c) for c in row))
    return "\n".join(parts)


_DOCUMENT_EXTRACTORS = {
    ".pdf": _extract_pdf_text,
    ".docx": _extract_docx_text,
    ".pptx": _extract_pptx_text,
    ".xlsx": _extract_xlsx_text,
}

# Formats with no text in them at all. Decoding these as UTF-8 with
# replacement produces pages of "�PNG IHDR ... IDATx^" — which is not an
# error, so nothing stops the model reading a second one. Measured: a session
# trying to see a sample document's layout read two page renders this way,
# spent two tool calls on the mojibake, and concluded it had seen the page.
_UNREADABLE_AS_TEXT = {
    ".png": "an image", ".jpg": "an image", ".jpeg": "an image",
    ".gif": "an image", ".bmp": "an image", ".webp": "an image",
    ".tif": "an image", ".tiff": "an image", ".ico": "an image",
    ".zip": "an archive", ".gz": "an archive", ".tar": "an archive",
    ".7z": "an archive", ".rar": "an archive",
    ".exe": "a binary", ".dll": "a binary", ".so": "a binary",
    ".pyc": "compiled bytecode", ".pyd": "a binary",
    ".mp3": "audio", ".wav": "audio", ".mp4": "video", ".mov": "video",
    ".woff": "a font", ".woff2": "a font", ".ttf": "a font", ".otf": "a font",
}


def _unreadable_error(path: Path, kind: str) -> str:
    """Say what it is and which tool can actually look at it."""
    if kind == "an image":
        how = ("Attach it to the conversation to have it looked at, or "
               "measure it with Pillow/PyMuPDF through run_command. To see a "
               "*document's* layout, measure the document itself — "
               "skills/document-style-match/run.py measure — not a picture "
               "of it")
    elif kind == "an archive":
        how = "list or extract it with run_command, then read the members"
    else:
        how = "inspect it with run_command if you need its bytes"
    return (f"Error: {path.name} is {kind}, not text. Reading it here would "
            f"return its raw bytes as replacement characters and nothing "
            f"else. {how}.")


#: Prepended to every extracted document. The extraction is genuinely useful
#: and stays — but a session read a PDF this way, saw its words in the right
#: order, and rebuilt it on US Letter paper at 1.15 line spacing with its
#: headings flush left, because none of that is in the text and nothing said
#: so. One line is cheaper than the six defects.
_EXTRACTED_TEXT_BANNER = (
    "[text only — extracted from {kind}. Page size, margins, fonts, sizes, "
    "alignment, bold and images are NOT in this output and cannot be judged "
    "from it. To reproduce this document's formatting, measure it: "
    "skills/document-style-match/run.py measure]"
)


def handle_read_file(params: dict) -> str:
    path = _resolve(params["file_path"])
    if not _safe(path):
        return _outside_project_error(path)
    if not path.exists():
        return _not_found_error(path)
    suffix = path.suffix.lower()
    unreadable = _UNREADABLE_AS_TEXT.get(suffix)
    if unreadable:
        return _unreadable_error(path, unreadable)
    offset = max(0, int(params.get("offset", 1)) - 1)
    limit = int(params.get("limit", 2000))
    banner = ""
    if suffix in _DOCUMENT_EXTRACTORS:
        try:
            text = _DOCUMENT_EXTRACTORS[suffix](path)
        except ImportError as e:
            return (f"Error: reading {suffix} files needs an extra package "
                     f"that isn't installed ({e}). Run: pip install -r requirements.txt")
        except Exception as e:
            return f"Error: could not extract text from {path.name}: {e}"
        # Only on the first page of a paginated read: repeating it at
        # offset=2000 would be noise, and the caller has already seen it.
        if offset == 0:
            banner = _EXTRACTED_TEXT_BANNER.format(kind=f"a {suffix[1:]}") + "\n"
        lines = [line + "\n" for line in text.split("\n")]
    else:
        with path.open("r", encoding="utf-8", errors="replace") as f:
            lines = f.readlines()
    selected = lines[offset:offset + limit]
    if not selected:
        return "(empty file)"
    # A line limit alone is not a size limit: a 2000-line read of a wide file
    # put 45 KB into the context in one tool result. Cap by characters too,
    # and say where to resume so the truncation is recoverable.
    out: list[str] = []
    used = 0
    for i, line in enumerate(selected):
        entry = f"{i + offset + 1:6}\t{line}"
        if used + len(entry) > MAX_READ_FILE_CHARS:
            next_line = i + offset + 1
            out.append(f"\n... [truncated at line {next_line} of {len(lines)} — "
                       f"re-read with offset={next_line} to continue] ...\n")
            break
        out.append(entry)
        used += len(entry)
    return banner + "".join(out)

# Formats whose files are containers, not text: .docx/.pptx/.xlsx are zip
# archives of XML parts, .pdf has an object table and xref. Writing a string
# to one of these paths produces a file that only *looks* right by extension.
# The read side of this was already fixed (see _DOCUMENT_EXTRACTORS); this is
# the write half. It reported "Successfully wrote 4933 chars" for a .docx
# that nothing could open, and the failure only surfaced several tool calls
# later as "Package not found", by which point the model was debugging its
# converter instead of its choice of tool.
_STRUCTURED_WRITE_FORMATS = {
    ".docx": "a Word document",
    ".doc": "a Word document",
    ".pptx": "a PowerPoint presentation",
    ".ppt": "a PowerPoint presentation",
    ".xlsx": "an Excel workbook",
    ".xls": "an Excel workbook",
    ".pdf": "a PDF file",
    ".odt": "an OpenDocument text file",
    ".ods": "an OpenDocument spreadsheet",
    ".odp": "an OpenDocument presentation",
    ".rtf": "an RTF document",
}


def _structured_write_error(path: Path, kind: str) -> str:
    """Say why this cannot work and what to reach for instead."""
    suffix = path.suffix.lower()
    if suffix in (".docx", ".doc"):
        how = ("the word-docs MCP tools (create_document, add_heading, "
               "add_paragraph, add_table), which build a real document")
    elif suffix == ".pdf":
        how = ("build the .docx first with the word-docs MCP tools, then "
               "convert_to_pdf — that keeps headings and bold. Generating a "
               "PDF from plain text loses all formatting")
    elif suffix in (".xlsx", ".xls"):
        how = "openpyxl through run_command"
    elif suffix in (".pptx", ".ppt"):
        how = "python-pptx through run_command"
    else:
        how = f"a library that writes {kind}, through run_command"
    return (f"Error: {path.name} is {kind} — a structured container, not text. "
            f"write_file would produce a file that has the right extension and "
            f"cannot be opened. Use {how}.")


def handle_write_file(params: dict) -> str:
    path = _resolve(params["file_path"])
    if not _safe(path, write=True):
        return _outside_project_error(path, write=True)
    kind = _STRUCTURED_WRITE_FORMATS.get(path.suffix.lower())
    if kind:
        return _structured_write_error(path, kind)
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(params["content"], encoding="utf-8")
    return f"Successfully wrote {len(params['content'])} chars to {path}"

def handle_edit_file(params: dict) -> str:
    path = _resolve(params["file_path"])
    if not _safe(path, write=True):
        return _outside_project_error(path, write=True)
    if not path.exists():
        return f"Error: file not found: {path}"
    content = path.read_text(encoding="utf-8")
    old = params["old_string"]
    replace_all = bool(params.get("replace_all", False))
    count = content.count(old)
    if count == 0:
        return f"Error: old_string not found in {path}"
    if count > 1 and not replace_all:
        # Ambiguity still fails loudly; replace_all is how you say "all of them"
        # instead of synthesising N disambiguating contexts by hand.
        return (f"Error: old_string matches {count} locations; be more specific, "
                f"or pass replace_all=true to replace all {count}.")
    new_content = content.replace(old, params["new_string"], -1 if replace_all else 1)

    # ── Generate a unified diff preview ──
    import difflib
    old_lines = content.splitlines(keepends=True)
    new_lines = new_content.splitlines(keepends=True)
    diff = difflib.unified_diff(
        old_lines, new_lines,
        fromfile=str(path),
        tofile=str(path),
        n=3,  # context lines
    )
    diff_text = "".join(diff)
    # Color the diff for terminal display
    colored_lines = []
    for line in diff_text.splitlines():
        if line.startswith('+') and not line.startswith('+++'):
            colored_lines.append(f'{GREEN}{line}{RESET}')
        elif line.startswith('-') and not line.startswith('---'):
            colored_lines.append(f'{RED}{line}{RESET}')
        elif line.startswith('@@'):
            colored_lines.append(f'{CYAN}{line}{RESET}')
        else:
            colored_lines.append(f'{DIM}{line}{RESET}')
    colored_diff = "\n".join(colored_lines)

    # Write the new content
    path.write_text(new_content, encoding="utf-8")

    # Return a summary with the diff
    n_add = sum(1 for l in diff_text.splitlines() if l.startswith('+') and not l.startswith('+++'))
    n_del = sum(1 for l in diff_text.splitlines() if l.startswith('-') and not l.startswith('---'))
    note = f", {count} replacements" if replace_all and count > 1 else ""
    return f"Successfully edited {path} (+{n_add} -{n_del} lines{note})\n\n{colored_diff}"

def handle_list_files(params: dict) -> str:
    path = _resolve(params.get("path", "."))
    if not _safe(path):
        return _outside_project_error(path)
    if not path.exists():
        return f"Error: directory not found: {path}"
    if not path.is_dir():
        # Listing a file used to reach iterdir() and surface the raw OS error
        # ("[WinError 267] The directory name is invalid"), which tells the
        # model nothing it can act on. Naming the right tool ends the guessing
        # in one call instead of two.
        return (f"Error: {path} is a file, not a directory. "
                f"Use read_file to read it.")
    entries = []
    for child in sorted(path.iterdir()):
        # skip noise
        if child.name in {".git", "__pycache__", ".agent"}:
            continue
        entries.append(f"{'[dir] ' if child.is_dir() else '      '}{child.name}")
    return "\n".join(entries) if entries else "(empty)"

# An interpreter token: `python`, `python3`, `python.exe`, or any of those with
# a directory in front of it — `.venv\Scripts\python.exe`, `/usr/bin/python3` —
# optionally quoted.
#
# **The prefix is part of the token, and leaving it out corrupted the command.**
# `_PYTHON_INLINE_RE` used to start matching at the word `python`, so the
# rewrite below spliced the replacement in *after* any leading path and left it
# dangling:
#
#     .venv\Scripts\python.exe -c "print('hi')"
#         →  .venv\Scripts\"C:\...\.venv\Scripts\python.exe" -u "…\_exec.py"
#         →  cmd.exe: The filename, directory name, or volume label syntax is
#            incorrect.
#
# The agent's own system prompt tells the model to run
# `C:\...\.venv\Scripts\python.exe`, so this broke the exact command shape the
# prompt asks for. One measured session spent 64 of its 117 `run_command`
# calls failing, most of them here.
#
# The directory part is optional but, when present, must be consumed: `[^\s"|&<>]*`
# stops at a shell separator so `cd X && python -c …` still matches only the
# interpreter.
_PY_EXE = (r'(?:"(?:[^"\n]*[\\/])?python3?(?:\.exe)?"'
           r'|(?:[^\s"|&<>]*[\\/])?python3?(?:\.exe)?)')

# Matches a `python -c "..."` payload at the end of a command line. cmd.exe
# cannot carry newlines or nested quotes through such a payload, so it is
# round-tripped via a temporary script file instead.
_PYTHON_INLINE_RE = re.compile(
    rf'(?P<exe>{_PY_EXE})\s+(?:-u\s+)?-c\s+"(?P<code>.*)"\s*$', re.S | re.I
)
# `python -c` without -u: cmd.exe swallows stdout of short-lived processes.
_PYTHON_DASH_C_RE = re.compile(rf'(?P<exe>{_PY_EXE})\s+-c\b', re.I)

# A bare `python3` on Windows is never the interpreter this project means.
# `.venv\Scripts` ships `python.exe` and `pythonw.exe` and no `python3.exe`, so
# `python3` falls through PATH to whatever shim answers first — here an msys64
# build with none of the project's packages, which is why a measured session
# saw `python3 analyze.py` return `ModuleNotFoundError: No module named 'fitz'`
# from a venv that has PyMuPDF installed. Models reach for `python3` out of
# POSIX habit and have no way to see that it resolves elsewhere.
_BARE_PYTHON3_RE = re.compile(r'(?<![\w\\/.-])python3(?!\.exe)\b', re.I)


def _normalise_windows_command(cmd: str) -> tuple[str, Optional[str]]:
    """Work around three cmd.exe/Windows defects around python invocations.

    Returns (command, temp_dir_to_clean). The temp directory is created
    outside the project so scratch files never land in the source tree and
    never collide with `unittest discover`.

    Each rewrite preserves what the caller named wherever it can. An
    interpreter given with a path is honoured, not replaced: a caller who
    types `C:\\Python313\\python.exe` has chosen an interpreter, and silently
    substituting the venv's would be a different bug from the one being fixed.
    Only a *bare* name is resolved, because a bare name on Windows resolves to
    whatever PATH answers with and the agent cannot see what that is.
    """
    if sys.platform != "win32":
        return cmd, None

    # 1. `python3` → the interpreter this process is running. See
    #    _BARE_PYTHON3_RE: there is no python3.exe in a Windows venv.
    cmd = _BARE_PYTHON3_RE.sub(lambda _: f'"{sys.executable}"', cmd)

    # 2. Force unbuffered output, keeping the interpreter the caller named.
    cmd = _PYTHON_DASH_C_RE.sub(lambda m: f"{m.group('exe')} -u -c", cmd)

    # 3. Multi-line or nested-quote payloads cannot survive cmd.exe tokenising.
    m = _PYTHON_INLINE_RE.search(cmd)
    if m and ("\n" in m.group("code") or "'" in m.group("code")):
        temp_dir = tempfile.mkdtemp(prefix="tomas_exec_")
        script = Path(temp_dir) / "_exec.py"
        # The code reached us still escaped for cmd.exe's outer double quotes
        # (a caller who needs a literal `"` inside `-c "..."` has to write
        # `\"`, e.g. `f'{d[\"x\"]}'`), but the temp file has no outer quoting
        # to survive — writing `\"` verbatim leaves a bare backslash inside
        # an f-string expression, a SyntaxError. Measured live: two separate
        # sessions each hit exactly this and burned a run_command call before
        # working around it by hand. Unescape before writing.
        code = m.group("code").replace('\\"', '"')
        script.write_text(code, encoding="utf-8")
        named = m.group("exe").strip('"')
        # A pathed interpreter is a choice; a bare one is a guess we can improve.
        exe = named if ("\\" in named or "/" in named) else sys.executable
        cmd = (cmd[:m.start()] + f'"{exe}" -u "{script}"' + cmd[m.end():])
        return cmd, temp_dir
    return cmd, None


# A command that hands off to a process it never waits for: Windows' `start`
# builtin, or a trailing POSIX-style `&`. subprocess.run(capture_output=True)
# waits for the stdout/stderr pipes to reach EOF, not for the shell to exit —
# and a detached grandchild (e.g. `start /b cmd /c "npm run dev"`) inherits
# the write end of those pipes and keeps it open for as long as it runs. A
# dev server launched this way blocked the tool call for the full timeout
# with nothing on screen to say why, because the shell itself had already
# returned instantly.
_BACKGROUND_START_RE = re.compile(r'(?:^|[&;]|\|\|)\s*start\s+', re.IGNORECASE)


def _looks_backgrounded(cmd: str) -> bool:
    if sys.platform == "win32" and _BACKGROUND_START_RE.search(cmd):
        return True
    stripped = cmd.rstrip()
    return stripped.endswith('&') and not stripped.endswith('&&')


def _run_command_background(cmd: str, env: dict) -> str:
    """Launch a detached process and return immediately, without capturing
    its output — there is no pipe left open for anything to block on."""
    kwargs = {}
    if sys.platform == "win32":
        kwargs["creationflags"] = (
            subprocess.CREATE_NEW_PROCESS_GROUP | subprocess.DETACHED_PROCESS
        )
    try:
        subprocess.Popen(
            cmd, shell=True, cwd=str(PROJECT_DIR), env=env,
            stdin=subprocess.DEVNULL, stdout=subprocess.DEVNULL,
            stderr=subprocess.DEVNULL, **kwargs,
        )
    except Exception as e:
        return f"Error: failed to launch background command: {e}"
    return ("[exit 0 — ok] Started in the background (detached) — its output "
            "is not captured. Check the port it binds to, or its own log "
            "file, to confirm it came up.")


def _kill_process_tree(proc: "subprocess.Popen") -> None:
    if sys.platform == "win32":
        # proc.kill() only signals the direct child (cmd.exe); /T also takes
        # down whatever it spawned, so a killed build doesn't leave orphans.
        subprocess.run(["taskkill", "/F", "/T", "/PID", str(proc.pid)],
                       capture_output=True)
    else:
        proc.kill()


def _console_codepages() -> list[str]:
    """Encodings a Windows console command might actually have written in.

    `PYTHONIOENCODING=utf-8` makes *Python* children emit UTF-8, and the pipe
    was decoded as UTF-8 on that basis. But `dir`, `type` and the rest of
    cmd.exe's builtins emit the OEM codepage regardless — 866 on a Ukrainian
    or Russian Windows — so every Cyrillic filename came back as replacement
    characters. Measured in two live sessions, both of which were looking at
    a directory whose one interesting file was named
    `Коваль_Олександр_Дмитрович_CURSOVA.docx` and saw
    `������_����ᠭ��_����஢��_CURSOVA.docx`.
    """
    pages = ["utf-8"]
    if sys.platform == "win32":
        try:
            import ctypes
            for fn in ("GetOEMCP", "GetACP"):
                cp = getattr(ctypes.windll.kernel32, fn)()
                if cp:
                    pages.append(f"cp{cp}")
        except Exception:
            pages += ["cp866", "cp1251"]
    return pages


def _decode_console(data: bytes) -> str:
    """Decode command output, line by line, best encoding first.

    Per line rather than per stream because one command can produce both:
    `dir & python -c "print('привіт')"` writes OEM bytes and then UTF-8 bytes
    into the same pipe, and any single whole-stream choice mangles one half.
    A line that decodes as UTF-8 was UTF-8 — the encoding is
    self-validating, and ASCII (most output) decodes identically either way.
    """
    if not data:
        return ""
    pages = _console_codepages()
    out = []
    for line in data.split(b"\n"):
        for page in pages:
            try:
                out.append(line.decode(page))
                break
            except (UnicodeDecodeError, LookupError):
                continue
        else:
            out.append(line.decode("utf-8", errors="replace"))
    return "\n".join(out)


def _run_command_foreground(cmd: str, env: dict, timeout: int) -> str:
    creationflags = subprocess.CREATE_NEW_PROCESS_GROUP if sys.platform == "win32" else 0
    # Binary pipes: the encoding is decided per line afterwards, because it
    # is not one encoding. See `_decode_console`.
    proc = subprocess.Popen(
        cmd, shell=True, cwd=str(PROJECT_DIR), env=env,
        stdout=subprocess.PIPE, stderr=subprocess.PIPE,
        creationflags=creationflags,
    )
    deadline = time.monotonic() + timeout
    stdout = stderr = ""
    while True:
        try:
            # Polled in short slices rather than one blocking wait, so an Esc
            # interrupt takes effect within a fraction of a second instead of
            # only at the timeout.
            raw_out, raw_err = proc.communicate(timeout=0.25)
            stdout, stderr = _decode_console(raw_out), _decode_console(raw_err)
            break
        except subprocess.TimeoutExpired:
            if _CURRENT_INTERRUPT is not None and _CURRENT_INTERRUPT.is_set():
                _kill_process_tree(proc)
                try:
                    proc.wait(timeout=5)
                except Exception:
                    pass
                return "[exit -1 — interrupted] Command killed (Esc pressed)."
            if time.monotonic() >= deadline:
                _kill_process_tree(proc)
                # communicate() raised, so stdout/stderr are still "". Drain
                # the pipes after the kill, or a 120 s test run reports the
                # bare fact that time ran out and nothing it printed.
                try:
                    raw_out, raw_err = proc.communicate(timeout=5)
                    stdout = _decode_console(raw_out)
                    stderr = _decode_console(raw_err)
                except Exception:
                    stdout = stderr = ""
                timed_out = []
                if (stdout or "").strip():
                    timed_out.append(stdout.rstrip())
                if (stderr or "").strip():
                    timed_out.append(f"[stderr]\n{stderr.rstrip()}")
                notice = f"Error: command timed out after {timeout}s"
                if not timed_out:
                    return notice
                partial = "\n".join(timed_out)
                if len(partial) > 30000:
                    partial = partial[:15000] + "\n\n... [truncated] ...\n\n" + partial[-15000:]
                return (f"{notice} — output produced before the kill:\n{partial}"
                        + _diagnose_disk_full(partial))

    parts = []
    if (stdout or "").strip():
        parts.append(stdout.rstrip())
    if (stderr or "").strip():
        parts.append(f"[stderr]\n{stderr.rstrip()}")
    body = "\n".join(parts) or "(no output)"
    if len(body) > 30000:
        body = body[:15000] + "\n\n... [truncated] ...\n\n" + body[-15000:]
    # The exit code is always reported. A command that fails while still
    # writing to stdout used to be indistinguishable from one that succeeded.
    if proc.returncode == 0:
        status = "ok"
    elif _is_compound(cmd):
        # A shell chain exits with the status of its *last* command, so
        # `dir existing & dir missing` is reported as a failure even though
        # the first half produced exactly what was asked for. Measured: a
        # session read that banner, treated the whole result as void, and
        # spent its next turn re-running the half that had already worked.
        status = (f"last command in the chain failed (exit {proc.returncode}) "
                  f"— earlier commands may have succeeded; their output is below")
    else:
        status = f"FAILED (exit {proc.returncode})"
    return f"[exit {proc.returncode} — {status}]\n{body}" + _diagnose_disk_full(body)


#: Shell separators that chain several commands into one call. `&&` and `||`
#: are covered by the `&` and `|` members: this only needs to know that more
#: than one command ran, not how they were joined.
_COMPOUND_SEPARATORS = ("&", "|", ";", "\n")


def _is_compound(cmd: str) -> bool:
    """Whether this call ran more than one command.

    Quoted separators do not count — `python -c "print('a&b')"` is one
    command — so anything inside quotes is blanked before looking.
    """
    outside = re.sub(r'"[^"]*"|\'[^\']*\'', "", cmd)
    return any(sep in outside for sep in _COMPOUND_SEPARATORS)


# ---------------------------------------------------------------------------
# run_command advisories — checked before execution, never block it
# ---------------------------------------------------------------------------
#
# A saved session spent 48 minutes hand-polling an `npm install` that was
# never going to finish: its own log showed each package taking longer than
# the last (7s -> 11s -> 57s -> 142s -> 223s) before stalling completely, and
# nothing ever checked the one fact that explained it — the volume npm's
# cache lives on was at 0.06 GB free. A write that slow on a near-full disk
# looks exactly like a hang; only the free-space number tells them apart.
# These checks put that number in front of the model before and after the
# fact, instead of dozens of `dir`/`tasklist` calls discovering it by hand.

#: Command shapes that write enough to disk for "how much room is there" to
#: matter. Not exhaustive — a warning that fires sometimes is still worth
#: more than the silence that let the session above run for 48 minutes.
_DISK_HUNGRY_CMD_RE = re.compile(
    r'\b(npm|pnpm|yarn)\s+(install|i|ci|add)\b'
    r'|\bpip3?\s+install\b'
    r'|\bgit\s+clone\b'
    r'|\bdocker\s+(build|pull)\b'
    r'|\b(choco|winget)\s+install\b'
    r'|\bapt(-get)?\s+install\b',
    re.IGNORECASE,
)

#: Below this, an install of any real size is more likely to stall or fail
#: from disk exhaustion than to just be slow. Not a hard floor — actual
#: installs vary by an order of magnitude — just the point below which the
#: warning is worth the noise.
LOW_DISK_WARN_GB = 2.0


def _disk_free_gb(path: Optional[Path] = None) -> float:
    """Free space, in GB, on the volume holding `path` (project root by
    default). -1 if it cannot be read — this must never raise into a tool
    result over a diagnostic that is itself optional."""
    try:
        return shutil.disk_usage(str(path or PROJECT_DIR)).free / (1024 ** 3)
    except Exception:
        return -1.0


def _low_disk_warning(cmd: str) -> str:
    """A prefix warning for install-shaped commands on a near-full disk, or
    "" if neither condition holds."""
    if not _DISK_HUNGRY_CMD_RE.search(cmd):
        return ""
    free = _disk_free_gb()
    if free < 0 or free >= LOW_DISK_WARN_GB:
        return ""
    return (f"[warning] Only {free:.2f} GB free on this drive. Installs "
            f"this size have previously stalled or failed here from disk "
            f"exhaustion rather than any problem with the command itself — "
            f"consider freeing space before this runs.\n\n")


#: Substrings that mean "the volume ran out of room," across the OSes and
#: tools this agent touches (Windows' own [Errno 28] wording, Linux's ENOSPC
#: text, npm/pip's phrasing of the same). Matched against combined
#: stdout+stderr so one check covers every command shape.
_DISK_FULL_SIGNATURES = (
    "no space left on device",
    "[errno 28]",
    "there is not enough space on the disk",
    "disk full",
)


def _diagnose_disk_full(output: str) -> str:
    """A diagnosis to append to a command's result, or "" if its output shows
    no sign of disk exhaustion.

    Also records the pattern through the learning system (`kind="note"`, the
    same gated path self-notes uses for auto-generated observations — see
    `self_notes._bridge_to_learning`) so a second occurrence on this machine
    starts building toward a standing fact instead of being rediscovered
    from scratch every session.
    """
    low = (output or "").lower()
    if not any(sig in low for sig in _DISK_FULL_SIGNATURES):
        return ""
    free = _disk_free_gb()
    where = f"{free:.2f} GB free" if free >= 0 else "free space unknown"
    try:
        learning.remember(
            "note",
            "This machine's project drive has run out of disk space during "
            "a shell command. Check free space before large installs "
            "(npm/pip/git clone/docker) rather than assuming a slow "
            "command is merely slow.",
            evidence=f"run_command hit a disk-full signature ({where})")
    except Exception:
        pass
    return (f"\n\n[diagnosis] The volume is out of space ({where}). That is "
            f"very likely the real cause, not the command itself — a "
            f"near-full disk makes writes pathologically slow before they "
            f"fail outright, which can look like a hang for a long time "
            f"first. Free up space before retrying.")


#: Not blocked — sometimes genuinely needed — but killing by process *name*
#: takes out every match on the machine, not just the one this session
#: started. `_kill_process_tree` already does this right (by PID) for
#: processes the tool itself launched; this only reaches commands the model
#: writes by hand.
_UNSCOPED_KILL_RE = re.compile(
    r'\btaskkill\b[^&|;\n]*/im\b|\bpkill\b|\bkillall\b', re.IGNORECASE)


def _unscoped_kill_warning(cmd: str) -> str:
    if not _UNSCOPED_KILL_RE.search(cmd):
        return ""
    return ("[warning] This kills every process matching that name on the "
            "whole machine, not just ones this session started. Prefer a "
            "PID-scoped kill (`taskkill /F /PID <pid>`) when the PID is "
            "known, e.g. from `tasklist` or a port lookup.\n\n")


def handle_run_command(params: dict) -> str:
    cmd = params["command"]
    for bad in BLOCKED_PATTERNS:
        if bad in cmd:
            return f"Error: blocked dangerous pattern: {bad}"
    timeout = int(params.get("timeout", 120))

    # Advisory only — neither check blocks execution, and both are judged
    # against what the model actually typed, before Windows normalisation
    # rewrites it.
    advisory = _low_disk_warning(cmd) + _unscoped_kill_warning(cmd)

    cmd, temp_dir = _normalise_windows_command(cmd)
    # Child processes must emit UTF-8 rather than the console codepage,
    # otherwise non-ASCII output is mangled beyond recovery on the way back.
    env = dict(os.environ, PYTHONIOENCODING="utf-8")
    try:
        if _looks_backgrounded(cmd):
            result = _run_command_background(cmd, env)
        else:
            result = _run_command_foreground(cmd, env, timeout)
    finally:
        if temp_dir:
            shutil.rmtree(temp_dir, ignore_errors=True)
    return advisory + result


#: Last-seen (size_bytes, checked_at) for a path `check_progress` has been
#: asked to watch. Keyed by the resolved absolute path so two different
#: spellings of the same file share one baseline.
_PROGRESS_BASELINES: dict[str, tuple[int, float]] = {}

#: Below this gap between checks, "no growth" is not evidence of a stall —
#: it is evidence the model checked too soon. A background install writes in
#: bursts (one fetch, one extract), not a steady drip.
PROGRESS_MIN_STALL_SECONDS = 5.0

#: Past this many files, a directory scan for `check_progress` stops and
#: reports what it had — a stall detector must not itself become the slow
#: part, which a full walk of a half-installed `node_modules` risks being.
_PROGRESS_DIR_FILE_CAP = 50_000


def _dir_size(path: Path) -> tuple[int, bool]:
    """Recursive byte size of `path`, and whether the scan hit the cap."""
    total = 0
    count = 0
    for root, _dirs, files in os.walk(path):
        for name in files:
            count += 1
            if count > _PROGRESS_DIR_FILE_CAP:
                return total, True
            try:
                total += (Path(root) / name).stat().st_size
            except OSError:
                continue
    return total, False


def _tail_file(path: Path, n_lines: int = 15, max_bytes: int = 8000) -> str:
    """The last few lines of a file, read from the end so a multi-GB log
    costs a few KB to check rather than a full read."""
    try:
        size = path.stat().st_size
        with path.open("rb") as f:
            f.seek(max(0, size - max_bytes))
            data = f.read()
        return "\n".join(data.decode("utf-8", errors="replace").splitlines()[-n_lines:])
    except OSError:
        return ""


def handle_check_progress(params: dict) -> str:
    """Report whether a background job is still making progress, without the
    model hand-rolling a polling loop out of run_command calls.

    A saved session tried to answer "is npm still working" across two dozen
    `dir`/`type`/`tasklist` calls, each starting from nothing, because no
    single one of them remembered what the last one saw. This keeps that
    baseline itself — one call answers "did anything change since I last
    asked" — and folds in the same disk-space check `_diagnose_disk_full`
    uses, since a stall and a full disk are usually the same event.
    """
    raw = (params or {}).get("path", "").strip()
    if not raw:
        return ("Error: check_progress needs a path — the log file a "
                 "backgrounded command is writing to, or a growing output "
                 "directory such as node_modules.")
    path = _resolve(raw)
    if not _safe(path):
        return _outside_project_error(path)
    if not path.exists():
        return f"'{raw}' does not exist yet — nothing to check."

    now = time.monotonic()
    tail = ""
    if path.is_dir():
        size, capped = _dir_size(path)
    else:
        try:
            size = path.stat().st_size
        except OSError as e:
            return f"Error: could not stat '{raw}': {e}"
        capped = False
        tail = _tail_file(path)

    key = str(path)
    prev = _PROGRESS_BASELINES.get(key)
    _PROGRESS_BASELINES[key] = (size, now)

    lines = [f"{'Directory' if path.is_dir() else 'File'}: {raw}",
             f"Size: {size:,} bytes"
             + (" (capped scan — larger than shown)" if capped else "")]

    if prev is None:
        lines.append("No earlier check on this path — this call is the "
                      "baseline; call again in a bit to see whether it moved.")
    else:
        prev_size, prev_time = prev
        elapsed = now - prev_time
        delta = size - prev_size
        lines.append(f"Since last check ({elapsed:.1f}s ago): {delta:+,} bytes")
        if delta == 0 and elapsed >= PROGRESS_MIN_STALL_SECONDS:
            free = _disk_free_gb(path)
            note = f"{free:.2f} GB free on that volume" if free >= 0 else "free space unknown"
            lines.append(f"[stalled] No growth in {elapsed:.1f}s — likely "
                        f"stuck rather than merely slow ({note}; a near-full "
                        f"disk is the most common cause seen here). Check "
                        f"disk space and the process list rather than "
                        f"waiting longer.")
        elif delta > 0:
            lines.append("Still growing — likely still working.")

    if tail:
        lines.append("Last lines:")
        lines.append(tail)
    return "\n".join(lines)

SEARCH_PAGE_SIZE = 50

#: Directories that never hold source worth grepping. Descending into them is
#: what made one search_code call read 914 MB across 26,383 files — 244 s to
#: return 16 matches. Pruning is measured relative to the *search root*, so
#: pointing search_code into .venv on purpose still works; only incidental
#: descent is skipped.
SEARCH_IGNORE_DIRS = frozenset({
    ".git", ".venv", "venv", "node_modules", "__pycache__", "site-packages",
    ".mypy_cache", ".pytest_cache", ".ruff_cache", ".next", ".tox",
    "dist", "build",
})
#: Past this a file is data, not code, and scanning it line-by-line is waste.
SEARCH_MAX_FILE_BYTES = 1_000_000


def handle_search_code(params: dict) -> str:
    pattern = params["pattern"]
    path = _resolve(params.get("path", "."))
    if not _safe(path):
        return _outside_project_error(path)
    file_glob = params.get("file_glob", "")
    offset = max(0, int(params.get("offset", 0)))

    # A file is a legitimate search target. rglob() on one yields nothing, so
    # the old code answered "no matches" for a path that was never searched —
    # a confident false negative the model had no reason to doubt.
    if path.is_file():
        candidates = [path]
    elif path.is_dir():
        candidates = path.rglob(file_glob) if file_glob else path.rglob("*")
    else:
        return f"Error: path does not exist: {path}"

    try:
        regex = re.compile(pattern)
    except re.error as e:
        return f"Error: invalid regex {pattern!r}: {e}"

    matches: list[str] = []
    total = 0
    skipped = 0
    for file in candidates:
        if not file.is_file():
            continue
        # Filtered here rather than in the glob so `file_glob` semantics are
        # untouched. The walk itself costs ~1 s; the 244 s was read_text().
        try:
            rel_dirs = set(file.relative_to(path).parts[:-1])
        except ValueError:
            rel_dirs = set()
        if rel_dirs & SEARCH_IGNORE_DIRS:
            skipped += 1
            continue
        try:
            if file.stat().st_size > SEARCH_MAX_FILE_BYTES:
                skipped += 1
                continue
            data = file.read_bytes()
            if b"\0" in data[:8192]:   # binary: a "match" here is unusable
                skipped += 1
                continue
            text = data.decode("utf-8", errors="replace")
        except Exception:
            continue
        for i, line in enumerate(text.splitlines(), 1):
            if regex.search(line):
                total += 1
                if offset < total <= offset + SEARCH_PAGE_SIZE:
                    matches.append(f"{file}:{i}: {line}")

    if total == 0:
        # A silent skip is the same confident false negative the file-path
        # branch above exists to prevent: say what was not searched.
        note = (f" ({skipped} files skipped: vendored, binary or over "
                f"{SEARCH_MAX_FILE_BYTES // 1000} KB)") if skipped else ""
        return f"No matches for pattern: {pattern}{note}"
    header = f"{total} match{'' if total == 1 else 'es'}"
    shown = len(matches)
    if total > offset + shown:
        # Report the true total so the model can decide whether to page,
        # instead of guessing what "truncated" hid.
        header += (f" — showing {offset + 1}-{offset + shown}; "
                   f"re-run with offset={offset + shown} for more")
    elif offset:
        header += f" — showing {offset + 1}-{offset + shown}"
    return header + "\n" + "\n".join(matches)

#: What `save_memory` needs, in the order the schema declares it.
_SAVE_MEMORY_FIELDS = ("key", "description", "content")


def handle_save_memory(params: dict) -> str:
    """Store a memory, or say what is missing.

    Indexing `params` directly raised `KeyError: 'key'` straight out of the
    handler, and the tool-call machinery handed the model the string `'key'` as
    the result — no tool name, no field list, no hint that anything was missing.
    A measured session called `save_memory({"content": ...})`, got back `'key'`,
    and never retried.

    Every other handler here returns `Error: …` for bad input. This one is the
    outlier, and the cost of the inconsistency is a lost memory write on the
    one turn of the corpus whose whole purpose is storing a rule.
    """
    params = params or {}
    missing = [f for f in _SAVE_MEMORY_FIELDS if not str(params.get(f, "")).strip()]
    if missing:
        return (f"Error: save_memory needs {', '.join(_SAVE_MEMORY_FIELDS)}; "
                f"missing or empty: {', '.join(missing)}. "
                f"Got: {sorted(params) or 'no arguments'}")
    store, why = save_memory(params["key"], params["description"],
                             params["content"],
                             store=str(params.get("store", "")).strip())
    # Naming the store and the reason is not decoration: the model is the
    # thing that can override a wrong guess next time, and it cannot learn
    # the routing from a result that only says "Saved".
    where = {
        STORE_INSTRUCTION: "instructions (cached; applies every turn, costs nothing per turn)",
        STORE_RULE: "standing rules (injected every turn — one of 10 slots)",
        STORE_FACT: "learned facts (retrieved when the topic comes up)",
        STORE_NOTE: "self-notes (reachable with /notes; not in the prompt)",
    }[store]
    return (f"Saved memory '{params['key']}' -> {where}.\n"
            f"Why: {why}.\n"
            f"If that is the wrong store, call save_memory again with "
            f"store=\"instruction|rule|fact|note\".")


def handle_read_mcp_resource(params: dict) -> str:
    """Read (or list) MCP resources.

    The MCP spec defines resources alongside tools, and much of the
    ecosystem's value sits there; the manager surfaced only tools until now.
    """
    if mcp_manager is None:
        return "Error: no MCP servers are connected."
    uri = (params or {}).get("uri", "").strip()
    server = (params or {}).get("server") or None
    if not uri:
        resources = mcp_manager.list_resources()
        if not resources:
            return "No MCP resources are available from the connected servers."
        lines = [f"{len(resources)} MCP resource(s) available:"]
        for r in resources[:100]:
            desc = (r.get("description") or r.get("name") or "")[:90]
            lines.append(f"  [{r['server']}] {r.get('uri', '?')}  {desc}")
        if len(resources) > 100:
            lines.append(f"  ... and {len(resources) - 100} more")
        return "\n".join(lines)
    return mcp_manager.read_resource(uri, server=server)


def handle_ask_user_question(params: dict) -> str:
    """Ask the user one or more multiple-choice questions and block for the
    answer(s), through the same arrow-key picker `/config` uses (`_arrow_menu`
    / `_arrow_checklist`, defined further down this file — resolved at call
    time, so the forward reference is fine).

    Risk tier "none" (see RISK_LEVELS) means this never goes through the
    permission prompt: the interactive question itself is the human-in-the-
    loop control, so gating it behind a separate "approve this tool?" would
    be a redundant prompt in front of the real one.
    """
    questions = (params or {}).get("questions") or []
    if not isinstance(questions, list) or not questions:
        return "Error: 'questions' must be a non-empty list of question objects."

    # Same test `agent_loop` uses to decide TerminalAdapter(interactive=...).
    # This tool's risk tier ("none") means it skips the permission gate that
    # would otherwise catch a non-interactive run and deny instead of
    # blocking — without this check, a headless run (CI, the simulation
    # harness) would hang forever on msvcrt.getwch() waiting for a keypress
    # that can never come.
    if not bool(getattr(sys.stdin, "isatty", lambda: False)()):
        return (
            "Error: cannot ask interactively — no interactive terminal is "
            "attached (a non-interactive/headless run). Proceed using your "
            "best judgement, state the assumption you made, and let the "
            "user correct it afterward if it was wrong."
        )

    OTHER_LABEL = "Other"

    def _ask_custom(prompt_text: str) -> Optional[str]:
        """Read a free-text answer, visibly. `None` means "go back to the list".

        Not `input()`. Three things went wrong with it, and only the first was
        the spinner's fault: the menu was left on screen so you typed
        underneath a dead list; there was no way back to the options once
        "Other" was chosen; and Esc did nothing. This reads the same way the
        main prompt does — one line, echoed as you type, Ctrl+U/Ctrl+W to
        correct, Esc to change your mind.
        """
        try:
            import msvcrt
        except ImportError:                      # not Windows: no raw reader
            try:
                got = input(f'  {DIM}{prompt_text}{RESET} ').strip()
            except (EOFError, KeyboardInterrupt):
                print()
                return None
            return got

        buf: list[str] = []
        with CONSOLE:
            sys.stdout.write(f'  {DIM}{prompt_text}{RESET}\n  {GREEN}›{RESET} ')
            sys.stdout.flush()
            while True:
                ch = msvcrt.getwch()
                if ch == '\r':
                    sys.stdout.write('\n')
                    sys.stdout.flush()
                    return ''.join(buf).strip()
                if ch == '\x1b':                 # Esc — back to the options
                    sys.stdout.write('\r\033[2K')
                    sys.stdout.flush()
                    return None
                if ch in ('\x08', '\x7f'):       # Backspace
                    if buf:
                        buf.pop()
                elif ch == '\x15':               # Ctrl+U — clear the line
                    buf.clear()
                elif ch == '\x17':               # Ctrl+W — delete a word
                    while buf and buf[-1].isspace():
                        buf.pop()
                    while buf and not buf[-1].isspace():
                        buf.pop()
                elif ch in ('\xe0', '\x00'):     # arrows: consume the pair
                    msvcrt.getwch()
                    continue
                elif ch == '\x03':               # Ctrl+C — treat as cancel
                    sys.stdout.write('\n')
                    return None
                elif ch.isprintable():
                    buf.append(ch)
                else:
                    continue
                # One redraw per keystroke, one line: `\r` + erase + reprint is
                # exact here because the answer is a single line by
                # construction (Enter ends it).
                sys.stdout.write(f'\r\033[2K  {GREEN}›{RESET} ' + ''.join(buf))
                sys.stdout.flush()

    # Each picker and the free-text editor take `CONSOLE` themselves, and the
    # adapter is told not to spin over this tool at all (ToolStarted.interactive),
    # so nothing else reaches for the keyboard while any of this is on screen.
    answered = []
    for q in questions:
        text = str(q.get("question") or "").strip()
        options = q.get("options") or []
        if not text or not isinstance(options, list) or len(options) < 2:
            continue
        header = str(q.get("header") or "").strip()
        multi = bool(q.get("multiSelect"))

        # Who is asking goes on its own line, with the header beside it rather
        # than after the question: a header trailing a long question wraps onto
        # the next line and reads as part of the question.
        print()
        who = 'TOMAS asks' + (f' {DIM}· {header}{RESET}' if header else '')
        print(f'  {MAGENTA}{BOLD}▌{RESET} {MAGENTA}{who}{RESET}')
        print(f'  {BOLD}{text}{RESET}')
        print()

        # Descriptions in their own column, not trailing off each label behind
        # an em dash. With three options of different name lengths the ragged
        # version gives the eye no line to follow, and the choice is the whole
        # point of the screen.
        from text_display import display_width

        raw = [str(opt.get("label", "")).strip() for opt in options]
        width = max((display_width(r) for r in raw + [OTHER_LABEL]), default=0)

        # `_fit_label` clips the description to what the terminal can show on
        # one row. A wrapped row is not merely ugly — it breaks the picker's
        # rewind arithmetic, which is the "I switch and it disappears" defect.
        labels = [_fit_label(label, str(opt.get("description", "")).strip(), width)
                  for opt, label in zip(options, raw)]
        labels.append(_fit_label(OTHER_LABEL, 'type your own answer', width))
        other_idx = len(options)

        ask_own = 'Type your answer · Esc goes back to the options'
        if multi:
            picked = _arrow_checklist(
                "", labels,
                footer='Space toggle · Enter confirm · Esc skip') or []
            chosen = []
            for i in picked:
                if i == other_idx:
                    custom = _ask_custom(ask_own)
                    if custom:
                        chosen.append(custom)
                else:
                    chosen.append(str(options[i].get("label", "")))
        else:
            # Loops so that Esc in the free-text field returns to the list
            # rather than silently skipping the question. Choosing "Other" by
            # accident used to be unrecoverable.
            chosen = []
            while True:
                idx = _arrow_menu(
                    "", labels,
                    footer='↑↓ move · Enter select · Esc skip',
                    erase_on_exit=True)
                if idx == other_idx:
                    custom = _ask_custom(ask_own)
                    if custom is None:
                        continue          # Esc — show the options again
                    if custom:
                        chosen.append(custom)
                elif idx != -1:
                    chosen.append(str(options[idx].get("label", "")))
                break

        if chosen:
            print(f'  {GREEN}→{RESET} ' + ', '.join(chosen))
        else:
            print(f'  {DIM}→ skipped{RESET}')
        answered.append({"question": text, "answers": chosen,
                         "skipped": not chosen})

    # Say what an empty answer *means*. `{"answers": []}` reads as "the user
    # answered nothing", which a model resolves by asking the same question
    # again or by silently picking for them. Declining to choose is itself an
    # instruction: proceed, decide it yourself, and say what you assumed.
    if any(a["skipped"] for a in answered):
        return json.dumps({
            "answers": answered,
            "note": ("The user skipped one or more questions. Do not ask them "
                     "again. Choose the most reasonable option yourself, "
                     "state which one you chose and why, and continue."),
        }, ensure_ascii=False)
    return json.dumps({"answers": answered}, ensure_ascii=False)


def _ssl_context_with_certifi():
    """Default SSL context, pointed at certifi's CA bundle when available.

    Bare `urlopen` on this project's Windows/MSYS2 Python builds fails every
    HTTPS request with CERTIFICATE_VERIFY_FAILED — that interpreter has no
    usable system CA store wired into `ssl`. Measured in two live sessions:
    every `fetch_url` call hit the error and fell back to the much slower
    `fetch_url_with_browser` (or was abandoned). certifi is already present
    as a transitive dependency; reuse its bundle instead of adding one.
    """
    import ssl
    try:
        import certifi
        return ssl.create_default_context(cafile=certifi.where())
    except ImportError:
        return ssl.create_default_context()


#: Default read ceiling. Was 50,000,000 — a number that never once stopped
#: anything, while the thing actually worth limiting went unbounded: measured
#: across three live sessions, tool results occupied 110-134k characters of a
#: ~250k-character transcript, and the largest single one was 15,681
#: characters of HTML whose text content was 2,152 (14%). The transcript is
#: re-sent on every subsequent step, so one fetch is paid for many times.
_FETCH_MAX_BYTES = 200_000

_UNREADABLE_ELEMENTS = ("script", "style", "noscript", "template")
_SCRIPT_STYLE_RES = {
    tag: re.compile(rf"<{tag}\b[^>]*>.*?</{tag}\s*>", re.IGNORECASE | re.DOTALL)
    for tag in _UNREADABLE_ELEMENTS
}
_TAG_RE = re.compile(r"<[^>]+>")
_HTML_COMMENT_RE = re.compile(r"<!--.*?-->", re.DOTALL)
_BLANK_LINES_RE = re.compile(r"\n\s*\n\s*\n+")


def _html_to_text(html: str) -> str:
    """Readable text from a page, with the parts nobody can read removed.

    `fetch_url` returned the raw bytes. On the VNTU department pages that
    every measured session went to, that is 10-14 `<script>` blocks, inline
    CSS and Google Analytics wrapped around the five names the agent came
    for. The names were genuinely there — this is not a correctness fix — but
    the model had to find them inside markup, and every later step paid for
    the markup again.

    Deliberately not a parser: no dependency, and a malformed page still
    yields its text. `<script>`/`<style>` bodies go first (their *contents*
    are the noise, so dropping only the tags would keep the JavaScript), then
    tags, then entities, then runs of blank lines.
    """
    text = _HTML_COMMENT_RE.sub(" ", html)
    for tag, pattern in _SCRIPT_STYLE_RES.items():
        # Guarded on the closing tag: a document with many *unclosed*
        # `<script>` makes the non-greedy body scan to end-of-input once per
        # opener — 2.0s on a 169 KB adversarial input, against 8ms on a
        # realistic 234 KB page. If nothing closes, there is no pair to
        # remove and the generic tag pass below handles the orphans.
        if f"</{tag}" in text or f"</{tag.upper()}" in text:
            text = pattern.sub(" ", text)
    # Keep block boundaries as newlines so a list of names does not come back
    # as one run-on line.
    text = re.sub(r"(?i)<(br|/p|/div|/li|/tr|/h[1-6])\b[^>]*>", "\n", text)
    text = _TAG_RE.sub(" ", text)
    text = html_module.unescape(text)
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r" ?\n ?", "\n", text)
    return _BLANK_LINES_RE.sub("\n\n", text).strip()


def _looks_like_html(data: bytes, content_type: str) -> bool:
    if "html" in content_type.lower():
        return True
    head = data[:2048].lstrip().lower()
    return head.startswith(b"<!doctype html") or head.startswith(b"<html")


def handle_fetch_url(params: dict) -> str:
    """Fetch content from a URL."""
    import urllib.request
    import urllib.error

    url = params["url"]
    timeout = int(params.get("timeout", 15))
    max_size = int(params.get("max_size", _FETCH_MAX_BYTES))
    # Opt-out for the rare caller that genuinely wants markup (scraping a
    # specific attribute, checking a meta tag). Text is the default because
    # every measured use of this tool wanted the text.
    want_raw = bool(params.get("raw"))

    # Basic URL validation
    if not url.startswith(("http://", "https://")):
        return f"Error: URL must start with http:// or https://"

    # Block dangerous URLs
    blocked_patterns = ["localhost", "127.0.0.1", "0.0.0.0", "::1", "169.254.169.254"]
    for pattern in blocked_patterns:
        if pattern in url:
            return f"Error: blocked URL pattern: {pattern}"

    def _get(context, note: str) -> str:
        req = urllib.request.Request(
            url,
            headers={"User-Agent": "Agent-for-TOM/1.0 (fetch_url tool)"},
        )
        with urllib.request.urlopen(req, timeout=timeout, context=context) as response:
            content_type = response.headers.get("Content-Type", "")
            # Read one byte past the limit so truncation is detectable, then
            # keep what fits rather than refusing the whole response. A hard
            # error here sent one session to a headless browser for a page it
            # already had the useful half of.
            data = response.read(max_size + 1)
            clipped = len(data) > max_size
            data = data[:max_size]

        if "pdf" in content_type.lower() or url.lower().endswith(".pdf"):
            # A PDF is not text and must not be decoded as if it were. The
            # browser tool cannot help either — it aborts with "Download is
            # starting" — so say what to do instead of returning mojibake.
            return (f"{note}[{len(data)} bytes of application/pdf — not decoded]\n"
                    f"Read it with PyMuPDF instead:\n"
                    f"  python -c \"import pymupdf,urllib.request,ssl;"
                    f"ctx=ssl.create_default_context();ctx.check_hostname=False;"
                    f"ctx.verify_mode=ssl.CERT_NONE;"
                    f"d=urllib.request.urlopen({url!r},context=ctx).read();"
                    f"doc=pymupdf.open(stream=d,filetype='pdf');"
                    f"print(doc[0].get_text())\"")

        text = data.decode("utf-8", errors="replace")
        if not want_raw and _looks_like_html(data, content_type):
            stripped = _html_to_text(text)
            saved = len(text) - len(stripped)
            if saved > 0:
                note += (f"[html→text: {len(text)}→{len(stripped)} chars; "
                         f"pass raw=true for markup]\n")
            text = stripped
        if clipped:
            note += (f"[clipped at {max_size} bytes — raise max_size if the "
                     f"rest is needed]\n")
        return note + text

    try:
        return _get(_ssl_context_with_certifi(), "")
    except urllib.error.HTTPError as e:
        return f"Error: HTTP {e.code} {e.reason}"
    except urllib.error.URLError as e:
        # A host whose certificate chain is broken is not a transient failure,
        # and retrying it verified never works. Measured across three
        # sessions: 12 fetch_url calls, 8 failures, 7 of them
        # CERTIFICATE_VERIFY_FAILED against the same host — after which one
        # session wrote its own urllib script with CERT_NONE and got the file,
        # and the other two abandoned the source. Retry once unverified and
        # label the result, so the model can weigh it rather than be blocked
        # by it. Labelled, not silent: an unauthenticated page is a weaker
        # source and the reader has to know which one they got.
        if "CERTIFICATE_VERIFY_FAILED" not in str(e.reason):
            return f"Error: {e.reason}"
        try:
            import ssl
            unverified = ssl.create_default_context()
            unverified.check_hostname = False
            unverified.verify_mode = ssl.CERT_NONE
            return _get(unverified,
                        "[cert-unverified: this host's TLS certificate could "
                        "not be validated; content fetched anyway]\n")
        except Exception as retry_error:
            return (f"Error: {e.reason} (retry without certificate "
                    f"verification also failed: {retry_error})")
    except Exception as e:
        return f"Error: {e}"


def handle_fetch_url_with_browser(params: dict) -> str:
    """Fetch content from a URL using a headless browser (Playwright)."""
    if not PLAYWRIGHT_AVAILABLE:
        return "Error: Playwright not available. Install with: pip install playwright && playwright install chromium"

    import asyncio

    url = params["url"]
    wait_for = params.get("wait_for")
    timeout = int(params.get("timeout", 30)) * 1000  # Convert to milliseconds
    max_size = int(params.get("max_size", 50000000))
    take_screenshot = params.get("screenshot", False)

    # Basic URL validation
    if not url.startswith(("http://", "https://")):
        return f"Error: URL must start with http:// or https://"

    # Block dangerous URLs
    blocked_patterns = ["localhost", "127.0.0.1", "0.0.0.0", "::1", "169.254.169.254"]
    for pattern in blocked_patterns:
        if pattern in url:
            return f"Error: blocked URL pattern: {pattern}"

    async def _fetch():
        from playwright.async_api import async_playwright  # deferred, see top
        async with async_playwright() as p:
            browser = await p.chromium.launch(headless=True)
            try:
                context = await browser.new_context(
                    user_agent="Agent-for-TOM/1.0 (fetch_url_with_browser tool)"
                )
                page = await context.new_page()
                
                # Set timeout
                page.set_default_timeout(timeout)
                
                # Navigate to the page
                response = await page.goto(url, wait_until="domcontentloaded")
                
                if response and response.status >= 400:
                    return f"Error: HTTP {response.status} {response.status_text}"
                
                # Wait for specific selector if provided
                if wait_for:
                    try:
                        await page.wait_for_selector(wait_for, timeout=timeout)
                    except Exception:
                        return f"Error: timeout waiting for selector '{wait_for}'"
                
                # Get page content
                content = await page.content()
                
                # Optionally take screenshot
                screenshot_data = None
                if take_screenshot:
                    screenshot_bytes = await page.screenshot(full_page=True)
                    import base64
                    screenshot_data = base64.b64encode(screenshot_bytes).decode('utf-8')
                
                # Close browser
                await browser.close()
                
                # Check size
                if len(content) > max_size:
                    content = content[:max_size] + f"\n\n... [truncated, max {max_size} bytes]"
                
                result = content
                if screenshot_data:
                    result += f"\n\n[Screenshot captured: {len(screenshot_data)} bytes base64]"
                
                return result
            except Exception as e:
                await browser.close()
                raise e

    try:
        return asyncio.run(_fetch())
    except Exception as e:
        return f"Error: {e}"


def _usable_result(title: str, href: str) -> bool:
    """Whether a search hit is something the model can act on.

    The `ddgs` backend rotates between engines, and the Startpage one
    occasionally yields its own click-tracking entry instead of a result:
    empty title, empty snippet, and a relative
    `/clev?event=StartpageResultClick&...&payload={...}` href carrying a few
    hundred characters of session JSON. Seen three times across three
    measured sessions, always as result #1, which is the one a model reads
    most carefully. Nothing downstream can fetch a relative URL, so this is
    pure cost.
    """
    return bool(title.strip()) and href.strip().startswith(("http://", "https://"))


def handle_search_web(params: dict) -> str:
    """Search the internet using DDGS / Playwright by default, with fallbacks."""
    query = params["query"]
    max_results = int(params.get("max_results", 5))

    # ── Primary: DDGS ──
    try:
        try:
            from ddgs import DDGS
        except ImportError:
            from duckduckgo_search import DDGS

        results = []
        skipped = 0
        with DDGS() as ddgs:
            for r in ddgs.text(query, max_results=max_results):
                title = r.get("title") or ""
                body = r.get("body") or ""
                href = r.get("href") or ""
                if not _usable_result(title, href):
                    skipped += 1
                    continue
                results.append(f"{len(results)+1}. {title}\n   {body}\n   URL: {href}")

        if results:
            # Counted rather than hidden: a query that returns two usable hits
            # out of eight is a query worth rewording, and silently showing
            # two looks like the engine simply had little to say.
            note = f" ({skipped} unusable result(s) skipped)" if skipped else ""
            return f"Search results for '{query}'{note}:\n\n" + "\n\n".join(results)

    except Exception:
        pass

    # ── Secondary: Playwright Bing / Google Chrome ──
    try:
        from playwright.async_api import async_playwright
        import urllib.parse
        import asyncio

        encoded_query = urllib.parse.quote_plus(query)
        search_url = f"https://www.bing.com/search?q={encoded_query}"

        async def _playwright_search():
            async with async_playwright() as p:
                browser = await p.chromium.launch(headless=True)
                page = await browser.new_page(
                    user_agent="Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/124.0.0.0 Safari/537.36"
                )
                await page.goto(search_url, timeout=15000, wait_until="domcontentloaded")
                
                results = []
                elements = await page.query_selector_all("li.b_algo")
                for el in elements:
                    title_el = await el.query_selector("h2 a")
                    snippet_el = await el.query_selector("p, div.b_caption")
                    if title_el:
                        title = (await title_el.inner_text()).strip()
                        href = await title_el.get_attribute("href") or ""
                        snippet = (await snippet_el.inner_text()).strip() if snippet_el else ""
                        if title and href.startswith("http"):
                            results.append(f"{len(results)+1}. {title}\n   {snippet}\n   URL: {href}")
                            if len(results) >= max_results:
                                break
                await browser.close()
                if results:
                    return f"Search results for '{query}' (via Playwright / Chrome):\n\n" + "\n\n".join(results)
                return None

        pw_results = asyncio.run(_playwright_search())
        if pw_results:
            return pw_results
    except Exception:
        pass

    # ── Tertiary fallback: DuckDuckGo HTML ──
    try:
        import urllib.request
        import urllib.parse
        import re

        encoded = urllib.parse.quote_plus(query)
        req = urllib.request.Request(
            f"https://html.duckduckgo.com/html/?q={encoded}",
            headers={"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64)"}
        )
        with urllib.request.urlopen(req, timeout=10, context=_ssl_context_with_certifi()) as resp:
            html = resp.read().decode('utf-8', errors='ignore')
            results = []
            matches = re.findall(r'<a class="result__a" href="([^"]+)">(.*?)</a>', html)
            for href, title_html in matches:
                clean_title = re.sub(r'<[^>]+>', '', title_html).strip()
                # Same filter as the primary path — this fallback scrapes a
                # results page and can pick up the same relative redirects.
                if not _usable_result(clean_title, href):
                    continue
                results.append(f"{len(results)+1}. {clean_title}\n   URL: {href}")
                if len(results) >= max_results:
                    break
            if results:
                return f"Search results for '{query}':\n\n" + "\n\n".join(results)
    except Exception as e:
        return f"Error searching the web: {e}"

    return f"No results found for '{query}'"


# ── Browser control ────────────────────────────────────────────────────────
# Thin by design: `core/browser.py` owns the connection, the loop thread and
# every message the model reads. These six functions do argument coercion and
# nothing else, so there is one place where "what does a failed click say"
# is answered.

def _browser_interrupt() -> bool:
    """Whether Esc has been pressed, asked the way core/browser.py asks it.

    A function and not a captured reference: `_CURRENT_INTERRUPT` is rebound
    per turn, so a closure over the *event object* would poll the previous
    turn's event and never fire.
    """
    return _CURRENT_INTERRUPT is not None and _CURRENT_INTERRUPT.is_set()


def _browser():
    """The engine, imported when it is used. See ACTIONS_FOR_SCHEMA."""
    from core import browser as core_browser
    return core_browser


def handle_tab_list(params: dict) -> str:
    select = params.get("select")
    return _browser().tabs(
        select=int(select) if select is not None else None,
        new_tab=params.get("new_tab") or None,
        start_browser=bool(params.get("start_browser")),
        interrupt=_browser_interrupt,
    )


def handle_tab_snapshot(params: dict) -> str:
    return _browser().snapshot(interrupt=_browser_interrupt)


def handle_tab_read(params: dict) -> str:
    engine = _browser()
    return engine.read(
        ref=params.get("ref") or None,
        max_chars=int(params.get("max_chars", engine.MAX_READ_CHARS)),
        interrupt=_browser_interrupt,
    )


def handle_tab_navigate(params: dict) -> str:
    return _browser().navigate(
        url=params.get("url") or None,
        action=(params.get("action") or "").strip().lower(),
        interrupt=_browser_interrupt,
    )


def handle_tab_act(params: dict) -> str:
    return _browser().act(
        action=(params.get("action") or "").strip().lower(),
        ref=params.get("ref") or None,
        text=params.get("text"),
        key=params.get("key") or None,
        option=params.get("option"),
        submit=bool(params.get("submit")),
        # Defaults to replacing the field, because a model that means to set
        # a value and gets an append writes into whatever the user had
        # already typed there.
        clear_first=params.get("clear_first", True) is not False,
        interrupt=_browser_interrupt,
    )


def handle_tab_screenshot(params: dict) -> str:
    return _browser().screenshot(
        path=params.get("path") or None,
        full_page=bool(params.get("full_page")),
        interrupt=_browser_interrupt,
    )


# ── Live Office documents ─────────────────────────────────────────────────
# Thin, like the tab_* handlers: `core/office.py` owns the COM thread and
# every message the model reads.

def _office():
    """The engine, imported when it is used. See EDIT_ACTIONS_FOR_SCHEMA."""
    from core import office as core_office
    return core_office


def handle_doc_list(params: dict) -> str:
    select = params.get("select")
    return _office().documents(
        select=int(select) if select is not None else None,
        new_document=bool(params.get("new_document")),
        start_app=bool(params.get("start_app")),
        interrupt=_browser_interrupt,
    )


def handle_doc_outline(params: dict) -> str:
    return _office().outline(interrupt=_browser_interrupt)


def handle_doc_read(params: dict) -> str:
    engine = _office()
    return engine.read(
        ref=params.get("ref") or None,
        max_chars=int(params.get("max_chars", engine.MAX_READ_CHARS)),
        interrupt=_browser_interrupt,
    )


def handle_doc_find(params: dict) -> str:
    return _office().find(
        text=params.get("text") or "",
        max_hits=int(params.get("max_hits", 20)),
        interrupt=_browser_interrupt,
    )


def handle_doc_edit(params: dict) -> str:
    return _office().edit(
        action=(params.get("action") or "").strip().lower(),
        ref=params.get("ref") or None,
        text=params.get("text"),
        find=params.get("find") or None,
        style=params.get("style") or None,
        interrupt=_browser_interrupt,
    )


def handle_doc_save(params: dict) -> str:
    return _office().save(path=params.get("path") or None,
                          interrupt=_browser_interrupt)


HANDLERS: dict[str, Callable[[dict], str]] = {
    "read_file": handle_read_file,
    "write_file": handle_write_file,
    "edit_file": handle_edit_file,
    "list_files": handle_list_files,
    "run_command": handle_run_command,
    "check_progress": handle_check_progress,
    "search_code": handle_search_code,
    "save_memory": handle_save_memory,
    "fetch_url": handle_fetch_url,
    "fetch_url_with_browser": handle_fetch_url_with_browser,
    "search_web": handle_search_web,
    "tab_list": handle_tab_list,
    "tab_snapshot": handle_tab_snapshot,
    "tab_read": handle_tab_read,
    "tab_navigate": handle_tab_navigate,
    "tab_act": handle_tab_act,
    "tab_screenshot": handle_tab_screenshot,
    "doc_list": handle_doc_list,
    "doc_outline": handle_doc_outline,
    "doc_read": handle_doc_read,
    "doc_find": handle_doc_find,
    "doc_edit": handle_doc_edit,
    "doc_save": handle_doc_save,
    "read_mcp_resource": handle_read_mcp_resource,
    "ask_user_question": handle_ask_user_question,
}

def execute_tool(name: str, params: dict) -> str:
    handler = HANDLERS.get(name)
    if handler:
        try:
            return handler(params)
        except KeyError as e:
            # A bare `Error: 'file_path'` (str(KeyError) is just the missing
            # key) tells a model something is wrong but not what to send
            # instead — measured live, a weaker model retried the identical
            # empty tool call three times before the loop guard stopped it.
            hint = ""
            if not params:
                # Every argument gone, not one forgotten. Measured twice in one
                # session: `write_file {}` where the model meant to send a
                # 12 KB script. Arguments that large are the ones the output
                # limit truncates mid-JSON, and the parser then substitutes an
                # empty object — so "you forgot an argument" is the wrong
                # diagnosis and re-sending the same call is the wrong fix.
                hint = (" No arguments arrived at all, which usually means "
                        "they were cut off at the output limit rather than "
                        "omitted. If you were sending a large value, write it "
                        "in several smaller calls.")
            return f"Error: '{name}' is missing required argument {e}.{hint}"
        except Exception as e:
            return f"Error: {e}"
    # Try MCP tool dispatch (with name mapping for renamed conflicting tools)
    if mcp_manager:
        mcp_name = MCP_TOOL_NAME_MAP.get(name, name)  # resolve renamed name -> original
        return _call_mcp_tool_interruptibly(mcp_name, params)
    return f"Error: unknown tool '{name}'"


def _call_mcp_tool_interruptibly(name: str, params: dict) -> str:
    """Run one MCP tool call so an Esc press can cut the *wait* short.

    mcp_manager.call_tool() blocks on a pipe read or an HTTP request with no
    cancellation point of its own — unlike handle_run_command, which polls
    _CURRENT_INTERRUPT directly inside its own wait loop. Before this, a slow
    MCP call (a browser action, a big fetch, OCR) was the one kind of "the
    agent is doing something" that Esc could not touch: the turn-level check
    only runs *between* tool calls, so the user sat through the entire call
    regardless of how many times they pressed it.

    The call still runs to completion on a background thread rather than
    being killed — an MCP server is a single long-lived process for the whole
    session (unlike run_command's one-shot subprocess), so tearing it down
    here would break every later call to it, and for playwright specifically
    would drop whatever page state the user was relying on. The one cost of
    abandoning the wait instead: mcp_manager.py takes a per-server lock for
    the duration of a call, so a server that never answers stays locked for
    the rest of the session. That is strictly better than today's failure
    mode for the same case, which is the whole agent hanging forever with no
    way to interrupt it at all.
    """
    outcome: list[str] = []

    def _run() -> None:
        outcome.append(mcp_manager.call_tool(name, params))

    thread = threading.Thread(target=_run, daemon=True)
    thread.start()
    while thread.is_alive():
        if _CURRENT_INTERRUPT is not None and _CURRENT_INTERRUPT.is_set():
            return (f"[interrupted] '{name}' was still waiting on its MCP "
                    f"server when Esc was pressed. The call keeps running in "
                    f"the background; this result was not.")
        thread.join(timeout=0.25)
    return outcome[0] if outcome else f"Error: '{name}' produced no result."

# ---------------------------------------------------------------------------
# Permission system
# ---------------------------------------------------------------------------

def check_permission(name: str, params: dict) -> bool:
    """Legacy entry point. The agent loop now asks through a
    PermissionResponder (core/permissions.py); this remains for any direct
    caller and shares the same session-scoped ApprovalStore.

    Note it no longer downgrades RISK_LEVELS on "always": that turned one
    approval of, say, `git status` into a blanket grant on every future
    run_command. Approval is scoped to the exact call the user saw.
    """
    if YOLO_MODE:
        return True  # YOLO mode approves everything
    if APPROVALS.is_approved(name, params):
        return True
    risk = risk_for(name, params)
    if risk == "low" and AUTO_APPROVE_LOW:
        return True

    from core.events import PermissionNeeded
    decision = TerminalAdapter().ask(
        PermissionNeeded("", name, params, risk))
    if decision == "always_allow_this_call":
        APPROVALS.approve(name, params)
        return True
    return decision == "allow"

# ---------------------------------------------------------------------------
# System prompt + project context re-injection
# ---------------------------------------------------------------------------

BASE_PROMPT = """You are a coding assistant that helps with software engineering tasks.
You have tools for reading, writing, and editing files, listing directories,
running shell commands, searching code, and saving memories.

Rules:
- Always read a file before editing it.
- Prefer edit_file over write_file for existing files.
- Keep responses concise. Focus on code, not lengthy explanations.
- Use absolute or project-relative paths.
- If a task is done, stop calling tools and summarize.
- What you have been asked to remember reaches you already: standing rules in the imperative section, everything else retrieved per message. There is no index to consult.

Clarifying questions:
- When a task has more than one reasonable interpretation, involves a choice
  with real consequences (which library, which approach, overwrite vs. merge,
  destructive vs. safe), or is missing information only the user has, call
  ask_user_question instead of guessing. Do not use it for something you can
  answer yourself by reading the code or the file system.
- Do not use it for a simple yes/no you could infer from context, and do not
  ask more questions than the task actually needs -- one well-chosen question
  beats four.
- Each question needs 2-4 short, genuinely distinct options (the user can
  always type a custom answer instead of picking one) and, for a "choose any
  that apply" question, multiSelect: true. You can ask several questions in
  one call when they are all blocking the same next step.
- The tool blocks until the user answers; use the returned answer(s) directly
  in the same turn rather than asking again or second-guessing the choice.

Memory and self-improvement:
- When the user states a durable preference, correction, or rule -- "always...",
  "never...", "from now on...", "I prefer...", or a correction to something you
  just did -- call save_memory in that same turn. Do not wait to be asked to
  remember it, and do not merely acknowledge it in prose: an acknowledgement is
  forgotten when the session ends, a saved memory is not.
- Phrase what you save the way the user phrased it. A rule meant for every reply
  ("always end with the date") must be saved in those words, because that is what
  makes it apply on every later turn rather than only when the topic comes up.
- Any standing rule already in your system prompt applies to EVERY reply in this
  session -- including replies about a completely unrelated topic, and including
  turns that are mostly tool calls. Following it on turn 1 and forgetting it by
  turn 20 is the failure to avoid.
- Before producing substantial work for a returning user, apply what you have
  been told about their conventions, terminology and formatting, rather than a
  generic default.
- A denied tool call will stay denied. Do not re-issue it with different wording.
- Reply in the language the user wrote in.
- Never say a file was created, saved, edited, or a task was completed unless
  you actually called the tool for it in this conversation and saw a
  successful result. If you only planned or described an action without
  calling the tool, say exactly that -- do not describe it as done. The turn
  ends the moment your reply has no tool call in it, so if you write "let me
  create..." or "I'll save this as...", call the tool in that same turn
  instead of stopping after the sentence.
- After creating or writing a file the user asked for, verify it exists
  (read_file or list_files) before telling them it is done. State the exact
  path you verified, not just the path you intended to use.
- A helper script you write only to carry out a step -- a converter, a
  generator, glue you will not need again -- is not a deliverable. Only what
  the user asked for belongs in their repo. Write it to `_scratch/` in the
  project root, which exists for exactly this, or skip the file and pass the
  code to `run_command` directly.
- Do NOT try to write to the system temp directory or anywhere under
  `~/.tomas`. Both are outside the writable sandbox and the call will be
  refused; `_scratch/` is the one scratch location you have.

Large deliverables:
- You have a per-reply output limit. A long document -- a methodichka, a
  report, a multi-section guide -- will not fit in one reply, and a reply that
  hits the limit is thrown away: the user gets nothing at all.
- So do not compose a long document in your reply text. Build it on disk,
  section by section: write_file the first section, then edit_file to append
  each following one. Each tool call is a separate turn with its own budget, so
  a document of any length is reachable this way, and progress survives even if
  a later turn fails.
- Say which section you are on as you go, so the user can see progress rather
  than waiting through several silent turns.
- Only after the file is complete and you have verified it with read_file do
  you convert it (e.g. word-docs convert_to_pdf) or summarise it to the user.
- .docx, .pdf, .xlsx and .pptx are structured containers, not text.
  write_file cannot produce them. Build a Word document with the word-docs
  MCP tools (create_document, then add_heading / add_paragraph / add_table),
  and get a PDF from convert_to_pdf on that document -- it keeps headings and
  bold, which re-typesetting the plain text into a new PDF throws away."""


def _classify_mcp_failures(failed: dict) -> tuple[list, list]:
    """Split MCP failures into 'needs credentials' and 'actually broken'.

    A 401 on an optional server is a permanent, expected state — it should
    not be styled like a crash.
    """
    needs_auth, broken = [], []
    for name, err in (failed or {}).items():
        text = str(err).lower()
        if any(m in text for m in ("401", "403", "unauthor", "forbidden",
                                   "authenticat", "api key", "token")):
            needs_auth.append(name)
        else:
            broken.append(name)
    return needs_auth, broken


def _environment_section() -> str:
    """Tell the model which shell it actually has.

    Without this it reached for `rm`, `ls` and `test -f` on cmd.exe — POSIX
    commands that simply do not exist there — and burned tool calls finding
    out. One line removes the whole class of mistake.
    """
    if sys.platform == "win32":
        return (
            "\n\n# Environment\n"
            "- OS: Windows. `run_command` runs through **cmd.exe**, not bash.\n"
            "- Use `dir`, `type`, `del`, `copy`, `move`, `findstr` — not "
            "`ls`, `cat`, `rm`, `cp`, `mv`, `grep`, `test -f`.\n"
            "- Prefer the built-in tools over shell equivalents: `list_files` "
            "over `dir`, `read_file` over `type`, `search_code` over `findstr`.\n"
            f"- Python interpreter: `{sys.executable}`"
        )
    return (
        "\n\n# Environment\n"
        f"- OS: {sys.platform}. `run_command` runs through a POSIX shell.\n"
        f"- Python interpreter: `{sys.executable}`"
    )

def _truncate_section(text: str, max_chars: int, label: str = "") -> str:
    """Truncate a section to max_chars, adding a notice if cut."""
    if len(text) <= max_chars:
        return text
    cut = text[:max_chars]
    # Try to cut at a line boundary for readability
    last_nl = cut.rfind('\n')
    if last_nl > max_chars * 0.8:
        cut = cut[:last_nl]
    notice = f"\n[... truncated: {len(text)} → {len(cut)} chars ...]"
    if label:
        notice = f"\n[... {label} truncated: {len(text)} → {len(cut)} chars ...]"
    return cut + notice


_warned_instructions: set = set()


def _warn_instructions_dropped(dropped: list, budget_chars: int) -> None:
    """Tell the *user* which instruction files did not fit, once per set.

    The old code appended "[... instructions truncated ...]" to the prompt and
    told nobody else. That notice went to the model, which cannot restore the
    missing text, cannot ask about it, and has no reason to mention it — so
    the person whose file was being ignored was the only one kept in the dark.
    """
    key = tuple(sorted(dropped))
    if key in _warned_instructions:
        return
    _warned_instructions.add(key)
    names = ", ".join(dropped)
    print(f'  {YELLOW}⚠{RESET}  instruction files not sent this session: '
          f'{BOLD}{names}{RESET}\n'
          f'     {DIM}They do not fit the {budget_chars:,}-char budget for a '
          f'{CONTEXT_WINDOW:,}-token window. Shorten them, or switch to a '
          f'model with a larger window.{RESET}')


# Maximum size for each system-prompt section (in characters)
#: Retained for callers and tests that ask "what is the ceiling": the *applied*
#: budget is now derived per-window by `core.budget.instructions_budget`, and
#: this is the largest it can ever be.
MAX_INSTRUCTIONS_CHARS = 40000      # AGENTS.md + CLAUDE.md + global instructions
MAX_DIRECTIVES_CHARS = 1000         # standing rules (bounded by MAX_DIRECTIVES too)
MAX_LEARNED_CHARS = 1500            # retrieved facts (bounded by k, not by store size)
MAX_SKILLS_CHARS = 4000             # skills section
#: Full body of a skill this message triggers. Raised twice, both times on the
#: same measurement: whatever SKILL.md does not state, the model goes and
#: reads out of the scripts, and that text then sits in the transcript and is
#: re-sent on every later step.
#:
#:   6000 -> 7200  the block schema. Three sessions each read 4-5 source
#:                 files — ~44,000 characters, ~11k tokens — to learn it.
#:   7200 -> 8600  the per-check thresholds. A later session read *ten*
#:                 files, 39,849 characters, and finished its turn having
#:                 produced no document at all.
#:
#: ~350 tokens of skill body, only when the skill triggers, against ~10k of
#: source reading is not a close call. It is still bounded by
#: MAX_TOTAL_SYSTEM_PROMPT; if this needs raising a third time, the answer is
#: probably to split the reference half into a file the model opens once.
#:
#:   8600 -> 9800  the retrieval gate. Triggers are substrings and they match
#:                 questions *about* a document as readily as orders to
#:                 produce one; the body is now prefixed with the instruction
#:                 to decide whether it applies, plus each skill's own
#:                 "use when / do NOT use when". ~700 characters of fixed
#:                 overhead that turns a keyword match into a decision the
#:                 model makes — cheaper than the extra round trip a separate
#:                 classifier call would cost, and it applies to every skill.
MAX_TRIGGERED_SKILL_CHARS = 9800
MAX_TOTAL_SYSTEM_PROMPT = 46000     # hard cap on the entire system prompt.
# ~11.5k tokens against the 200k context the default models actually have.
#
# Raised from 44,000 when the stable half reached 44,530 and the truncator
# started eating the end of the skills catalogue — the section that lets the
# model *find* a skill — to stay under a round number. The trade is stated two
# paragraphs down and had already been made: this content sits in the stable
# prefix, so it is served from cache rather than re-tokenised, and 2,000 more
# characters there cost almost nothing per turn. Trimming conventions to fit an
# arbitrary constant would have cost more.
#
# It is not headroom to spend freely. AGENTS.md and CLAUDE.md overlap by an
# estimated 2,000-3,000 characters in this repo; dedupe them before raising
# this again.
#
# Raised from 8000/28000 when CLAUDE.md started being loaded alongside AGENTS.md
# (it had been documented as loaded for a long time without being loaded at
# all). Together they are 16,259 chars here, so the old instructions budget was
# discarding 8,259 of them — half the project's conventions, silently, with the
# truncation notice as the only trace.
#
# Affordable because of where this content now sits: the instructions are part
# of the stable prefix, so they are re-read from cache rather than re-tokenised
# every turn. The per-turn cost of a larger stable prefix is close to nothing;
# the per-turn cost of a larger *volatile* tail is not, which is why those
# budgets are unchanged.
#
# Worth knowing: AGENTS.md and CLAUDE.md overlap substantially in this repo —
# both describe the tool layers, the prompt load order and the quirks. Deduping
# them would recover roughly 2,000-3,000 chars of prompt with no loss.


# The smallest the stable prefix may be squeezed to when the volatile tail is
# unusually large. BASE_PROMPT plus the environment section is what the agent
# *is*; below this the rules it operates by start disappearing.
MIN_STABLE_CHARS = 6000

#: Ceiling on the *stable* half, and the only thing the truncator reads.
#:
#: The cut has to be decided by a constant. It used to be
#: `MAX_TOTAL_SYSTEM_PROMPT - len(tail)`, which is message-dependent — so once
#: the stable half grew past the cap, every message trimmed it in a different
#: place. Measured after adding 2.3 KB to CLAUDE.md: "what time is it in
#: Tokyo" kept 43,822 characters of it and "зроби схожий файл" — whose tail
#: carries a triggered skill body — kept 35,230. Two prompts with no shared
#: prefix, and prefix caching matches on an exact byte prefix, so *every* turn
#: billed the whole conversation at full price. That is precisely the failure
#: the ordering rule in CLAUDE.md exists to prevent, arriving through the
#: truncation path instead of the ordering one.
#:
#: The tail then sits on *top* of this rather than inside it, so the worst
#: case is this plus the tail's own budgets (standing rules 1000 + retrieved
#: facts 1500 + a triggered skill body 8600) — about 14k tokens against the
#: 200k context the default models have. Subtracting the reserve instead was
#: tried and is worse: it caps the stable half at 32.5k whether or not a skill
#: triggered, and the first thing that falls off the end is the skills
#: catalogue — so the prompt loses the ability to *find* the skill in exchange
#: for room to quote one.
MAX_STABLE_PROMPT_CHARS = MAX_TOTAL_SYSTEM_PROMPT

# Cached stable prefix, and the filesystem signature it was built from.
_stable_prefix: Optional[str] = None
_stable_signature: Optional[tuple] = None


def _stable_inputs() -> list[Path]:
    """Every file whose contents can change the stable prefix."""
    import skills_manager
    paths = [PROJECT_DIR / name for name in
             (*instructions_manager.PROJECT_INSTRUCTION_FILES,
              "AGENT_INSTRUCTIONS.md", "BEHAVIOR.md")]
    paths.append(instructions_manager.PROJECT_INSTRUCTIONS_DIR
                 / f"{PROJECT_DIR.name}.md")
    try:
        paths.extend(sorted(instructions_manager.GLOBAL_INSTRUCTIONS_DIR.glob("*.md")))
    except OSError:
        pass
    paths.extend(skills_manager.SKILL_DIRS)
    return paths


def _budget_sections() -> frozenset:
    """Which system-prompt sections the active budget allows.

    Falls back to "everything" rather than "nothing": a budget that cannot be
    read must not silently strip the agent of its instructions.
    """
    try:
        return active_budget().enabled_sections
    except Exception:
        return frozenset(core_budget.SECTION_KEYS)


def _stable_fingerprint() -> tuple:
    """(path, mtime, size) for each input — cheap, and changes when they do.

    Rebuilding the stable half cost 18.8 ms and 98 `exists()` calls on every
    turn, re-reading files that change at most once a session. Keying on stat
    rather than a timer means an edit to AGENTS.md still takes effect on the
    very next message, which is what makes the file worth editing.
    """
    signature = []
    for path in _stable_inputs():
        try:
            st = path.stat()
            signature.append((str(path), int(st.st_mtime_ns), st.st_size))
        except OSError:
            signature.append((str(path), 0, -1))    # absent is a state too
    # Two sections of the stable half are switchable (instructions, skills
    # catalogue), and switching one moves no file on disk. Without this the
    # memoised prefix would outlive the setting that built it, and a user who
    # turned the catalogue off would watch the token count refuse to move.
    signature.append(("budget:sections", tuple(sorted(_budget_sections())), 0))
    return tuple(signature)


def invalidate_prompt_cache() -> None:
    """Drop the cached prefix. For callers that change instructions in-process."""
    global _stable_prefix, _stable_signature
    _stable_prefix = None
    _stable_signature = None


def build_system_prompt(user_message: str = "") -> str:
    """Build the system prompt for this turn, stable part first.

    `user_message` is what learned knowledge is retrieved against. An empty
    query falls back to the most recently confirmed facts, so callers that
    have no message yet still get something sensible.

    ── Why the order is what it is ──

    Everything that does not depend on `user_message` is emitted before
    anything that does. Prefix caching — DeepSeek's automatic kind, Anthropic's
    `cache_control`, OpenAI's — matches on an exact byte prefix, and the system
    prompt is serialised *before* the messages. So the first byte that differs
    from last turn ends the cache hit for the system prompt **and for the whole
    conversation history behind it**.

    The catalogue used to sit last, after two sections rebuilt per message.
    Measured on Zen/DeepSeek over five turns with a populated fact store: 52.0%
    of prompt tokens served from cache, 14,864 billed as new. With the stable
    content moved ahead of the volatile content, the same five turns cached
    83.9% and billed 4,752 — 68% less. In a long session the saving is larger
    still, because it is the history that stops being re-read.
    """
    # ── Stable: identical every turn until a file on disk changes ──
    global _stable_prefix, _stable_signature
    signature = _stable_fingerprint()
    if _stable_prefix is not None and signature == _stable_signature:
        stable = _stable_prefix
    else:
        # Which sections are allowed at all. Read once here rather than per
        # section, so the stable half is built against a single consistent
        # answer — and note `_stable_fingerprint` carries the same flags, or
        # a toggle would not survive the memoisation below.
        allowed = _budget_sections()
        stable = BASE_PROMPT + _environment_section()
        # project-level instructions from AGENTS.md / agent.md / CLAUDE.md
        # plus ~/.tomas/instructions/
        # Budgeted as a share of the real window and cut by whole files —
        # never at a character offset. See `instructions_manager.fit_instructions`
        # for why the old blunt slice was the bug behind "the agent does not
        # follow my AGENT.md": the cut landed mid-document and the only party
        # told about it was the model.
        if "instructions" in allowed:
            import core.budget as _budget
            import instructions_manager as _im
            budget_chars = _budget.instructions_budget(CONTEXT_WINDOW)
            instructions_section, dropped = _im.fit_instructions(
                _im.instruction_parts(PROJECT_DIR), budget_chars)
            if dropped:
                _warn_instructions_dropped(dropped, budget_chars)
        else:
            instructions_section = ""
        if instructions_section:
            stable += f"\n\n{instructions_section}"
        # legacy support: AGENT_INSTRUCTIONS.md or BEHAVIOR.md
        for candidate in [PROJECT_DIR / "AGENT_INSTRUCTIONS.md",
                          PROJECT_DIR / "BEHAVIOR.md"]:
            if candidate.exists():
                legacy = candidate.read_text(encoding="utf-8")
                legacy = _truncate_section(legacy, 2000, candidate.name)
                stable += f"\n\n# Agent Instructions ({candidate.name})\n{legacy}"
                break
        # installed skills — budgeted by whole entries, not by slicing the
        # joined string at a character offset (which used to cut mid-skill-
        # name). Names only; it does not vary with the message, so it belongs
        # up here where it can be cached rather than re-sent every turn.
        skills_section = (build_skills_section(max_chars=MAX_SKILLS_CHARS)
                          if "skills_catalogue" in allowed else "")
        if skills_section:
            stable += f"\n\n{skills_section}"
        # Capped here, once, rather than per message: the cut is decided by a
        # constant (TAIL_RESERVE_CHARS), so doing it inside the memoised half
        # keeps the cached value equal to what is actually emitted. It used to
        # be applied below against `MAX_TOTAL_SYSTEM_PROMPT - len(tail)`, and a
        # message-dependent cut is a message-dependent prefix.
        room = max(MIN_STABLE_CHARS, MAX_STABLE_PROMPT_CHARS)
        if len(stable) > room:
            stable = _truncate_section(stable, room, "system prompt")
        _stable_prefix = stable
        _stable_signature = signature

    # ── Volatile: rebuilt against this turn's message ──
    tail = ""
    # Standing rules — always on, never retrieved. These lead the tail for a
    # measured reason. A 30-turn session ran two standing rules at once: "end
    # with My Lord" from AGENT.md (static, imperative heading) was obeyed 29/29;
    # "append the date" from the fact store (retrieved, filed under "What I've
    # learned") was obeyed 0/29 — while being present in the prompt on all 29
    # turns. The rule was never the problem; the channel was. Directives keep
    # their own imperative heading and still precede the retrieved facts.
    tail_allowed = _budget_sections()
    try:
        directives = (learning.directives_for_prompt()
                      if "standing_rules" in tail_allowed else "")
        if directives:
            directives = _truncate_section(directives, MAX_DIRECTIVES_CHARS,
                                           "standing rules")
            tail += (
                "\n\n# Standing rules from the user — these apply to EVERY reply\n"
                "You MUST follow every rule below on every single turn. They "
                "apply even when the current message is about something else "
                "entirely, and even on turns where you mostly call tools. These "
                "are not background information about the user; they are "
                "instructions to you, and the user checks whether you followed "
                "them.\n\n"
                f"{directives}")
    except Exception:
        pass
    # ── What the agent has learned — retrieved, not dumped ──
    # This replaces the old memory-index dump, the notes dump and the tips
    # block. Those three grew with everything ever learned until entries
    # silently fell off the end of the budget; retrieval keeps the prompt
    # flat in size no matter how much is stored.
    try:
        learned = (learning.recall(user_message, k=5)
                   if "learned_facts" in tail_allowed else "")
        if learned:
            learned = _truncate_section(learned, MAX_LEARNED_CHARS, "learned")
            # Heading rewritten from "What I've learned about this user and
            # project". That phrasing read as a dossier and the model treated it
            # as trivia. Same data, actionable framing.
            tail += ("\n\n# Context retrieved for this message — apply what is "
                     "relevant\n"
                     f"{learned}")
    except Exception:
        pass
    # A skill this message triggers goes in full. The catalogue above only
    # names skills — enough to pick one, useless for following one — and
    # nothing ever read `triggers`, so a procedure written for a job never
    # reached the model while it was doing that job.
    try:
        triggered = (build_triggered_skills(user_message, MAX_TRIGGERED_SKILL_CHARS)
                     if "triggered_skills" in tail_allowed else "")
        if triggered:
            tail += f"\n\n{triggered}"
    except Exception:
        pass
    # NOTE: the self-improvement tips/session-context block used to be injected
    # here. It was template text addressed to a human developer ("Consider
    # creating shortcuts or aliases for this tool") that consumed context and
    # changed nothing about the model's behaviour. Reflection replaces it; the
    # generator code is still in self_improve.py pending deletion.

    # The total cap is applied to the stable half where it is built and
    # memoised, not here: it is decided by a constant, so it belongs with the
    # thing it bounds. The tail is never trimmed for it — it is what the model
    # must act on *now*, the user's standing rules and the procedure for the
    # job in hand, and trimming from the end of the joined string (as this
    # once did) made the triggered skill body the first thing dropped.
    return stable + tail

# ---------------------------------------------------------------------------
# Three-layer memory system
# ---------------------------------------------------------------------------

def load_memory_index() -> str:
    idx = MEMORY_DIR / "MEMORY.md"
    if idx.exists():
        return idx.read_text(encoding="utf-8")
    return ""


# ── Where a remembered thing belongs ───────────────────────────────────────
#
# Four stores exist and they are not interchangeable — they differ in what
# they cost per turn and in whether the model sees them at all:
#
#   instruction  ~/.tomas/instructions/  stable prompt half, memoised.
#                Read on EVERY turn and billed once, because the stable half
#                is what prefix caching matches on. Permanent identity.
#   rule         directive fact          volatile tail, every turn, capped at
#                MAX_DIRECTIVES and MAX_DIRECTIVE_CHARS. Unconditional.
#   fact         explicit fact           volatile tail, top-5 by relevance to
#                the message. Conditional preferences.
#   note         ~/.tomas/self-notes/    NOT in the prompt at all — a
#                user-facing scratchpad reachable through /notes.
#
# Routing used to be a two-way test (`looks_like_directive` → directive, else
# explicit), which put "your name is TOMAS" and "prefer PowerShell on
# Windows" in the same channel and gave the first one a per-turn price it
# never needed to pay. The distinction that matters is not conditional vs
# unconditional; it is *permanent* vs *revisable*.

#: Identity: who the agent is, what it calls the user, what language it
#: defaults to, how it signs off. These change about once and then never, so
#: they belong in the half of the prompt that is cached rather than in the
#: ten-slot budget that is re-sent every turn.
_IDENTITY_RE = re.compile(
    r"(?i)\b(?:your\s+name\s+is|call\s+me|address\s+(?:me|the\s+user)\s+as"
    r"|refer\s+to\s+me\s+as|sign\s+(?:off|every\s+\w+)\s+with"
    r"|end\s+(?:every|each|all)\s+\w+\s+with"
    r"|(?:default|always\s+(?:reply|respond|write|answer))\s+(?:language|in)"
    r"|звертайся\s+до\s+мене|називай\s+мене|твоє?\s+ім'?я"
    r"|закінчуй\s+(?:кожн\w+|усі)|підписуй"
    r"|обращайся\s+ко\s+мне|называй\s+меня|твоё?\s+имя)\b")

#: Long or multi-paragraph material is reference, not a rule. Putting it in
#: the tail would spend the standing-rule budget on something no turn needs
#: in full.
_NOTE_MIN_CHARS = 400


def route_memory(text: str) -> tuple[str, str]:
    """Which store this belongs in, and why. Returns (store, reason).

    Pure and cheap on purpose: the model may override it, the user can see
    it, and a wrong guess is one `/rules` command away from being fixed.
    """
    text = (text or "").strip()
    if not text:
        return STORE_FACT, "empty"
    unconditional = learning.looks_like_directive(text)
    if _IDENTITY_RE.search(text):
        return (STORE_INSTRUCTION,
                "identity or form of address — permanent, so it goes in the "
                "cached instructions rather than spending a standing-rule "
                "slot on every turn")
    if len(text) >= _NOTE_MIN_CHARS or text.count("\n") >= 3:
        return (STORE_NOTE,
                "long reference material — too big for a prompt section that "
                "is re-sent every turn")
    if unconditional:
        return (STORE_RULE,
                "unconditional wording ('always'/'never'/'from now on') — it "
                "applies to every turn, so it goes in the standing rules")
    return (STORE_FACT,
            "a conditional preference — retrieved when the topic comes up "
            "rather than repeated on turns it has nothing to do with")


# ── /rules ─────────────────────────────────────────────────────────────────
#
# A rule the user set lives in one of two places depending on how permanent
# it is, and before this they could only *see* one of them and only delete
# from it. Managing rules across both stores from one command is the point:
# the split is an implementation detail of where the prompt puts them, not
# something the user should have to hold in their head.
#
# Instruction rules are addressed as i1, i2 … (position in the file, which is
# what a person can see); directives keep their content ids, which is what
# reflection and /forget already use.

def _rules_snapshot() -> tuple[list[str], list[dict]]:
    try:
        managed = instructions_manager.read_managed_rules()
    except Exception:
        managed = []
    try:
        directives = learning.load_directives()
    except Exception:
        directives = []
    return managed, directives


def _resolve_rule(ref: str):
    """('instruction', index) | ('directive', fact) | (None, None)."""
    ref = (ref or "").strip()
    managed, _ = _rules_snapshot()
    if re.fullmatch(r"(?i)i\d+", ref):
        index = int(ref[1:]) - 1
        if 0 <= index < len(managed):
            return "instruction", index
        return None, None
    fact = learning.find_fact(ref)
    if fact and fact.get("kind") == learning.KIND_DIRECTIVE:
        return "directive", fact
    return None, None


def _handle_rules(rest: str) -> str:
    verb, _, argument = rest.partition(" ")
    verb, argument = verb.strip().lower(), argument.strip()

    if verb in ("add", "new", "+"):
        if not argument:
            return (f'  {YELLOW}Usage:{RESET} {CYAN}/rules add <text>{RESET}\n'
                    f'  {DIM}Or just type{RESET} {CYAN}#<text>{RESET}')
        return _rules_add(argument)

    if verb in ("edit", "change", "set"):
        ref, _, new_text = argument.partition(" ")
        if not ref or not new_text.strip():
            return (f'  {YELLOW}Usage:{RESET} {CYAN}/rules edit <id> <new text>{RESET}')
        kind, target = _resolve_rule(ref)
        if kind == "instruction":
            managed, _ = _rules_snapshot()
            was = managed[target]
            managed[target] = new_text.strip()
            instructions_manager.write_managed_rules(managed)
            invalidate_prompt_cache()
            return (f'  {GREEN}✓{RESET} Rule {ref} updated.\n'
                    f'     {DIM}was:{RESET} {was[:70]}\n'
                    f'     {DIM}now:{RESET} {new_text.strip()[:70]}')
        if kind == "directive":
            updated = learning.edit_fact(target["id"], new_text.strip())
            if not updated:
                return f'  {RED}✗{RESET} Could not update {ref}.'
            demoted = ("\n     " + DIM + "No longer worded unconditionally — "
                       "it is now a retrieved preference, not a standing rule."
                       + RESET) if updated.get("kind") != learning.KIND_DIRECTIVE else ""
            return (f'  {GREEN}✓{RESET} Rule {updated["id"]} updated.\n'
                    f'     {DIM}was:{RESET} {target.get("fact", "")[:70]}\n'
                    f'     {DIM}now:{RESET} {updated.get("fact", "")[:70]}{demoted}')
        return f'  {DIM}No rule {ref}. Run {RESET}{CYAN}/rules{RESET}{DIM} for the ids.{RESET}'

    if verb in ("rm", "remove", "delete", "forget", "-"):
        ref = argument.split(maxsplit=1)[0] if argument else ""
        if not ref:
            return (f'  {YELLOW}Usage:{RESET} {CYAN}/rules rm <id>{RESET}')
        kind, target = _resolve_rule(ref)
        if kind == "instruction":
            managed, _ = _rules_snapshot()
            removed = managed.pop(target)
            instructions_manager.write_managed_rules(managed)
            invalidate_prompt_cache()
            return (f'  {GREEN}✓{RESET} Rule removed: {removed[:70]}\n'
                    f'  {DIM}Ids below it shift up; run /rules to see them.{RESET}')
        if kind == "directive":
            removed = learning.forget(target["id"])
            # Precise about what a tombstone does: it stops *reflection*
            # inferring the rule again. Stating it yourself still brings it
            # back, which is the behaviour you want — otherwise one /forget
            # would permanently blacklist a rule you later change your mind on.
            return (f'  {GREEN}✓{RESET} Rule removed: '
                    f'{(removed or {}).get("fact", "")[:70]}\n'
                    f'  {DIM}Reflection will not infer it again. Telling me the '
                    f'rule yourself still restores it.{RESET}')
        return f'  {DIM}No rule {ref}. Run {RESET}{CYAN}/rules{RESET}{DIM} for the ids.{RESET}'

    if verb in ("help", "?"):
        return _rules_help()

    if verb:
        # Unrecognised verb, but the user clearly typed a rule: "/rules always
        # answer in Ukrainian" should not be an error message.
        return _rules_add(rest)

    return _rules_list()


def _rules_add(text: str) -> str:
    text = text.strip()
    store, why = route_memory(text)
    # `/rules add` is the user calling it a rule. Honour that: the only choice
    # left is *which* rule store, never demotion to a retrieved preference.
    if store not in (STORE_INSTRUCTION, STORE_RULE):
        store = STORE_RULE
        why = "you added it with /rules, so it is kept as a standing rule"
    if store == STORE_INSTRUCTION:
        if not instructions_manager.add_managed_rule(text):
            return f'  {DIM}Already a rule.{RESET}'
        invalidate_prompt_cache()
        where = "instructions — applies every turn, costs nothing per turn"
    else:
        directives = learning.load_directives()
        if len(directives) >= learning.MAX_DIRECTIVES:
            return (f'  {RED}✗{RESET} {learning.MAX_DIRECTIVES} standing rules '
                    f'already — that cap is real, they are re-sent every turn.\n'
                    f'  {DIM}Drop one with{RESET} {CYAN}/rules rm <id>{RESET}'
                    f'{DIM}, or reword this as identity so it lands in the '
                    f'cached instructions.{RESET}')
        learning.remember(learning.KIND_DIRECTIVE, text,
                          evidence="added with /rules add", scope="global")
        where = f"standing rules — one of {learning.MAX_DIRECTIVES} slots"
    return (f'  {GREEN}✓{RESET} Rule added to {where}.\n'
            f'  {DIM}{why}.{RESET}')


def _rules_help() -> str:
    return "\n".join([
        f'  {BOLD}Managing rules{RESET}',
        f'  {DIM}{"─" * 60}{RESET}',
        f'  {CYAN}/rules{RESET}                    — list every rule, with ids',
        f'  {CYAN}/rules add <text>{RESET}         — add one ({CYAN}#<text>{RESET} does the same)',
        f'  {CYAN}/rules edit <id> <text>{RESET}   — reword one, keeping its id',
        f'  {CYAN}/rules rm <id>{RESET}            — remove one',
        '',
        f'  {DIM}Rules live in two places, and /rules manages both:{RESET}',
        f'  {DIM}  i1, i2 …  instructions — permanent identity, cached, free per turn{RESET}',
        f'  {DIM}  hex ids   standing rules — re-sent every turn, max '
        f'{learning.MAX_DIRECTIVES}{RESET}',
    ])


def _rules_list() -> str:
    managed, directives = _rules_snapshot()
    if not managed and not directives:
        return (f'  {DIM}No rules yet. Tell me one ("always end every reply '
                f'with the date"), type{RESET} {CYAN}#<text>{RESET}{DIM}, or run{RESET} '
                f'{CYAN}/rules add <text>{RESET}{DIM}.{RESET}')

    today = time.strftime("%Y-%m-%d")
    lines = [f'  {BOLD}Rules ({len(managed) + len(directives)}){RESET}',
             f'  {DIM}{"─" * 60}{RESET}']
    if managed:
        lines.append(f'  {BOLD}Instructions{RESET} {DIM}— permanent, cached, '
                     f'free per turn{RESET}')
        for i, text in enumerate(managed, 1):
            lines.append(f'  {text[:96]}')
            lines.append(f'     {DIM}i{i}{RESET}')
    if directives:
        if managed:
            lines.append('')
        conflicts = learning.find_conflicts(directives)
        in_conflict = {i for pair in conflicts for i in pair}
        lines.append(f'  {BOLD}Standing rules{RESET} {DIM}— re-sent every turn '
                     f'({len(directives)}/{learning.MAX_DIRECTIVES}){RESET}')
        for fact in directives:
            text = (fact.get("fact") or "").replace("\n", " ")[:96]
            fid = fact.get("id", "?")
            lines.append(f'  {text}')
            notes = []
            expired = learning.stale_dates(fact.get("fact", ""), today)
            if expired:
                notes.append(f'names {", ".join(expired)}, today is {today}')
            if fid in in_conflict:
                notes.append('conflicts with another rule')
            note = f'  {YELLOW}⚠ {"; ".join(notes)}{RESET}' if notes else ''
            lines.append(f'     {DIM}{fid}{RESET}{note}')
    lines.append('')
    lines.append(f'  {DIM}Manage:{RESET} {CYAN}/rules add|edit|rm{RESET}'
                 f'{DIM} — see{RESET} {CYAN}/rules help{RESET}')
    return "\n".join(lines)

def save_memory(key: str, description: str, content: str,
                store: str = "") -> tuple[str, str]:
    """Persist an explicit "remember this" from the user. Returns (store, why).

    Writes the markdown file (still useful to read and edit by hand) and then
    routes the same content to whichever store it belongs in. The user said it
    outright — no inference — so it goes active immediately rather than
    through the evidence gate.

    Which store it lands in is what determines whether it gets followed, and
    at what price. See `route_memory`: identity goes to the cached
    instructions, unconditional rules to the standing-rule budget, conditional
    preferences to relevance-retrieved facts, long reference material to a
    note. `store` overrides the guess when the caller knows better.
    """
    MEMORY_DIR.mkdir(parents=True, exist_ok=True)
    filename = f"{key.replace(' ', '-').lower()}.md"
    filepath = MEMORY_DIR / filename
    filepath.write_text(
        f"---\nname: {key}\ndescription: {description}\n---\n\n{content}",
        encoding="utf-8",
    )
    idx = MEMORY_DIR / "MEMORY.md"
    lines = idx.read_text(encoding="utf-8").splitlines() if idx.exists() else []
    new_entry = f"- [{key}]({filename}) - {description}"
    updated = False
    for i, line in enumerate(lines):
        if filename in line:
            lines[i] = new_entry
            updated = True
            break
    if not updated:
        lines.append(new_entry)
    idx.write_text("\n".join(lines) + "\n", encoding="utf-8")

    # Classify on the whole record: the rule wording can live in either half
    # ("Never delete logs" / "see the runbook"). But when the content already
    # *is* the rule — a whole sentence — prefixing it with the index summary
    # produces "How to address the user: Address the user as MY KING", which
    # is what then gets injected into every turn and read back to the user.
    content_is_a_sentence = (content.strip().endswith((".", "!", "?"))
                             and len(content.strip()) > 25)
    body = (content.strip() if content_is_a_sentence
            else f"{description}: {content}".strip(": "))
    chosen, why = route_memory(body)
    if store in MEMORY_STORES:
        chosen, why = store, "you asked for this store"

    if chosen == STORE_INSTRUCTION:
        try:
            instructions_manager.add_managed_rule(body)
            invalidate_prompt_cache()
            return chosen, why
        except Exception:
            # A file that cannot be written must not lose the memory: fall
            # through to the store that needs nothing but the fact log.
            chosen, why = STORE_RULE, "instructions file unwritable — kept as a standing rule"

    if chosen == STORE_NOTE:
        try:
            self_notes.create_note(key, body, tags=["remembered"])
            return chosen, why
        except Exception:
            chosen, why = STORE_FACT, "note store unavailable — kept as a fact"

    if learning.is_enabled():
        try:
            kind = (learning.KIND_DIRECTIVE if chosen == STORE_RULE
                    else learning.KIND_EXPLICIT)
            learning.remember(kind, body,
                              evidence=f"user asked to remember '{key}'",
                              scope="global")
        except Exception:
            pass
    return chosen, why


# ---------------------------------------------------------------------------
# Learning system wiring
# ---------------------------------------------------------------------------

def _learning_call_model(model: str, system: str, messages: list,
                         max_tokens: int) -> str:
    """Adapter so learning/ can reach the model without importing agent."""
    resp = _get_client().messages.create(
        model=model, max_tokens=max_tokens, system=system, messages=messages,
    )
    return "".join(b.text for b in resp.content if hasattr(b, "text"))


def init_learning() -> None:
    """Scope the store to this project and import the pre-Phase-3 stores."""
    try:
        learning.set_project(PROJECT_DIR)
        imported = learning.migrate_legacy_stores()
        if imported:
            print(f'  {GREEN}✦{RESET} {DIM}Imported {imported} existing '
                  f'memories/notes into the learning store{RESET}')
        repaired = learning.repair_frontmatter_facts()
        if repaired.get("repaired"):
            print(f'  {GREEN}✦{RESET} {DIM}Repaired {repaired["repaired"]} '
                  f'learned fact(s) that had note metadata baked in{RESET}')
    except Exception:
        pass


def reflect_on_session_end(messages: list) -> None:
    """Learn from the finished session. Never raises, never blocks the user.

    Runs and persists silently — results surface on request via
    /self-improve facts and /self-improve reflect, not as unprompted
    chatter at the moment a session ends.
    """
    if not learning.is_enabled():
        return
    try:
        # Corrections first: deterministic, free, and independent of whether
        # the reflection model call succeeds. A provider outage at session end
        # should not cost the user the one signal they gave explicitly.
        learning.promote_corrections(messages)
    except Exception:
        pass
    try:
        learning.run_session_reflection(messages, call_model=_learning_call_model)
        for scope in ("global", "project"):
            learning.decay(scope)
    except Exception:
        pass


def _ago(timestamp: float) -> str:
    if not timestamp:
        return "never"
    seconds = max(0, time.time() - timestamp)
    for unit, size in (("d", 86400), ("h", 3600), ("m", 60)):
        if seconds >= size:
            return f"{int(seconds // size)}{unit} ago"
    return "just now"


def _fmt_duration(seconds: float) -> str:
    """How long a turn took, in the shortest form that stays readable."""
    seconds = max(0.0, seconds)
    if seconds < 60:
        return f"{seconds:.1f}s"
    minutes, rest = divmod(int(seconds), 60)
    return f"{minutes}m {rest:02d}s"


def _render_learned_facts() -> str:
    """Everything the agent believes, with the evidence behind it."""
    try:
        rows = []
        for scope in ("global", "project"):
            for fact in learning.load_facts(scope):
                fact = dict(fact)
                fact["scope"] = scope
                rows.append(fact)
    except Exception as e:
        return f'  {RED}✗{RESET} Could not read the learning store: {e}'

    if not rows:
        return (f'  {DIM}Nothing learned yet. Facts appear here as sessions '
                f'accumulate evidence.{RESET}')

    order = {"active": 0, "candidate": 1, "observed": 2}
    rows.sort(key=lambda f: (order.get(f.get("status"), 3),
                             -f.get("evidence_count", 0)))
    active = sum(1 for f in rows if f.get("status") == "active")

    lines = [
        f'  {BOLD}What I have learned{RESET}',
        f'  {DIM}{"─" * 56}{RESET}',
        f'  {DIM}{active} active · {len(rows) - active} still gathering evidence '
        f'· promotes at {learning.PROMOTE_AT}{RESET}',
        '',
    ]
    for fact in rows:
        status = fact.get("status", "observed")
        mark = {"active": f'{GREEN}●{RESET}',
                "candidate": f'{YELLOW}◐{RESET}'}.get(status, f'{DIM}○{RESET}')
        scope = fact.get("scope", "global")
        text = (fact.get("fact") or "").replace("\n", " ")[:100]
        lines.append(f'  {mark} {text}')
        lines.append(f'      {DIM}{fact.get("id", "?")} · {scope} · '
                     f'{fact.get("evidence_count", 0)}× · '
                     f'confirmed {_ago(fact.get("last_seen", 0))}{RESET}')
        evidence = (fact.get("evidence") or [])
        if evidence:
            lines.append(f'      {DIM}└ {str(evidence[-1])[:90]}{RESET}')
    lines.append('')
    lines.append(f'  {DIM}/forget <id> removes one permanently.{RESET}')
    return '\n'.join(lines)


def _render_reflection_log(limit: int = 5) -> str:
    """What reflection would have learned — the shadow-mode review screen."""
    path = learning.LEARNED_DIR / "reflection-log.jsonl"
    if not path.exists():
        return (f'  {DIM}No reflection runs yet. Reflection happens when a '
                f'session ends.{RESET}')
    try:
        entries = [json.loads(ln) for ln in
                   path.read_text(encoding="utf-8").splitlines() if ln.strip()]
    except (OSError, json.JSONDecodeError) as e:
        return f'  {RED}✗{RESET} Could not read the reflection log: {e}'
    if not entries:
        return f'  {DIM}Reflection log is empty.{RESET}'

    lines = [
        f'  {BOLD}Reflection log{RESET} {DIM}(last {min(limit, len(entries))} '
        f'of {len(entries)}){RESET}',
        f'  {DIM}{"─" * 56}{RESET}',
    ]
    for entry in entries[-limit:]:
        result = entry.get("result", {})
        lines.append(f'  {DIM}{_ago(entry.get("at", 0))} · mode={entry.get("mode")} '
                     f'· {entry.get("signals", 0)} correction signals{RESET}')
        for key, label in (("user_preferences", "preference"),
                           ("corrections", "lesson"),
                           ("project_notes", "project"),
                           ("skill_candidates", "skill")):
            for item in (result.get(key) or []):
                text = (item.get("fact") or item.get("lesson")
                        or item.get("name") or "")
                if text:
                    lines.append(f'      {DIM}[{label}]{RESET} {str(text)[:90]}')
        if not any(result.get(k) for k in
                   ("user_preferences", "corrections", "project_notes",
                    "skill_candidates")):
            lines.append(f'      {DIM}(learned nothing — the correct answer '
                         f'for most sessions){RESET}')
    mode_now = learning.reflect_mode()
    if mode_now == "shadow":
        lines.append('')
        lines.append(f'  {DIM}Shadow mode: nothing above was written to the '
                     f'store. Unset TOMAS_REFLECT to enable (active is the '
                     f'default).{RESET}')
    elif mode_now == "off":
        lines.append('')
        lines.append(f'  {DIM}Reflection is off (TOMAS_REFLECT=off). Nothing is '
                     f'learned from finished sessions.{RESET}')
    return '\n'.join(lines)

# ---------------------------------------------------------------------------
# Context management — auto-compaction
# ---------------------------------------------------------------------------

def estimate_tool_tokens(tools: list) -> int:
    """Token cost of a tool block, counted as the JSON it is serialised to."""
    try:
        return int(sum(len(json.dumps(t)) for t in tools) / CHARS_PER_TOKEN_JSON)
    except Exception:
        return 0


def _estimate_tokens(messages: list) -> int:
    """Rough token estimate for a message list."""
    total_chars = 0
    for m in messages:
        content = m.get("content", "")
        if isinstance(content, list):
            # tool_results: each item has a 'content' field
            for item in content:
                if isinstance(item, dict) and "content" in item:
                    total_chars += len(str(item["content"]))
                else:
                    total_chars += len(str(item))
        else:
            total_chars += len(str(content))
    return total_chars // CHARS_PER_TOKEN_PROSE

def _estimate_system_prompt_tokens(system_prompt: str) -> int:
    """Estimate tokens for the system prompt string."""
    return len(system_prompt) // CHARS_PER_TOKEN_PROSE

# ── Pruning old tool results ───────────────────────────────────────────
# Measured on a real 55-call session: history was 48,655 tokens of the ~60,000
# sent per call, and 29,394 of those — 60% of the history — were tool results.
# A file read on turn 2 was still being re-sent on turn 50. The model has long
# stopped reading it; the transcript is paying for it every single turn.

#: How many of the most recent tool-result batches are kept verbatim.
#:
#: Counted in batches, not user turns. The first version counted user turns and
#: found nothing to prune in the very session that motivated the feature: it
#: made 56 tool calls inside *two* user turns, so "older than 3 turns" never
#: matched anything. Long agentic work happens *within* a turn, which is
#: exactly where the history piles up.
TOOL_RESULT_KEEP_TURNS = _env_int("TOMAS_KEEP_RESULT_BATCHES", 8, minimum=1)

#: Results smaller than this are left alone — stubbing a 300-char result costs
#: information and saves nothing.
TOOL_RESULT_STUB_OVER = _env_int("TOMAS_STUB_RESULTS_OVER", 2_000, minimum=200)

#: Only prune once there is this much to gain, in characters.
#:
#: This is what keeps pruning cache-friendly. Stubbing rewrites a message in
#: the middle of the transcript, which invalidates the prefix cache from that
#: point on — so doing it every turn would trade one saving for a permanent
#: stream of cache misses. Pruning in occasional large batches pays that cost
#: rarely. Stubbing is monotonic (a stub is never re-stubbed), so everything
#: before the pruned point stays byte-identical.
TOOL_RESULT_PRUNE_AT = _env_int("TOMAS_PRUNE_RESULTS_AT", 20_000, minimum=2_000)

_PRUNED_MARK = "[older tool result released from context"


def _is_user_turn(message: dict) -> bool:
    """True for a real user message, as opposed to a tool-result carrier.

    Both have role "user" — the tool-result ones carry a list of tool_result
    blocks and no prose, which is what separates one turn from the next.
    """
    if message.get("role") != "user":
        return False
    content = message.get("content")
    if isinstance(content, str):
        return True
    if isinstance(content, list):
        return any(isinstance(b, dict) and b.get("type") != "tool_result"
                   for b in content)
    return False


def prunable_chars(messages: list, keep_turns: int = TOOL_RESULT_KEEP_TURNS,
                   stub_over: int = TOOL_RESULT_STUB_OVER) -> int:
    """How much could be reclaimed right now, without changing anything."""
    return sum(saved for _, _, saved in _prune_targets(messages, keep_turns,
                                                       stub_over))


def _prune_targets(messages: list, keep_batches: int, stub_over: int):
    """Yield (message_index, block_index, chars_saved) for each stubbable result.

    Walks backwards counting *batches of tool results* — one message carrying
    tool_result blocks is one batch — and leaves the newest `keep_batches`
    alone. Those are what the model is still reasoning over; everything behind
    them is reference material it can re-fetch far more cheaply than it can
    carry.
    """
    batches_seen = 0
    for i in range(len(messages) - 1, -1, -1):
        content = messages[i].get("content")
        if not isinstance(content, list):
            continue
        results = [(j, b) for j, b in enumerate(content)
                   if isinstance(b, dict) and b.get("type") == "tool_result"]
        if not results:
            continue
        batches_seen += 1
        if batches_seen <= keep_batches:
            continue
        for j, block in results:
            body = block.get("content")
            if not isinstance(body, str) or body.startswith(_PRUNED_MARK):
                continue
            if len(body) > stub_over:
                yield i, j, len(body)


def prune_tool_results(messages: list,
                       keep_turns: int = TOOL_RESULT_KEEP_TURNS,
                       stub_over: int = TOOL_RESULT_STUB_OVER,
                       prune_at: int = TOOL_RESULT_PRUNE_AT) -> int:
    """Replace the body of old, large tool results with a stub, in place.

    Returns the characters reclaimed (0 if it was not worth doing yet).

    The `tool_use`/`tool_result` pairing is preserved exactly — only the body
    text is replaced — because dropping the block outright leaves a dangling
    tool_call that upstreams reject.

    The stub says what was there and how to get it back, so this is a
    *release*, not a loss: the model can re-read the file if it turns out to
    still need it, which is far cheaper than carrying it for fifty turns.
    """
    targets = list(_prune_targets(messages, keep_turns, stub_over))
    if sum(saved for _, _, saved in targets) < prune_at:
        return 0
    reclaimed = 0
    for i, j, _ in targets:
        block = messages[i]["content"][j]
        body = block["content"]
        block["content"] = (
            f"{_PRUNED_MARK} — it was {len(body):,} characters. "
            f"Re-run the tool if you need it again.]")
        reclaimed += len(body) - len(block["content"])
    return reclaimed


def maybe_prune(messages: list) -> list:
    """Prune if it is worth it, and say so. Mirrors `maybe_compact`'s shape."""
    before = _estimate_tokens(messages)
    reclaimed = prune_tool_results(messages)
    if reclaimed > 0:
        print(f'  {DIM}[context] released {reclaimed:,} chars '
              f'(~{reclaimed // 4:,} tokens) of old tool results{RESET}')
        # Recorded on the same footing as compaction: pruning is the cheap
        # tier of the same job, and a session that pruned enough never to
        # need a summary should be able to show that.
        _record_compaction("prune", before, _estimate_tokens(messages), None)
    return messages


def maybe_compact(messages: list, system_prompt: str = "",
                  force: bool = False) -> list:
    """Compact the conversation if it's getting too large.

    Now accounts for system_prompt + tool definitions + max_tokens in the budget.

    `force` is what `/compact` passes, and it exists because the two are
    genuinely different questions. Turning automatic compaction off means "do
    not spend a model call on this without asking me" — it cannot also mean
    "and refuse when I do ask", which is what a single code path would have
    delivered: the moment the user set Never, the command they were told to use
    instead became a silent no-op.
    """
    msg_tok = _estimate_tokens(messages)
    sys_tok = _estimate_system_prompt_tokens(system_prompt) if system_prompt else 0
    # The fit rule and the cost rule are asked separately (see
    # core.context.compaction_plan): a 1,000,000-token window used to compact
    # at 120,000 because one `min()` answered both questions with the smaller.
    plan = core_context.compaction_plan(
        used_tokens=msg_tok,
        window_tokens=CONTEXT_WINDOW,
        reserve_tokens=sys_tok + TOOL_TOKENS + output_reserve(),
        # `force` asks the fit question at the default threshold rather than
        # not asking it at all, so /compact reports honest arithmetic while
        # still doing what it was told.
        fit_fraction=(core_context.DEFAULT_FIT_FRACTION if force
                      else COMPACTION_THRESHOLD),
        cost_limit=COMPACTION_COST_LIMIT,
    )
    if plan.reason == "overhead" and not force:
        # Not silence: this is the state where the agent cannot work properly
        # and the cause is fixable, so saying nothing would leave the user
        # watching every turn run slowly for no visible reason. Once per
        # session — it is a configuration problem, not a per-turn event.
        global _warned_overhead
        if not _warned_overhead:
            _warned_overhead = True
            print(f'  {YELLOW}⚠{RESET}  {plan.reserve:,} tokens of tools, prompt and '
                  f'output reserve leave no room in a {plan.window:,}-token window '
                  f'{DIM}(compaction cannot help — it only shrinks the '
                  f'conversation){RESET}\n'
                  f'     {DIM}Fix it in{RESET} {CYAN}/budget{RESET}{DIM} — or the '
                  f'Context Budget page in agent_cli.py{RESET}')
        return messages
    # `force` compacts a conversation that does not strictly need it — which
    # is the whole point of a manual command — but not one so short that the
    # summary plus its two framing messages would be larger than the original.
    # "Compacted" that grows the transcript is a worse answer than "nothing to
    # gain", and the caller reports the difference.
    if not plan.needed and not (force and len(messages) > 4):
        return messages
    before_tok = plan.used
    # The line itself is not optional whatever the diagnostics switch says:
    # this is a model call the user is sitting through, with no spinner over
    # it because it happens before the turn starts, and silence there is the
    # dead screen `Thinking` exists to prevent. Only the arithmetic behind it
    # is a diagnostic — "why now" is the question you ask when you are already
    # asking why the agent is slow.
    why = (f' ({plan.used:,} ≥ {plan.trigger:,} tokens, {plan.reason} limit)'
           if features().enabled("advanced_diagnostics") else '')
    print(f'  {DIM}[context] compacting conversation{why}...{RESET}')
    try:
        resp = _get_client().messages.create(
            model=_get_model(),
            max_tokens=4096,
            system=("Summarize this conversation. Keep all file paths, decisions made, "
                    "errors encountered, and current task state. Be specific about what "
                    "was changed and why."),
            messages=messages,
        )
        summary = "".join(b.text for b in resp.content if hasattr(b, "text"))
        print(f'  {GREEN}✓{RESET} {DIM}compacted to summary ({len(summary)} chars){RESET}')
        compacted = [
            {"role": "user", "content": f"[Conversation summary]\n{summary}"},
            {"role": "assistant", "content": "I have the context from our previous conversation. What should I work on next?"},
        ]
        compacted.extend(messages[-4:])
        _record_compaction("summary", before_tok, _estimate_tokens(compacted), plan)
        return compacted
    except Exception as e:
        print(f'  {RED}⚠{RESET} {DIM}compaction failed: {e}{RESET}')
        # Fallback: leave room for tools + max_tokens
        budget = int((CONTEXT_WINDOW - TOOL_TOKENS - output_reserve()) * 0.5)
        keep = [messages[-1]] if messages else []
        running = _estimate_tokens(keep)
        for m in reversed(messages[:-1]):
            if running + _estimate_tokens([m]) > budget:
                break
            keep.insert(0, m)
            running = _estimate_tokens(keep)
        print(f'  {YELLOW}⚠{RESET} {DIM}truncated to {len(keep)} messages ({running} est. tokens){RESET}')
        _record_compaction("truncate", before_tok, running, plan, error=str(e)[:200])
        return keep


def _prepare_turn_context(messages: list, user_input: str) -> tuple[list, str]:
    """Prune, build the system prompt, then compact exactly once against its
    real size.

    Order matters: compacting before the system prompt is known forces a
    second check with an understated reserve, and on a heavy system prompt
    (many facts/skills/MCP tools) close to the trigger boundary that second
    check can independently fire its own real summarization call — two
    network round-trips instead of at most one. `build_system_prompt` depends
    only on `user_input` and on-disk state, never on `messages` (see its
    docstring), so building it first costs nothing extra and removes the need
    for a second compaction check entirely.
    """
    messages = maybe_prune(messages)
    system_prompt = build_system_prompt(user_input)
    messages = maybe_compact(messages, system_prompt)
    return messages, system_prompt

# ---------------------------------------------------------------------------
# The agent loop
# ---------------------------------------------------------------------------

MAX_TOOL_CALLS_PER_TURN = int(os.environ.get("TOMAS_MAX_TOOL_CALLS", "40"))
# How many times bypass mode may extend the budget without asking. Nine
# extensions plus the initial budget is 400 tool calls at the default — more
# than any observed task (the session that motivated the mode used 56) while
# still being a number, so an unattended runaway ends instead of billing on.
MAX_AUTO_CONTINUATIONS = _env_int("TOMAS_MAX_AUTO_CONTINUATIONS", 9, minimum=1)
# Ceiling on a single tool result, kept in step with core.state's fail-safe.
# Raise it when a task genuinely needs one enormous result in context.
MAX_RESULT_CHARS = _env_int("TOMAS_MAX_RESULT_CHARS", 30_000, minimum=2_000)
_streaming_disabled = False  # set True if provider doesn't support streaming

# The loop itself now lives in core/loop.py. These aliases keep older call
# sites working while this shim is still in place.
_is_retryable_error = core_loop.is_retryable_error
_is_client_error = core_loop.is_client_error

# Session-scoped tool approvals. Answering "always" records the exact call the
# user saw; it no longer rewrites RISK_LEVELS for the rest of the process.
APPROVALS = ApprovalStore()


def init_mcp(config: Optional[dict] = None) -> dict:
    """Connect MCP servers and wire them into this module's tool state.

    Extracted from `main()`, where it was inlined. That mattered: the
    simulation harness needed the same startup and had to reimplement thirty
    lines of it, which is how a harness and its app drift until the harness is
    testing something the app does not do. One caller is a coincidence; two
    is a function.

    Returns a summary — servers connected, servers that failed, tool count,
    renames — so a caller can print whatever its front end prints instead of
    this deciding. Never raises: a broken MCP config degrades to the built-in
    tools, because losing `read_file` because a browser server has no
    credentials is not a trade anyone would choose.
    """
    global COMBINED_TOOLS, TOOL_TOKENS, MCP_TOOL_NAME_MAP, ALL_TOOLS, mcp_manager
    global _tool_pool_version

    from mcp_manager import read_mcp_servers

    summary = {"servers": [], "failed": {}, "disabled": [], "tools": 0,
               "renamed": 0, "dropped": 0, "budget": 0, "error": ""}
    try:
        manager = MCPManager()
        all_config = read_mcp_servers() if config is None else dict(config)
        summary["disabled"] = sorted(n for n, c in all_config.items()
                                     if c.get("disabled"))
        manager.discover_and_connect(config=all_config)
        mcp_manager = manager
        summary["servers"] = sorted(manager.servers)
        summary["failed"] = dict(getattr(manager, "failed_servers", {}) or {})

        if manager.tools:
            mcp_tools, MCP_TOOL_NAME_MAP, renamed = resolve_mcp_tool_conflicts(
                manager.tools)
            ALL_TOOLS = TOOLS + mcp_tools
            budget = tool_ceiling()
            COMBINED_TOOLS, dropped = apply_tool_cap(mcp_tools, max_allowed=budget)
            summary.update(tools=len(mcp_tools), renamed=renamed,
                           dropped=dropped, budget=budget)
        else:
            ALL_TOOLS = list(TOOLS)
            COMBINED_TOOLS = TOOLS
            MCP_TOOL_NAME_MAP = {}
        TOOL_TOKENS = estimate_tool_tokens(COMBINED_TOOLS)
    except Exception as e:
        summary["error"] = f"{type(e).__name__}: {e}"
        mcp_manager = MCPManager()
        ALL_TOOLS = list(TOOLS)
        COMBINED_TOOLS = TOOLS
        MCP_TOOL_NAME_MAP = {}
        TOOL_TOKENS = estimate_tool_tokens(TOOLS)
    # All three branches above converge here with ALL_TOOLS finalized — the
    # one place _avg_tool_tokens()'s cache needs invalidating.
    _tool_pool_version += 1
    return summary


def _tool_origin(name: str) -> str:
    """Human-readable provenance for a tool name."""
    if name in HANDLERS:
        return "built-in"
    if mcp_manager:
        srv = mcp_manager.get_server_for_tool(MCP_TOOL_NAME_MAP.get(name, name))
        if srv:
            return f"MCP: {srv}"
    return "built-in"


_EXIT_CODE_RE = re.compile(r'^\[exit (-?\d+) ')


def _record_tool_call(name: str, args: dict, preview: str,
                      duration_ms: int = 0, ok: bool = True) -> None:
    try:
        self_improve.record_tool_call(name, args, preview)
    except Exception:
        pass
    entry = {
        "turn": len(_turn_timings) + 1,
        "tool": name,
        "duration_sec": round(duration_ms / 1000, 3),
        "ok": ok,
    }
    # run_command now leads with its exit code, so the log can record the real
    # process status rather than "the tool returned a string".
    m = _EXIT_CODE_RE.match(preview or "")
    if m:
        entry["exit"] = int(m.group(1))
    elif not ok or (preview or "").startswith("Error:"):
        entry["exit"] = 1
        entry["error"] = (preview or "")[:200]
    else:
        entry["exit"] = 0
    if entry["exit"] != 0 and "error" not in entry:
        entry["error"] = (preview or "")[:200]
    # Arguments and results are already in `messages`; duplicating them here
    # is how a 6-turn session file reached 190 KB.
    _tool_log.append(entry)


def _last_user_text(messages: list) -> str:
    """The most recent user message as plain text, for tool selection."""
    for m in reversed(messages or []):
        if m.get("role") != "user":
            continue
        content = m.get("content")
        if isinstance(content, str):
            return content
        if isinstance(content, list):
            # Tool-result turns carry no user intent; keep looking back.
            texts = [b.get("text", "") for b in content
                     if isinstance(b, dict) and b.get("type") == "text"]
            if texts:
                return "\n".join(texts)
    return ""


#: User turns of context tool selection reads, newest first, and how many
#: times the newest is repeated to keep it outweighing them.
TOOL_CONTEXT_TURNS = 8
TOOL_CONTEXT_WEIGHT = 3
TOOL_CONTEXT_CHARS = 2000


def _recent_user_text(messages: list, turns: int = TOOL_CONTEXT_TURNS) -> str:
    """What the user has been asking about, not just their last keystroke.

    Selecting on the newest message alone loses the task the moment the user
    answers a question about it. Observed in session 20260804_144250: the
    request naming "docx" and "pdf" selected the whole word-docs working set,
    then "зроби новий", "2" and "зроби це, все вірно" each selected none of it
    — three words carry no keywords, so selection fell back to list order.
    By the time the user said "yes, do it", the tools to do it were gone, and
    the model concluded it would have to write a text file named .docx.

    The newest message is repeated because `extract_keywords` ranks by
    frequency: the current turn should still outweigh what came before it,
    or a conversation could never change subject.
    """
    collected: list[str] = []
    for m in reversed(messages or []):
        if m.get("role") != "user":
            continue
        content = m.get("content")
        if isinstance(content, str):
            text = content
        elif isinstance(content, list):
            # Tool-result turns carry no user intent; keep looking back.
            text = "\n".join(b.get("text", "") for b in content
                             if isinstance(b, dict) and b.get("type") == "text")
        else:
            text = ""
        if text.strip():
            collected.append(text)
        if len(collected) >= turns:
            break
    if not collected:
        return ""
    weighted = [collected[0]] * TOOL_CONTEXT_WEIGHT + collected
    return "\n".join(weighted)[:TOOL_CONTEXT_CHARS]


#: How many of the most recent messages are scanned for a failed tool result.
#: Small: the point is the failure the turn is *stuck on* right now, not every
#: error the session ever produced.
FAILURE_CONTEXT_MESSAGES = 4

#: Error text → words describing the capability that gets past it.
#:
#: Selection scores tools against the user's message, and a user asking for a
#: lab report says nothing about TLS. So when a call fails, the capability
#: needed to recover is exactly the one that cannot be scored — measured
#: across three sessions: 7 CERTIFICATE_VERIFY_FAILED against one host, after
#: which one session wrote its own urllib script with CERT_NONE and two
#: abandoned the source, while `stealthy_fetch` sat in the pool unselected
#: the whole time.
#:
#: These add *words to the query*, never tool names to the result. Relevance
#: still decides, so a mapping that guesses wrong costs a few keywords rather
#: than a wrong tool in the payload.
_FAILURE_HINTS: tuple[tuple[str, str], ...] = (
    ("CERTIFICATE_VERIFY_FAILED", "stealthy fetch browser page html certificate tls"),
    ("SSLError", "stealthy fetch browser page html certificate tls"),
    ("Download is starting", "pdf download extract text document"),
    ("HTTP 403", "stealthy fetch browser page bypass"),
    ("HTTP 429", "stealthy fetch browser page"),
    ("HTTP 999", "stealthy fetch browser page"),
    ("no such file", "list files directory search"),
    ("file not found", "list files directory search"),
    ("FileNotFoundError", "list files directory search"),
    ("ModuleNotFoundError", "install package documentation library"),
    ("is not recognized as an internal", "list files directory search"),
)


def _failure_context(messages: list) -> str:
    """Words describing what would get past the most recent tool failures.

    Returns "" when nothing recent failed, which is the common case and must
    stay free: adding keywords on a healthy turn would pull unrelated tools
    into a payload that is deliberately under-filled.
    """
    hints: list[str] = []
    for m in reversed((messages or [])[-FAILURE_CONTEXT_MESSAGES:]):
        content = m.get("content")
        if not isinstance(content, list):
            continue
        for block in content:
            if not isinstance(block, dict) or block.get("type") != "tool_result":
                continue
            body = block.get("content")
            if not isinstance(body, str):
                continue
            for needle, words in _FAILURE_HINTS:
                if needle.lower() in body.lower() and words not in hints:
                    hints.append(words)
    return " ".join(hints)


def tools_for_turn(messages: list) -> tuple[list[dict], list[dict]]:
    """Pick this turn's tool payload.

    Re-selected per turn rather than per session: a user who starts with file
    edits and moves to browser automation should get browser tools when they
    ask for them, not whichever server connected first at startup.
    """
    pool = enabled_tools()
    context = _recent_user_text(messages)
    try:
        context = f"{context}\n{self_improve.get_session_analysis().get('purpose', '')}"
    except Exception:
        pass
    # Appended, not substituted: the turn is still about what the user asked
    # for, and the recovery capability is an addition to that, not a
    # replacement for it.
    failure = _failure_context(messages)
    if failure:
        context = f"{context}\n{failure}"
    # A skill that declares the tools it needs narrows the payload to them.
    # Keyed off the user's own message, the same string `build_triggered_skills`
    # matches on, so the allowlist and the skill body arrive together or not
    # at all.
    allowlist = set()
    try:
        allowlist = triggered_tool_allowlist(_recent_user_text(messages))
    except Exception:
        pass
    return select_tools(pool, context, tool_ceiling(), allowlist=allowlist or None)


TEXT_TOOL_PROTOCOL = """

## Tool use (text protocol)

This endpoint does not support native tool calling, so tools are invoked by
writing a fenced block. To call one, emit exactly:

```tool_call
{"name": "<tool name>", "input": {...}}
```

Emit one block per call and then stop; the result is returned to you in the
next message. Available tools:
"""

_TEXT_TOOL_RE = re.compile(r"```tool_call\s*\n(.*?)\n?```", re.S)


def describe_tools_as_text(tools: list[dict]) -> str:
    """Render the tool list into the system prompt for the text protocol."""
    lines = [TEXT_TOOL_PROTOCOL]
    for t in tools:
        schema = t.get("input_schema", {}) or {}
        params = ", ".join(schema.get("properties", {}).keys()) or "none"
        lines.append(f"- `{t['name']}`({params}) — {t.get('description', '')[:160]}")
    return "\n".join(lines)


def parse_text_tool_calls(text: str) -> list[dict]:
    """Extract tool calls a model wrote as text. Malformed blocks are skipped."""
    calls = []
    for block in _TEXT_TOOL_RE.findall(text or ""):
        try:
            payload = json.loads(block.strip())
        except json.JSONDecodeError:
            continue
        if isinstance(payload, dict) and payload.get("name"):
            calls.append({"name": payload["name"],
                          "input": payload.get("input") or {}})
    return calls


def cache_marked_system(system_prompt: str):
    """Mark the stable prefix so the provider caches it instead of re-reading it.

    Every turn resends ~6,300 identical tokens — the system prompt plus the
    tool definitions — and re-tokenising and re-attending them is most of what
    a short turn costs. One `cache_control` breakpoint at the end of the system
    block covers the whole prefix, because Anthropic's cache hierarchy is
    tools → system → messages: marking system caches the tools ahead of it too.

    Returns the block form when caching is on, the plain string otherwise, so
    the call site stays a single assignment. OpenAI-wire providers get the
    string: they do prefix caching server-side and have no `cache_control`,
    and `anthropic_to_openai` flattens block lists anyway.
    """
    if not system_prompt:
        return system_prompt
    return [{"type": "text", "text": system_prompt,
             "cache_control": {"type": "ephemeral"}}]


#: How many `cache_control` breakpoints Anthropic accepts in one request.
#: One is spent on the system block; the rest are available to the history.
MAX_CACHE_BREAKPOINTS = 4

#: Don't spend a breakpoint on a trivial amount of history — each one has a
#: write cost, and caching 200 tokens never repays it.
MIN_CACHED_HISTORY_CHARS = 4_000


def _measure(message: dict) -> int:
    """Characters the model actually sees in one message.

    Deliberately not `json.dumps`: that escapes Cyrillic to `\\uXXXX`, six
    characters for one, and inflated a real Ukrainian transcript four-fold —
    which would make every size threshold in this module fire at a quarter of
    its intended size on exactly the sessions this project is used for.
    """
    content = message.get("content")
    if isinstance(content, str):
        return len(content)
    total = 0
    for block in content or []:
        if not isinstance(block, dict):
            total += len(str(block))
            continue
        kind = block.get("type")
        if kind == "text":
            total += len(block.get("text") or "")
        elif kind == "tool_result":
            total += len(str(block.get("content") or ""))
        elif kind == "tool_use":
            try:
                total += len(json.dumps(block.get("input") or {},
                                        ensure_ascii=False))
            except Exception:
                total += len(str(block.get("input") or ""))
        else:
            total += len(str(block))
    return total


def clear_history_cache_marks(messages: list) -> None:
    """Remove every `cache_control` this module put on the history."""
    for message in messages:
        content = message.get("content")
        if not isinstance(content, list):
            continue
        for block in content:
            if isinstance(block, dict):
                block.pop("cache_control", None)


def mark_history_for_caching(messages: list, breakpoints: int = 3) -> int:
    """Put `cache_control` on the history so it is re-read, not re-processed.

    Only the system block was ever marked. But the system prompt is ~7,000
    tokens while the history reached 48,655 in a real session — so the part
    that dominated the bill was the part paying full price on every turn,
    because Anthropic caches a *prefix* and the prefix ended where the marks
    did.

    Mutates in place and returns how many marks were placed. In place because
    `core.loop` appends to this very list: handing back a copy would send the
    marked version to the model while the assistant's replies accumulated in
    an object nobody reads — the no-memory bug, reintroduced silently.

    Moving a breakpoint between turns is free: `cache_control` is a directive
    about where to cut, not content, so it does not change the prefix the cache
    is keyed on. That is what makes re-marking every turn safe.

    Only messages whose content is *already* a block list are marked. A string
    would have to be reshaped into blocks to carry the key, and that rewrites
    the user's own turns in the saved transcript for a few tokens of benefit.
    Tool-heavy sessions — the ones with history worth caching — are mostly
    block-list messages anyway.
    """
    clear_history_cache_marks(messages)
    if len(messages) < 2 or breakpoints < 1:
        return 0

    # Candidates are any block-list message except the last. Turn boundaries
    # are deliberately *not* used: a turn that makes 56 tool calls offers one
    # boundary and one breakpoint, leaving the bulk of its own history uncached
    # — which is the case that needs caching most.
    running, sized = 0, []
    for i, message in enumerate(messages):
        running += _measure(message)
        if (i < len(messages) - 1
                and isinstance(message.get("content"), list)
                and running >= MIN_CACHED_HISTORY_CHARS):
            sized.append(i)
    if not sized:
        return 0

    # The latest candidate matters most — it caches the longest prefix. The
    # others are spread evenly behind it so that when the tail shifts, there is
    # still a usable cached prefix further back instead of an all-or-nothing
    # miss.
    chosen = {sized[-1]}
    for n in range(1, breakpoints):
        chosen.add(sized[max(0, len(sized) * (breakpoints - n) // (breakpoints + 1))])

    placed = 0
    for i in sorted(chosen)[-breakpoints:]:
        tail = next((b for b in reversed(messages[i]["content"])
                     if isinstance(b, dict)), None)
        if tail is not None:
            tail["cache_control"] = {"type": "ephemeral"}
            placed += 1
    return placed


def build_state(system_prompt: str, messages: list, responder) -> AgentState:
    """Assemble the turn context the core needs out of this module's globals."""
    global _CURRENT_INTERRUPT
    # handle_run_command has no access to AgentState (execute_tool is a bare
    # (name, params) -> str callable) so it reads this module global directly
    # to kill its subprocess the moment Esc is pressed, instead of only at
    # the next loop checkpoint — which for a shell command can be as far away
    # as the whole timeout.
    _CURRENT_INTERRUPT = getattr(responder, "esc_interrupt", None)

    active_features = features()
    selected, withheld = tools_for_turn(messages)
    # Compacted here rather than at discovery so `ALL_TOOLS` keeps the full
    # schemas: selection scores against descriptions, and scoring clipped text
    # would quietly change which tools a message retrieves.
    selected = compact_tool_schemas(selected)
    if withheld:
        system_prompt = system_prompt + withheld_tools_notice(withheld)
    caps = _active_capabilities()

    # ── Degradations. Each costs a feature; none costs the session. ──
    if not caps.tool_use:
        # No native tool calling: describe the tools and parse a text
        # protocol out of the reply instead (see agent_loop).
        system_prompt = system_prompt + describe_tools_as_text(selected)
    if not caps.system_prompt:
        # No system role: prepend it as the first user message.
        messages = ([{"role": "user", "content": system_prompt},
                     {"role": "assistant", "content": "Understood."}]
                    + list(messages))
        system_prompt = ""

    # Last, so it wraps the finished text: the withheld-tools notice and the
    # text-protocol description are both appended above, and a breakpoint
    # placed before them would cache a prefix that no longer matches.
    if caps.prompt_caching and caps.system_prompt:
        system_prompt = cache_marked_system(system_prompt)
    if caps.prompt_caching:
        # The system block spends one of the four breakpoints; the history
        # gets the rest. In place — see the note in mark_history_for_caching
        # about why a copy would lose the transcript.
        mark_history_for_caching(messages,
                                 breakpoints=MAX_CACHE_BREAKPOINTS - 1)
        # ...and again before every model call the turn goes on to make. Marking
        # once here covered only the history that existed before the first call;
        # a turn that then makes 29 tool calls appends 58 messages behind those
        # marks and pays full price for all of them, which is the case
        # mark_history_for_caching's own docstring names as needing it most.
        # Re-cutting is free: cache_control says where to cut, so moving it does
        # not change the bytes the cache is keyed on.
        remark_cache = lambda msgs: mark_history_for_caching(
            msgs, breakpoints=MAX_CACHE_BREAKPOINTS - 1)
    else:
        remark_cache = None

    # ── The every-3rd-reply cap ──
    # Applied here rather than inside `run_turn` for the same reason the mode
    # flags are: the core is handed numbers and a flag, and decides with them.
    # `reply_capped` is what stops the truncation escalation undoing this on
    # the very turn it applies to — see `core.loop._can_escalate`.
    capped = core_features.caps_this_reply(active_features,
                                           _replies_so_far(messages))
    max_output = (core_features.SHORT_REPLY_MAX_TOKENS if capped
                  else effective_max_tokens(caps))

    return AgentState(
        system_prompt=system_prompt,
        messages=messages,
        get_client=_get_client,
        get_model=_get_model,
        tools=selected if caps.tool_use else [],
        max_tokens=max_output,
        reply_capped=capped,
        temperature=effective_temperature(),
        max_turn_seconds=MAX_TURN_SECONDS,
        execute_tool=execute_tool,
        risk_of=risk_for,
        origin_of=_tool_origin,
        parallel_safe=parallel_safe,
        before_model_call=remark_cache,
        describe_endpoint=_describe_endpoint,
        responder=responder,
        approvals=APPROVALS,
        auto_approve_low=AUTO_APPROVE_LOW,
        yolo=YOLO_MODE,
        auto_continue=BYPASS_MODE,
        max_auto_continuations=MAX_AUTO_CONTINUATIONS,
        tool_budget=MAX_TOOL_CALLS_PER_TURN,
        max_result_chars=MAX_RESULT_CHARS,
        # Three separate answers to "may this turn stream?", and all three have
        # to say yes: the user's switch, the provider's probed capability, and
        # the runtime fallback that turns it off after a provider rejects a
        # stream mid-session. The switch is listed first because it is the only
        # one the user can see.
        streaming_enabled=(active_features.enabled("streaming")
                           and (not _streaming_disabled) and caps.streaming),
        on_tool_call=_record_tool_call,
        interrupted=getattr(responder, "is_interrupted", lambda: False),
    )


TEXT_PROTOCOL_MAX_ROUNDS = 6


def _run_text_protocol(state, adapter, reply: str, messages: list) -> str:
    """Drive tool calls for providers with no native tool support.

    The model writes ```tool_call blocks; we execute them, feed the results
    back, and let it continue. Slower and less reliable than native tool use
    — which is the point: the missing capability costs a feature, not the
    session.

    Permission is asked through `state` and the live `adapter`, the same way
    `core.loop.run_turn` asks it. This called a `request_permission()` that
    does not exist anywhere in the program, so the whole path died with
    `NameError` the moment a model actually emitted a tool call — on the one
    route that exists *because* the provider cannot call tools natively, i.e.
    exactly the small and self-hosted models it was written for. Nothing
    caught it because nothing tested it.

    Not `check_permission()`: that legacy helper reads module globals and
    constructs its own `TerminalAdapter`, so a headless run would have been
    handed an interactive prompt with no one to answer it, and the turn's own
    mode (yolo/bypass) would have been ignored.
    """
    from core.events import PermissionNeeded

    for _ in range(TEXT_PROTOCOL_MAX_ROUNDS):
        calls = parse_text_tool_calls(reply)
        if not calls:
            return reply
        results = []
        for call in calls:
            name, params = call["name"], call["input"]
            approved = True
            if state.needs_permission(name, params):
                decision = adapter.ask(PermissionNeeded(
                    "", name, params, state.risk_of(name, params)))
                if decision == "always_allow_this_call":
                    state.approvals.approve(name, params)
                else:
                    approved = decision == "allow"
            if not approved:
                results.append(
                    f"[{name}] Error: the user denied this tool call. Retrying "
                    f"the same call will be denied again — do not re-issue it.")
                continue
            t0 = time.perf_counter()
            try:
                # Through `state`, like everything else here. `build_state`
                # injects the module-level `execute_tool`/`_record_tool_call`,
                # so this is the same call in production — but reading the
                # globals directly meant the one path a host cannot override
                # was the fallback path, and the function was untestable
                # without touching the real filesystem.
                out = state.execute_tool(name, params)
                ok = not (isinstance(out, str)
                          and out.lstrip().startswith("Error:"))
            except Exception as e:
                out, ok = f"Error: tool raised {type(e).__name__}: {e}", False
            if state.on_tool_call:
                try:
                    state.on_tool_call(name, params, str(out)[:200],
                                       int((time.perf_counter() - t0) * 1000),
                                       ok)
                except Exception:
                    pass
            results.append(f"[{name}]\n{str(out)[:state.max_result_chars]}")
        messages.append({"role": "assistant", "content": reply})
        messages.append({"role": "user",
                         "content": "Tool results:\n\n" + "\n\n".join(results)})
        state.messages = messages
        reply = adapter.run(state)
    return reply


# ── Deterministic rule capture ─────────────────────────────────────────
# A live run exposed the other half of the problem. Told "Rule one: always end
# every reply with the date. Save that to memory", the model replied "Saved."
# and never called save_memory — three times in a row, in a session whose tool
# log shows ten calls and not one of them a write. BASE_PROMPT already forbids
# claiming an action was taken without calling the tool; a weak model ignored
# it anyway.
#
# Retrieval was only ever half the loop. If capture depends on the model
# choosing to call a tool, then on a model that does not, the user's rule is
# lost while being told it was kept — which is worse than refusing outright.
# So when the user states a rule *and* asks for it to be saved, save it here,
# deterministically, before the model gets a chance not to.
#
# Both signals are required. Either alone is too loose: "always run the tests
# first" is conversation, and "save this file" is not a rule.
_SAVE_REQUEST = re.compile(
    r"\b(?:save|remember|store|keep)\b.{0,30}?"
    r"\b(?:that|this|it|them|memory|rules?|preferences?)\b"
    # Ukrainian writes this with U+02BC (ʼ) or U+2019 (’) far more often than
    # with an ASCII quote, and matching only ' silently drops the whole
    # language — the one this agent defaults to.
    r"|запам['’ʼʹ]?ятай|\b(?:збережи|запомни|сохрани)\b"
    r"|\bfrom now on\b|\bвідтепер\b|\bотныне\b",
    re.IGNORECASE,
)


def capture_stated_rule(user_text: str) -> Optional[str]:
    """Persist a rule the user stated and asked to have saved.

    Returns the text captured, or None. Never raises — nothing in the learning
    path may break the user's turn.
    """
    text = (user_text or "").strip()
    if not text or len(text) > 600:
        return None
    if not _SAVE_REQUEST.search(text):
        return None
    try:
        if not learning.is_enabled() or not learning.looks_like_directive(text):
            return None
        learning.remember(learning.KIND_DIRECTIVE, text,
                          evidence="user stated it and asked to save it",
                          scope="global")
        return text
    except Exception:
        return None


# ── Standing-rule reinforcement ────────────────────────────────────────
# Compliance in the transcripts is bimodal *per session*, not per turn: whatever
# the model did on turn 1 it kept doing for the next 29, whether or not that
# matched the rule. That is transcript momentum, and a system prompt alone does
# not beat it — the conversation is simply the louder signal by turn 20. So the
# rules are also restated in-context, periodically and briefly.
STANDING_RULE_REMINDER_EVERY = 10
_REMINDER_OPEN = "<standing-rules>"
_REMINDER_CLOSE = "</standing-rules>"


def _reinforce_standing_rules(messages: list) -> None:
    """Every Nth user turn, restate the standing rules inside the conversation.

    Appended to the user's own message rather than sent as a separate turn, so
    it cannot be mistaken for something the user said on its own and cannot
    desynchronise the user/assistant alternation some providers require.
    """
    if not messages or messages[-1].get("role") != "user":
        return
    user_turns = sum(1 for m in messages
                     if m.get("role") == "user" and isinstance(m.get("content"), str))
    if user_turns == 0 or user_turns % STANDING_RULE_REMINDER_EVERY != 0:
        return
    try:
        rules = learning.directives_for_prompt()
    except Exception:
        return
    if not rules:
        return
    content = messages[-1].get("content")
    if not isinstance(content, str) or _REMINDER_OPEN in content:
        return
    messages[-1]["content"] = (
        f"{content}\n\n{_REMINDER_OPEN}\n"
        f"Still in force — these applied to every reply so far and apply to "
        f"this one:\n{rules}\n{_REMINDER_CLOSE}"
    )


def _text_of_message(message: dict) -> str:
    """The plain text of a message, whatever block shape it arrived in."""
    content = message.get("content")
    if isinstance(content, str):
        return content
    parts = []
    for block in content or []:
        if isinstance(block, dict) and block.get("type") == "text":
            parts.append(block.get("text", ""))
    return "\n".join(parts)


def _intercept_slash_command(messages: list) -> Optional[str]:
    """Run a `/command` that arrived as a *message* rather than as REPL input.

    `handle_slash_command` used to be reachable only from the REPL's own input
    path, so a harness, a subagent, or a pasted script that sent "/status" as a
    message got it forwarded to the model as prose. The model then either
    invented the output or tried to shell out to reimplement the command — both
    observed, repeatedly, in the session transcripts.

    Unknown commands deliberately fall through to the model: a message that
    merely begins with a path must not be swallowed. `//` escapes a literal
    leading slash.

    The name is matched against SLASH_COMMANDS *before* dispatching, rather
    than relying on the handler's return value — the handler answers an
    unrecognised command with the help screen, not with None, so dispatching
    first would turn "/usr/local/bin is where it lives" into a help dump.
    """
    if not messages or messages[-1].get("role") != "user":
        return None
    text = _text_of_message(messages[-1]).strip()
    if not text.startswith("/") or text.startswith("//"):
        return None
    # Only a single-line command — a code block that happens to open with a
    # slash is not an instruction to the REPL.
    if "\n" in text:
        return None
    name = text[1:].split(maxsplit=1)[0].lower() if len(text) > 1 else ""
    if name not in SLASH_COMMANDS:
        return None
    try:
        result = handle_slash_command(text[1:], messages)
    except Exception:
        return None
    if result is None:
        return None
    # The command consumed the turn; drop the user message so the transcript
    # does not carry a prompt that never reached the model.
    messages.pop()
    return result


def _try_zen_fallback(reason: str) -> str:
    """Switch to a different free, served Zen model after this one turned out
    to be listed but not actually served, so a task can continue instead of
    just stopping.

    Zen only: Ollama has no equivalent "try a different served model"
    catalogue to fall back to, and every other provider is outside the
    current working set (see provider_manager.VISIBLE_PROVIDER_TYPES).
    `mark_unavailable` then `default_free_model` compose correctly with no
    extra plumbing — `default_free_model` filters through `_stamp_served`,
    which re-reads the unavailable-models file fresh on every call, so the
    just-failed model is already excluded from the candidate it returns.

    Returns the new model name, or "" when no fallback applies (a different
    provider is active, nothing else free is currently served, or switching
    itself failed) — the caller's existing "model unavailable" message stands
    unchanged in that case. A failure here must not cost the user *more* than
    the original error already did, so nothing above this returning "" is
    allowed to propagate — a config write that fails partway through must
    leave the turn ending the way it always has, not crash it worse.
    """
    try:
        import provider_manager
        active = provider_manager.get_active()
        if active is None or active.type != "zen":
            return ""
        import zen_catalog
        failed_model = _get_model()
        zen_catalog.mark_unavailable(failed_model, reason)
        candidate = zen_catalog.default_free_model()
        if not candidate or candidate == failed_model:
            return ""
        active.model = candidate
        provider_manager.save(active, activate_it=True)
        provider_manager.activate(active.name)
        reinit_client()
        _refresh_context_window()
    except Exception:
        return ""
    return candidate


def agent_loop(system_prompt: str, messages: list) -> str:
    """Shim — drives core.loop.run_turn through the terminal adapter.

    Kept so the REPL and agent_cli.py keep working unchanged. New front ends
    should drive run_turn directly with their own adapter instead.
    """
    global _streaming_disabled

    handled = _intercept_slash_command(messages)
    if handled is not None:
        return handled

    # Capture before the model call, so the rule is stored whether or not the
    # model bothers to call save_memory — and so it is already in the prompt
    # this same turn rather than only from the next one.
    if messages and messages[-1].get("role") == "user":
        captured = capture_stated_rule(_text_of_message(messages[-1]))
        if captured:
            system_prompt = build_system_prompt(_text_of_message(messages[-1]))

    _reinforce_standing_rules(messages)

    interactive = bool(getattr(sys.stdin, "isatty", lambda: False)())
    adapter = TerminalAdapter(interactive=interactive,
                              show_status=features().enabled("status_indicator"),
                              diagnostics=features().enabled("advanced_diagnostics"))
    state = build_state(system_prompt, messages, adapter)

    # A local model is about to be loaded into VRAM and will stay there for
    # five minutes after this session stops talking to it. Claimed before the
    # call rather than after, so a session killed mid-turn still leaves a
    # record another session can see and sweep — see ollama_runtime.py.
    _claim_local_model()

    turn_index = len(_turn_timings) + 1
    t0 = time.perf_counter()
    error: Optional[BaseException] = None
    reply = ""
    # Snapshot rather than a counter: _tool_log is appended by the same
    # callback every front end uses, so the difference is the tool count for
    # this turn without threading anything new through the loop.
    tools_before = len(_tool_log)
    try:
        # messages already carries the user's turn, so no user_message here.
        reply = adapter.run(state)
        if not _active_capabilities().tool_use:
            reply = _run_text_protocol(state, adapter, reply, messages)
        # A free Zen model that is listed but not actually served used to just
        # end the task here — the red error above already printed, and the
        # user had to notice, run /model, and hope the next pick was not also
        # dead. `messages` is untouched by the failure (the core never appends
        # anything before the request itself fails), so whatever the turn had
        # already gathered — tool results included — is still there to retry
        # with. One bounded attempt, never a chain: this is a fallback, not a
        # search.
        # A model that exists but refuses tool definitions. The agent already
        # has a protocol for that — it just had never been reached this way,
        # because the refusal arrived as an unclassified 4xx and ended the
        # turn. Degrading first means the retry is on the text protocol, and
        # `degrade_capability` persists it, so the next session starts already
        # knowing. Measured on Groq: five of its ten chat models answer
        # `400 "tool calling" is not supported with this model`.
        if not (reply or "").strip() and getattr(state, "tool_use_rejected", False):
            degrade_capability("tool_use", "rejected the tool definitions")
            print(f'  {YELLOW}⟳{RESET} {DIM}retrying with the text tool '
                  f'protocol …{RESET}')
            state = build_state(system_prompt, messages, adapter)
            reply = adapter.run(state)
            reply = _run_text_protocol(state, adapter, reply, messages)

        if not (reply or "").strip() and getattr(state, "model_unavailable", False):
            fallback_model = _try_zen_fallback(state.last_error or "")
            if fallback_model:
                print(f'  {YELLOW}⟳{RESET} {DIM}unavailable — retrying with'
                     f'{RESET} {CYAN}{fallback_model}{RESET}{DIM} …{RESET}')
                state = build_state(system_prompt, messages, adapter)
                reply = adapter.run(state)
                if not _active_capabilities().tool_use:
                    reply = _run_text_protocol(state, adapter, reply, messages)
    except BaseException as e:
        error = e
        raise
    finally:
        _turn_timings.append(time.perf_counter() - t0)
        # A model the provider will not serve is worth writing down. The core
        # only reports it — it cannot import a catalogue without giving up
        # being dependency-free — so the host does the writing, here, where it
        # already reads the rest of the turn's findings back.
        if getattr(state, "model_unavailable", False):
            try:
                import zen_catalog

                zen_catalog.mark_unavailable(_get_model(),
                                             state.last_error or "")
            except Exception:
                pass
        # Propagate what the turn learned back into module state.
        was_streaming = not _streaming_disabled
        _streaming_disabled = not state.streaming_enabled
        if was_streaming and _streaming_disabled:
            # The core already fell back. Persist that only when the failure
            # says something about the provider: a 429 or a 5xx means the
            # endpoint was busy, not that it cannot stream, and writing it
            # down would disable streaming for good over a transient blip.
            if state.streaming_error_retryable:
                _streaming_disabled = False   # retry streaming next turn
            else:
                degrade_capability("streaming", state.streaming_error or "stream failed")
        _last_turn_usage["input"] = state.usage.get("input", 0)
        _last_turn_usage["output"] = state.usage.get("output", 0)
        _last_turn_usage["cached_input"] = state.usage.get("cached_input", 0)
        _last_turn_diag["stop_reason"] = state.last_stop_reason or ""
        _last_turn_diag["error"] = state.last_error or ""
        _last_turn_diag["tool_calls"] = len(_tool_log) - tools_before
        # A call has now gone through the shim, which is the only moment
        # Ollama's real served window can be read rather than estimated.
        _verify_ollama_window()
        _session_tokens["input"] += state.usage.get("total_input", 0)
        _session_tokens["output"] += state.usage.get("total_output", 0)
        _session_tokens["cached_input"] += state.usage.get("total_cached_input", 0)
        _session_tokens["calls"] += state.usage.get("calls", 0)
        # What the streamed call cost before being thrown away and re-issued
        # non-streamed. Spent on every tool step and counted by nothing until
        # now — see core.loop._record_discarded_stream.
        for key in ("duplicate_input", "duplicate_calls",
                    "would_have_served", "stream_malformed_tool_args"):
            _session_tokens[key] = (_session_tokens.get(key, 0)
                                    + state.usage.get(key, 0))
        # Bucketed by why the stream was discarded (see
        # core.loop._record_discarded_stream) — dynamic keys, so copied by
        # prefix rather than added to the fixed tuple above.
        for key, value in state.usage.items():
            if key.startswith("duplicate_reason_"):
                _session_tokens[key] = _session_tokens.get(key, 0) + value
        # A turn that produced nothing is recorded as such. Silence here is
        # what let a session with eight prompts and zero replies be saved,
        # and then be reported as eight turns of completed work.
        if error is not None or not (reply or "").strip():
            # Prefer the reason the core recorded — "empty_reply" with no
            # explanation is the unreadable record P6-11 exists to prevent.
            core_error = getattr(state, "last_error", None)
            if error is not None:
                reason, detail = type(error).__name__, str(error)
            elif core_error:
                reason, detail = "turn_error", core_error
            else:
                reason, detail = "empty_reply", ""
            _failed_turns.append({
                "turn": turn_index,
                "reason": reason,
                "error": detail[:300],
            })
        elif (len(_tool_log) == tools_before
              and len((reply or "").strip()) < LOW_CONTENT_REPLY_CHARS):
            # Finished cleanly, called nothing, said almost nothing.
            _low_content_turns.append({
                "turn": turn_index,
                "reply_chars": len((reply or "").strip()),
                "reply": (reply or "").strip()[:80],
            })
    return reply


# ---------------------------------------------------------------------------
# Slash commands
# ---------------------------------------------------------------------------

SLASH_COMMANDS = {
    "help":         {"desc": "Show this help message",            "icon": "ℹ"},
    "clear":        {"desc": "Clear conversation history",        "icon": "✧"},
    "status":       {"desc": "Show current model and connection", "icon": "◈"},
    "version":      {"desc": "Show TOMAS's version and last-updated date", "icon": "ℹ"},
    "model":        {"desc": "Show/switch model: /model [list|<name>]", "icon": "◎"},
    "mode":         {"desc": "Show/change mode: /mode [auto|default]", "icon": "⚙"},
    "config":       {"desc": "Interactive menu: provider, model, mode", "icon": "🛠"},
    "compact":      {"desc": "Force compact conversation now",    "icon": "⚙"},
    "budget":       {"desc": "Context budget: /budget [economy|balanced|full|auto|tools N|output N|on|off <section>]", "icon": "▣"},
    "settings":     {"desc": "Feature switches: /settings [<name>] to toggle", "icon": "⚙"},
    "debug":        {"desc": "Raw request/response JSON: /debug [on|off|N|schemas]", "icon": "🐞"},
    "export":       {"desc": "Save the conversation: /export [txt|json]", "icon": "📤"},
    "setup":        {"desc": "Tell TOMAS about you; tunes instructions and defaults", "icon": "🧭"},
    "skills":       {"desc": "List installed skills",            "icon": "⚡"},
    "skill":        {"desc": "Run a skill: /skill <name>",        "icon": "⚡"},
    "mcp-prompt":   {"desc": "MCP prompt templates: /mcp-prompt [name]", "icon": "◈"},
    "mcp-resources": {"desc": "List resources published by MCP servers", "icon": "◈"},
    "provider":     {"desc": "Show/switch provider: /provider [list|<name>|probe]", "icon": "◎"},
    "pdf-report":   {"desc": "Generate AI news PDF report",      "icon": "📄"},
    "zen":          {"desc": "OpenCode Zen proxy status",         "icon": "◉"},
    "self-improve": {"desc": "What the agent has learned",        "icon": "🧠"},
    "si":           {"desc": "Alias for /self-improve",           "icon": "🧠"},
    "forget":       {"desc": "Forget a learned fact: /forget <id>", "icon": "🗑"},
    "private":      {"desc": "Toggle incognito (learn nothing)",  "icon": "🕶"},
    "save":         {"desc": "Save current session",              "icon": "💾"},
    "load":         {"desc": "Load a saved session: /load <id>",  "icon": "📂"},
    "session":      {"desc": "Session mgmt: list/save/continue",  "icon": "📋"},
    "sessions":     {"desc": "Alias for /session list",            "icon": "📋"},
    "rules":        {"desc": "Rules: list, add, edit, rm (or type #<text>)", "icon": "📌"},
    "note":         {"desc": "Create a self-note: /note <title> <content>", "icon": "📝"},
    "notes":        {"desc": "List all self-notes",               "icon": "📒"},
    "exit":         {"desc": "Exit TOMAS",                        "icon": "✕"},
}

def _describe_endpoint() -> str:
    """Where requests are going, in words — for error messages only.

    `core.loop` reports an unreachable endpoint and cannot name it: it must
    not know that `provider_manager` exists. This is the host's answer, and it
    is what turns "[WinError 10061] ... actively refused it" into a sentence
    naming the provider and address that refused.
    """
    try:
        import provider_manager
        active = provider_manager.get_active()
        if active is None:
            return os.environ.get("ANTHROPIC_BASE_URL", "")
        where = active.base_url or os.environ.get("ANTHROPIC_BASE_URL", "")
        if active.type == "zen" and not where:
            where = "opencode.ai/zen"     # reached in-process, not via base_url
        return f"{active.name} at {where}" if where else active.name
    except Exception:
        return ""


def _get_model() -> str:
    """Read model from environment or active provider configuration."""
    model = os.environ.get("AGENT_MODEL")
    if model:
        return model
    try:
        import provider_manager
        active = provider_manager.get_active()
        if active and active.model:
            return active.model
    except Exception:
        pass
    return "Not set"



# The keys the prompt actually binds. Kept next to the help renderer so a new
# binding cannot be added without the place users look to find it.
KEY_HELP = [
    ("Enter",        "send · accepts the highlighted /command"),
    ("Tab",          "complete a /command · otherwise cycle mode"),
    ("↑ ↓",          "history · moves the /command selection"),
    ("Esc",          "clear the line · during a reply or tool call, stops it"),
    ("Ctrl+W",       "delete the last word"),
    ("Ctrl+U",       "clear the line"),
    ("Ctrl+Z",       "undo the last clear, paste or history recall"),
    ("Ctrl+Y",       "copy the line to the clipboard · ↑ first to copy an older one"),
    ("Ctrl+L",       "clear the screen, keep what is typed"),
    ("Ctrl+C",       "cancel"),
    ("⇧+Space",      "toggle auto-approve"),
    ("Ctrl+Alt+X",   "debug view — the last raw request and response"),
    ("F5 F6 F7 F8 F9", "auto · default · strict · yolo · bypass"),
]


def _show_commands(match: str = "") -> str:
    """Build a formatted help string for matching slash commands."""
    from text_display import rule

    lines = []
    lines.append(f'  {BOLD}Available commands{RESET}')
    lines.append(f'  {DIM}{rule()}{RESET}')
    for cmd, info in sorted(SLASH_COMMANDS.items()):
        if match and not cmd.startswith(match.lower()):
            continue
        icon = info["icon"]
        desc = info["desc"]
        padded = " " * max(1, 12 - len(cmd))
        lines.append(f'    {DIM}{icon}{RESET}  {CYAN}/{cmd}{RESET}{padded}{DIM}{desc}{RESET}')

    # Only when showing everything: filtered output is an answer to a specific
    # question, and a key table underneath it is just noise.
    if not match:
        lines.append('')
        lines.append(f'  {BOLD}Keys{RESET}')
        lines.append(f'  {DIM}{rule()}{RESET}')
        width = max(len(k) for k, _ in KEY_HELP)
        for key, desc in KEY_HELP:
            lines.append(f'    {CYAN}{key}{RESET}{" " * (width - len(key) + 2)}{DIM}{desc}{RESET}')

    lines.append('')
    lines.append(f'  {DIM}Type{RESET} {CYAN}/command{RESET} {DIM}to run — or just{RESET} {CYAN}/{RESET} {DIM}to see all{RESET}')
    return '\n'.join(lines)


def _numbered_menu_fallback(title: str, items: list[str], footer: str = "") -> int:
    """Digit-driven stand-in for `_arrow_menu` when msvcrt is unavailable."""
    print(f'  {BOLD}{title}{RESET}')
    for i, label in enumerate(items, 1):
        print(f'    {CYAN}{i}{RESET}  {label}')
    try:
        sel = input(f'  {DIM}{footer or "Number to choose, Enter to cancel"}:{RESET} ').strip()
    except (EOFError, KeyboardInterrupt):
        print()
        return -1
    if not sel or not sel.isdigit():
        return -1
    n = int(sel)
    return n - 1 if 1 <= n <= len(items) else -1


def _arrow_menu(title: str, items: list[str], footer: str = "",
                erase_on_exit: bool = False) -> int:
    """A minimal arrow-key picker for slash-command menus (e.g. /config).

    Deliberately not `agent_cli.arrow_menu` — that lives in the TUI
    entrypoint, and importing it here would pull `agent_cli` into the
    headless REPL (`agent.py` already reads raw keystrokes itself for the
    main prompt — see `read_input_with_suggestions` — so this reuses the
    exact same two-byte codes msvcrt reports for arrow keys, '\\xe0'/'\\x00'
    prefix then 'H'/'P', rather than a second decoding of the same keys).
    Returns the selected index, or -1 on Esc/cancel.
    """
    if not items:
        return -1
    try:
        import msvcrt
    except ImportError:
        return _numbered_menu_fallback(title, items, footer)

    # `\n` is a *line feed*. With virtual-terminal processing enabled — which
    # is how this prompt renders colour at all — it moves down one row and
    # leaves the cursor in the column it was already in. `\033[2K` erases the
    # row but does not move the cursor either. So every item after the first
    # was drawn starting at the column the previous item ended on, and every
    # redraw pushed it further right: the staircase of repeated, ever-more-
    # indented labels a user reported was not a redraw bug, it was the cursor
    # never returning to column 0. Same for the rewind below — `\033[{n}A`
    # preserves the column too.
    RETURN = '\r'
    CLEAR_LINE = '\033[2K'
    CURSOR_UP_N = '\033[{}A'
    ERASE_DOWN = '\033[J'
    HIDE_CURSOR = '\033[?25l'
    SHOW_CURSOR = '\033[?25h'

    selected = 0
    last_rows = 0

    def draw() -> int:
        rows = 0
        # An empty title draws no row at all. `ask_user_question` prints the
        # question itself -- with its header badge and its own spacing -- and
        # then opens the picker; passing the same text through as a title put
        # that sentence on screen twice, one line apart.
        if title:
            line = f'  {BOLD}{title}{RESET}'
            sys.stdout.write(RETURN + CLEAR_LINE + line + '\n')
            rows += _drawn_rows(line)
        for i, label in enumerate(items):
            marker = f'{GREEN}▶{RESET} ' if i == selected else '  '
            body = f'{BOLD}{label}{RESET}' if i == selected else label
            line = f'  {marker}{body}'
            sys.stdout.write(RETURN + CLEAR_LINE + line + '\n')
            # Physical rows, not one per item: a wrapped label made the rewind
            # below land mid-menu and erase the rest. See `_drawn_rows`.
            rows += _drawn_rows(line)
        foot = footer or '↑↓ move · Enter select · Esc cancel'
        line = f'  {DIM}{foot}{RESET}'
        sys.stdout.write(RETURN + CLEAR_LINE + line + '\n')
        rows += _drawn_rows(line)
        return rows

    def redraw():
        nonlocal last_rows
        if last_rows:
            sys.stdout.write(CURSOR_UP_N.format(last_rows) + RETURN + ERASE_DOWN)
        last_rows = draw()
        sys.stdout.flush()

    # One reader on the console at a time. Without this the adapter's Esc
    # poller keeps calling getwch() from its spinner thread and swallows a
    # byte of every two-byte arrow sequence, so the selection stops moving
    # and the screen fills with half-drawn menus. See core/console.py.
    CONSOLE.acquire()
    sys.stdout.write(HIDE_CURSOR)
    try:
        redraw()
        while True:
            ch = msvcrt.getwch()
            if ch in ('\xe0', '\x00'):
                # Extended-key prefix: arrows report 'H' (up) / 'P' (down)
                # for the byte that follows, same as the main prompt reader.
                ext = msvcrt.getwch()
                if ext == 'H':
                    selected = (selected - 1) % len(items)
                    redraw()
                elif ext == 'P':
                    selected = (selected + 1) % len(items)
                    redraw()
                continue
            if ch == '\r':
                return selected
            if ch in ('\x1b', 'q', 'Q'):
                return -1
            if ch.isdigit() and ch != '0':
                n = int(ch)
                if n <= len(items):
                    return n - 1
    finally:
        # Leaving the list on screen is right for /config, where the next
        # thing printed is the result of the choice. It is wrong when the
        # caller is about to prompt again — you end up typing underneath a
        # menu that no longer accepts input.
        if erase_on_exit and last_rows:
            sys.stdout.write(CURSOR_UP_N.format(last_rows) + RETURN + ERASE_DOWN)
        sys.stdout.write(SHOW_CURSOR)
        sys.stdout.flush()
        CONSOLE.release()


def _drawn_rows(line: str) -> int:
    """How many physical rows one written line will occupy.

    Both pickers rewind with `\\033[{n}A` and counted `n` as one per item —
    true only while nothing wraps. `ask_user_question` builds its labels as
    `label + padding + description`, which routinely exceeds the terminal
    width; the rewind then landed mid-menu and `\\033[J` erased whatever was
    below it, so each keypress ate more of the screen. That is the "I switch
    and it disappears" report.

    The main prompt reader already solves this exactly — see `_rows()` in
    `read_input_with_suggestions`, whose docstring explains that an exact row
    count is "what makes the cursor arithmetic safe". This is the same
    measurement for the menus.
    """
    from text_display import display_width, strip_ansi, term_columns
    width = max(20, term_columns() - 1)
    plain = strip_ansi(line)
    if not plain:
        return 1
    return max(1, -(-display_width(plain) // width))   # ceil-divide


def _fit_label(label: str, desc: str, pad: int = 0) -> str:
    """A menu row that fits the terminal, description clipped rather than wrapped.

    Wrapping is not merely ugly here: a wrapped row breaks the rewind
    arithmetic above. `_drawn_rows` makes a wrapped row survivable; this makes
    it rare.
    """
    from text_display import display_width, shorten, term_columns
    padding = " " * max(0, pad - display_width(label))
    if not desc:
        return f'{label}{padding}'.rstrip()
    room = term_columns() - display_width(label) - len(padding) - 10
    if room < 12:
        return f'{label}{padding}'.rstrip()
    return f'{label}{padding}   {DIM}{shorten(desc, room)}{RESET}'


def _numbered_checklist_fallback(title: str, items: list[str],
                                 footer: str = "") -> Optional[list[int]]:
    """Digit-driven stand-in for `_arrow_checklist` when msvcrt is unavailable."""
    print(f'  {BOLD}{title}{RESET}')
    for i, label in enumerate(items, 1):
        print(f'    {CYAN}{i}{RESET}  {label}')
    try:
        sel = input(f'  {DIM}{footer or "Comma-separated numbers, Enter for none"}:{RESET} ').strip()
    except (EOFError, KeyboardInterrupt):
        print()
        return None
    if not sel:
        return []
    picked = []
    for part in sel.split(','):
        part = part.strip()
        if part.isdigit() and 1 <= int(part) <= len(items):
            picked.append(int(part) - 1)
    return picked


def _arrow_checklist(title: str, items: list[str], footer: str = "") -> Optional[list[int]]:
    """Multi-select sibling of `_arrow_menu`.

    ↑↓ moves, Space toggles the box under the cursor, a digit toggles that
    item directly, Enter confirms whatever is checked (possibly nothing),
    Esc cancels — returning `None` there rather than `[]` lets a caller tell
    "confirmed an empty selection" from "backed out" apart.
    """
    if not items:
        return []
    try:
        import msvcrt
    except ImportError:
        return _numbered_checklist_fallback(title, items, footer)

    RETURN = '\r'          # see `_arrow_menu`: LF does not return the cursor
    CLEAR_LINE = '\033[2K'
    CURSOR_UP_N = '\033[{}A'
    ERASE_DOWN = '\033[J'
    HIDE_CURSOR = '\033[?25l'
    SHOW_CURSOR = '\033[?25h'

    cursor = 0
    checked = [False] * len(items)
    last_rows = 0

    def draw() -> int:
        rows = 0
        if title:                      # see `_arrow_menu.draw` for why
            line = f'  {BOLD}{title}{RESET}'
            sys.stdout.write(RETURN + CLEAR_LINE + line + '\n')
            rows += _drawn_rows(line)
        for i, label in enumerate(items):
            box = f'{GREEN}[x]{RESET}' if checked[i] else '[ ]'
            marker = f'{GREEN}▶{RESET} ' if i == cursor else '  '
            body = f'{BOLD}{label}{RESET}' if i == cursor else label
            line = f'  {marker}{box} {body}'
            sys.stdout.write(RETURN + CLEAR_LINE + line + '\n')
            rows += _drawn_rows(line)      # physical rows — see `_drawn_rows`
        foot = footer or 'Space toggle · Enter confirm · Esc cancel'
        line = f'  {DIM}{foot}{RESET}'
        sys.stdout.write(RETURN + CLEAR_LINE + line + '\n')
        rows += _drawn_rows(line)
        return rows

    def redraw():
        nonlocal last_rows
        if last_rows:
            sys.stdout.write(CURSOR_UP_N.format(last_rows) + RETURN + ERASE_DOWN)
        last_rows = draw()
        sys.stdout.flush()

    # One reader on the console at a time. Without this the adapter's Esc
    # poller keeps calling getwch() from its spinner thread and swallows a
    # byte of every two-byte arrow sequence, so the selection stops moving
    # and the screen fills with half-drawn menus. See core/console.py.
    CONSOLE.acquire()
    sys.stdout.write(HIDE_CURSOR)
    try:
        redraw()
        while True:
            ch = msvcrt.getwch()
            if ch in ('\xe0', '\x00'):
                ext = msvcrt.getwch()
                if ext == 'H':
                    cursor = (cursor - 1) % len(items)
                    redraw()
                elif ext == 'P':
                    cursor = (cursor + 1) % len(items)
                    redraw()
                continue
            if ch == ' ':
                checked[cursor] = not checked[cursor]
                redraw()
                continue
            if ch == '\r':
                return [i for i, c in enumerate(checked) if c]
            if ch in ('\x1b', 'q', 'Q'):
                return None
            if ch.isdigit() and ch != '0':
                n = int(ch)
                if n <= len(items):
                    checked[n - 1] = not checked[n - 1]
                    redraw()
    finally:
        sys.stdout.write(SHOW_CURSOR)
        sys.stdout.flush()
        CONSOLE.release()


def _handle_settings(arg: str) -> str:
    """Show the feature switches, or toggle one by name.

    Names are matched on a prefix so `/settings stream` works; an ambiguous
    prefix lists the candidates rather than picking one, because silently
    toggling the wrong switch is a worse outcome than a second keystroke.
    """
    current = features(refresh=True)
    name = (arg or "").strip().lower().replace("-", "_")

    if name:
        matches = [f for f in core_features.FEATURES
                   if f["key"] == name or f["key"].startswith(name)]
        if not matches:
            return (f'  {RED}✗{RESET} No such setting: {name}\n'
                    f'  {DIM}One of: '
                    f'{", ".join(core_features.FEATURE_KEYS)}{RESET}')
        if len(matches) > 1:
            names = ", ".join(m["key"] for m in matches)
            return (f'  {YELLOW}⚠{RESET} "{name}" matches several settings: '
                    f'{names}')
        key = matches[0]["key"]
        updated = core_features.toggle(current, key)
        save_features(updated)
        now_on = updated.enabled(key)
        state = f'{GREEN}on{RESET}' if now_on else f'{DIM}off{RESET}'
        note = ""
        if key == "debug_view":
            note = (f'\n  {DIM}Press Ctrl+Alt+X (or /debug) to view the last '
                    f'request and response.{RESET}' if now_on else
                    f'\n  {DIM}Recorded payloads discarded.{RESET}')
        elif key == "prefill_context":
            note = f'\n  {DIM}Applies to the next new session.{RESET}'
        elif key == "advanced_diagnostics":
            # Says what changes, not that something changed. The switch is
            # invisible until a turn does something unusual, so "on" with no
            # further word looks like nothing happened.
            note = (f'\n  {DIM}Retries, output limits, cache and token counts, '
                    f'why each turn ended, and the detail behind every '
                    f'error.{RESET}' if now_on else
                    f'\n  {DIM}Back to errors, permissions and the reply '
                    f'itself.{RESET}')
        return f'  {GREEN}✓{RESET} {matches[0]["label"]}: {state}{note}'

    lines = [
        f'  {BOLD}Settings{RESET}',
        f'  {DIM}{"─" * 58}{RESET}',
    ]
    for spec in core_features.FEATURES:
        on = current.enabled(spec["key"])
        mark = f'{GREEN}✓{RESET}' if on else f'{RED}✕{RESET}'
        lines.append(f'  {mark} {spec["label"]:<22}{DIM}{spec["detail"]}{RESET}')
    lines.append('')
    lines.append(f'  {DIM}Toggle with{RESET} {CYAN}/settings <name>{RESET}'
                 f'{DIM} — e.g. /settings streaming{RESET}')
    return '\n'.join(lines)


def _open_debug_window() -> str:
    """Open a second console that tails the debug log live.

    The chat REPL owns the console it runs in — it reads raw keystrokes and
    repaints its own prompt — so a live view cannot share that screen without
    fighting it for the cursor. A separate window is not a workaround for
    that, it is the only arrangement in which "watch the traffic while you use
    the session" is a thing you can actually do.

    PowerShell's `Get-Content -Wait` is the tail: it is present on every
    supported Windows install, so this adds no dependency. The window stays up
    after the session ends (`-NoExit`) because the last exchange is usually
    the interesting one.
    """
    if sys.platform != "win32":
        return (f'  {YELLOW}⚠{RESET} A separate debug window is Windows-only.\n'
                f'  {DIM}The log is at {debug_log.live_file()} — tail it with'
                f'{RESET} {CYAN}tail -f{RESET}')

    path = debug_log.live_file()
    if not path:
        return f'  {RED}✗{RESET} No debug log file is active.'

    # Single-quoted for PowerShell and doubled to escape, so a path containing
    # an apostrophe cannot end the string and run the rest as a command.
    ps_path = str(path).replace("'", "''")
    command = (f"$Host.UI.RawUI.WindowTitle = 'TOMAS — live debug'; "
               f"Write-Host 'Tailing {ps_path}' -ForegroundColor Cyan; "
               f"Get-Content -LiteralPath '{ps_path}' -Wait -Tail 200")
    try:
        subprocess.Popen(
            ["cmd", "/c", "start", "TOMAS Debug", "powershell",
             "-NoProfile", "-NoExit", "-Command", command],
            creationflags=getattr(subprocess, "CREATE_NEW_CONSOLE", 0))
    except Exception as exc:
        return (f'  {RED}✗{RESET} Could not open the debug window: {exc}\n'
                f'  {DIM}The log is still being written to {path}{RESET}')
    return (f'  {GREEN}✓{RESET} Debug window opened — it updates live as this '
            f'session talks to the model.\n'
            f'  {DIM}{path}{RESET}')


def _handle_debug(arg: str) -> str:
    """Show what actually went over the wire on recent model calls.

    Off by default and explicitly switched on, because recording holds a full
    copy of every request — see `core/debug_log.py`. Turning it on here rather
    than only in Settings is deliberate: the moment someone wants this, they
    want it for the call that just happened, and making them leave the chat to
    arm it guarantees they miss it.
    """
    arg = (arg or "").strip().lower()

    if arg in ("on", "off"):
        updated = core_features.Features(
            **{**features(refresh=True).to_dict(), "debug_view": arg == "on"})
        save_features(updated)
        if arg == "on":
            return (f'  {GREEN}✓{RESET} Debug recording on.\n'
                    f'  {DIM}Ctrl+Alt+X opens a live window; /debug shows the '
                    f'last payload here.{RESET}')
        return f'  {GREEN}✓{RESET} Debug recording off; captured payloads discarded.'

    if arg == "window":
        if not debug_log.is_enabled():
            _handle_debug("on")
        return _open_debug_window()

    if not debug_log.is_enabled():
        return (f'  {DIM}Debug recording is off.{RESET}\n'
                f'  {CYAN}/debug on{RESET}{DIM} to start capturing requests and '
                f'responses, then{RESET} {CYAN}Ctrl+Alt+X{RESET}{DIM} for a live '
                f'window.{RESET}')

    captured = debug_log.entries()
    if not captured:
        return (f'  {DIM}Nothing captured yet — recording is on, but no model '
                f'call has been made since.{RESET}')

    if arg == "schemas":
        tools = (captured[-1].request or {}).get("tools") or []
        if not tools:
            return f'  {DIM}No tool schemas were sent on the last call.{RESET}'
        return (f'  {BOLD}Tool schemas sent on call #{captured[-1].seq}{RESET} '
                f'{DIM}({len(tools)} tools){RESET}\n'
                + json.dumps(tools, indent=2, ensure_ascii=False, default=str))

    if arg.isdigit():
        wanted = int(arg)
        entry = next((e for e in captured if e.seq == wanted), None)
        if entry is None:
            available = ", ".join(f"#{e.seq}" for e in captured)
            return (f'  {RED}✗{RESET} No captured call #{wanted}.\n'
                    f'  {DIM}Available: {available}{RESET}')
        return debug_log.as_json(entry)

    if arg in ("", "last"):
        return debug_log.as_json(captured[-1])

    if arg == "list":
        lines = [f'  {BOLD}Captured calls{RESET}',
                 f'  {DIM}{"─" * 58}{RESET}']
        for entry in captured:
            outcome = (f'{RED}{entry.error[:28]}{RESET}' if entry.error
                       else (entry.response.get("stop_reason") or "—"))
            lines.append(
                f'  {CYAN}#{entry.seq}{RESET} {entry.path:<13}'
                f'{DIM}{entry.message_count} msg · {entry.tool_count} tools · '
                f'{entry.elapsed_ms:,}ms · {RESET}{outcome}')
        lines.append('')
        lines.append(f'  {DIM}/debug <n> for one call · /debug schemas for the '
                     f'tool block{RESET}')
        return '\n'.join(lines)

    return (f'  {DIM}Usage:{RESET} /debug [on|off|last|list|<n>|schemas]')


def _context_controls_disabled(what: str) -> str:
    """The one message both context commands give when switched off.

    Shared so the two cannot drift into describing the same switch
    differently, and so it always names the way back — a command that refuses
    without saying which setting refused it is a dead end.
    """
    return (f'  {YELLOW}⚠{RESET} /{what} is off — "Context controls" is '
            f'disabled in Settings.\n'
            f'  {DIM}Turn it back on with{RESET} {CYAN}/settings context{RESET}'
            f'{DIM}, or in the Settings menu.{RESET}')


#: Subfolder offered as the tidy default for exports. Relative to EXPORT_DIR,
#: created on demand — an export that has to be filed by hand afterwards is
#: one the user stops doing.
EXPORT_SUBFOLDER = "exports"

#: Whether `/export` may stop and ask where to save.
#:
#: A separate switch from `sys.stdin.isatty()`, because that test does not
#: actually answer the question here: the picker reads the console through
#: `msvcrt.getwch()`, which bypasses stdin entirely, so a redirected or closed
#: stdin does not stop it blocking. Measured — `tests/test_command_surface`
#: calls every advertised command, reached `/export`, and hung on a keypress
#: that no test can send. Anything driving TOMAS without a person at the
#: keyboard sets this False.
EXPORT_PROMPT = True


def _ask_export_destination(filename: str) -> Optional[Path]:
    """Where should this export go? Returns the directory, or None to cancel.

    Asked rather than assumed. Writing straight into the project root is fine
    once and litter by the fifth time, and the answer genuinely varies — the
    repo for something to be committed, a subfolder to keep it out of the way,
    somewhere else entirely to hand to another program.

    Falls back to the project root without prompting whenever asking is not
    possible, because an export that blocks on input nobody can give is worse
    than one filed in the obvious place.
    """
    root = Path(EXPORT_DIR)
    if not EXPORT_PROMPT:
        return root
    if not bool(getattr(sys.stdin, "isatty", lambda: False)()):
        return root

    options = [
        f'  Project root        {DIM}{root}{RESET}',
        f'  Subfolder           {DIM}{root / EXPORT_SUBFOLDER}{RESET}',
        f'  Another path…       {DIM}type it in{RESET}',
    ]
    print(f'\n  {BOLD}Save {filename} where?{RESET}')
    choice = _arrow_menu('', options, footer='↑↓ choose · Enter save · Esc cancel',
                         erase_on_exit=True)
    if choice < 0:
        return None
    if choice == 0:
        return root
    if choice == 1:
        target = root / EXPORT_SUBFOLDER
        try:
            target.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            print(f'  {YELLOW}⚠{RESET} Could not create {target}: {exc} — '
                  f'using the project root instead.')
            return root
        return target

    try:
        typed = input(f'  {CYAN}Path:{RESET} ').strip().strip('"').strip("'")
    except (EOFError, KeyboardInterrupt):
        return None
    if not typed:
        return None
    target = Path(typed).expanduser()
    # A path that names a file is a common way to answer "where?", so the
    # parent is taken rather than refusing — but only when the name matches
    # what is about to be written, or an unrelated trailing segment would be
    # silently discarded.
    if target.suffix and target.name == filename:
        target = target.parent
    try:
        target.mkdir(parents=True, exist_ok=True)
    except OSError as exc:
        print(f'  {RED}✗{RESET} Cannot use {target}: {exc}')
        return None
    return target


def _handle_export(messages: list, arg: str) -> str:
    """Write the conversation to a file the user can open.

    Sessions are already persisted as JSON under `~/.tomas/sessions/`, but
    that is the agent's own record: named by timestamp, in a schema built for
    reloading, somewhere the user does not look. This writes a copy where they
    are working, in a format they choose — which is the difference between
    "the data exists" and "the user has it".
    """
    if not features().enabled("context_controls"):
        return _context_controls_disabled("export")
    if not messages:
        return f'  {DIM}Nothing to export — the conversation is empty.{RESET}'

    # `/export [txt|json] [path]` — a path given here skips the picker, which
    # is what makes the command usable from a script as well as by hand.
    parts = (arg or "").strip().split(maxsplit=1)
    fmt = (parts[0] or "txt").lower() if parts else "txt"
    where = parts[1].strip().strip('"').strip("'") if len(parts) > 1 else ""
    if fmt not in ("txt", "json"):
        return f'  {DIM}Usage:{RESET} /export [txt|json] [path]'

    stamp = time.strftime("%Y%m%d_%H%M%S")
    name = f"tomas_conversation_{stamp}.{fmt}"
    if where:
        directory = Path(where).expanduser()
        if directory.suffix and directory.name == name:
            directory = directory.parent
        try:
            directory.mkdir(parents=True, exist_ok=True)
        except OSError as exc:
            return f'  {RED}✗{RESET} Cannot use {directory}: {exc}'
    else:
        directory = _ask_export_destination(name)
    if directory is None:
        return f'  {DIM}Export cancelled — nothing was written.{RESET}'
    path = directory / name

    try:
        if fmt == "json":
            payload = {
                "exported": time.strftime("%Y-%m-%d %H:%M:%S"),
                "model": _get_model(),
                "message_count": len(messages),
                "system_prompt": build_system_prompt(),
                "messages": messages,
            }
            path.write_text(json.dumps(payload, indent=2, ensure_ascii=False,
                                       default=str), encoding="utf-8")
        else:
            lines = [
                f"TOMAS conversation — {time.strftime('%Y-%m-%d %H:%M:%S')}",
                f"Model: {_get_model()}",
                f"Messages: {len(messages)}",
                "=" * 70,
                "",
                "--- SYSTEM PROMPT " + "-" * 52,
                build_system_prompt(),
                "",
            ]
            for message in messages:
                role = str(message.get("role", "?")).upper()
                lines.append("-" * 70)
                lines.append(f"[{role}]")
                lines.append(_render_message_for_export(message))
                lines.append("")
            path.write_text("\n".join(lines), encoding="utf-8")
    except OSError as exc:
        return f'  {RED}✗{RESET} Could not write {path.name}: {exc}'

    size = path.stat().st_size
    return (f'  {GREEN}✓{RESET} Exported {len(messages)} messages to '
            f'{CYAN}{path.name}{RESET}\n'
            f'  {DIM}{path}  ·  {size:,} bytes{RESET}')


def _render_message_for_export(message: dict) -> str:
    """One message as readable text, tool calls and results included.

    The whole point of a .txt export is that it can be read without a JSON
    viewer, so a content list is unpacked into labelled sections rather than
    dumped — a transcript whose interesting half is `[{'type': 'tool_use',
    ...}]` has not been exported into text, only into a file with a .txt name.
    """
    content = message.get("content")
    if isinstance(content, str):
        return content
    if not isinstance(content, list):
        return str(content)

    parts = []
    for block in content:
        if not isinstance(block, dict):
            parts.append(str(block))
            continue
        kind = block.get("type")
        if kind == "text":
            parts.append(block.get("text", ""))
        elif kind == "tool_use":
            args = json.dumps(block.get("input", {}), ensure_ascii=False,
                              default=str)
            parts.append(f"  → tool call: {block.get('name', '?')}({args})")
        elif kind == "tool_result":
            body = block.get("content", "")
            if isinstance(body, list):
                body = "\n".join(
                    b.get("text", "") if isinstance(b, dict) else str(b)
                    for b in body)
            parts.append(f"  ← tool result: {body}")
        else:
            parts.append(json.dumps(block, ensure_ascii=False, default=str))
    return "\n".join(p for p in parts if p)


def _config_menu(messages: list) -> str:
    """Arrow-key /config menu: pick provider, model or mode without leaving chat.

    Every actual switch is delegated to `/provider`, `/model` and `/mode`
    (via a recursive `handle_slash_command` call), so this command carries no
    switching logic of its own that could drift out of sync with theirs — it
    is purely a picker sitting in front of commands that already work.
    """
    import provider_manager

    while True:
        active = provider_manager.get_active()
        providers = provider_manager.visible_providers()
        mode = current_mode_name()

        top_items = [
            f'Provider   {active.name if active else "(none configured)"}',
            f'Model      {_get_model()}',
            f'Mode       {mode}',
            'Close',
        ]
        top = _arrow_menu('TOMAS Config', top_items,
                          footer='↑↓ move · Enter select · Esc close')
        if top in (-1, 3):
            return f'  {GREEN}✓{RESET} Config closed.'

        if top == 0:  # ── Provider ──
            if not providers:
                print(f'  {DIM}No providers configured yet — use agent_cli.py to add one.{RESET}')
                continue
            labels = [
                f'{"●" if active and p.name == active.name else "○"} {p.name}  '
                f'{DIM}({p.type} · {p.model or "no model set"}){RESET}'
                for p in providers
            ] + ['◀ Back']
            sel = _arrow_menu('Switch Provider', labels)
            if sel in (-1, len(providers)):
                continue
            print(handle_slash_command(f'provider {providers[sel].name}', messages))

        elif top == 1:  # ── Model ──
            if active is None:
                print(f'  {RED}No active provider — set one first.{RESET}')
                continue
            try:
                models = provider_manager.available_models(active.name)
            except Exception as e:
                models = []
                print(f'  {YELLOW}Could not fetch the model list: {e}{RESET}')
            current = _get_model()
            if models:
                labels = [f'{"●" if m == current else "○"} {m}' for m in models]
                labels += ['✎ Type a model name…', '◀ Back']
                sel = _arrow_menu(f'Switch Model  ({active.name})', labels)
                if sel in (-1, len(labels) - 1):
                    continue
                if sel == len(models):  # "type a model name" entry
                    target = input(f'  {DIM}Model name:{RESET} ').strip()
                    if not target:
                        continue
                else:
                    target = models[sel]
            else:
                target = input(f'  {DIM}No model list from the endpoint — type a model '
                                f'name (Enter to cancel):{RESET} ').strip()
                if not target:
                    continue
            print(handle_slash_command(f'model {target}', messages))

        elif top == 2:  # ── Mode ──
            modes = ["auto", "default", "strict", "yolo", "bypass"]
            labels = [f'{"●" if m == mode else "○"} {m}' for m in modes] + ['◀ Back']
            sel = _arrow_menu('Switch Mode', labels)
            if sel in (-1, len(modes)):
                continue
            print(handle_slash_command(f'mode {modes[sel]}', messages))


def handle_slash_command(cmd_args: str, messages: list) -> str | None:
    """
    Handle a slash command (text after the '/').

    Returns:
      - response text to display
      - "__exit__"     to break the main loop
      - "__continue__" to skip agent processing and continue
      - None           if command produced no output
    """
    global mcp_manager, AUTO_APPROVE_LOW, YOLO_MODE, BYPASS_MODE
    parts = cmd_args.strip().split(maxsplit=1)
    cmd = parts[0].lower() if parts else ""

    # ── Just "/" alone — show all commands ──
    if not cmd:
        return _show_commands()

    # ── Handle aliases ──
    if cmd in ("q", "quit"):
        return "__exit__"
    if cmd in ("h", "?"):
        cmd = "help"

    # ── Match against known commands ──
    exact = [name for name in SLASH_COMMANDS if name == cmd]
    partial = [name for name in SLASH_COMMANDS if name.startswith(cmd)]

    if not exact and not partial:
        # No matches at all — show error with hint
        close = [name for name in SLASH_COMMANDS
                 if cmd in name or any(cmd[i:i+2] in name for i in range(len(cmd) - 1))]
        if close:
            return (
                f'  {YELLOW}Unknown command "{cmd}". Did you mean:{RESET}\n'
                + _show_commands()
            )
        return f'  {RED}Unknown command "{cmd}". Type{RESET} {CYAN}/{RESET} to see all commands.'

    if partial and not exact and len(partial) == 1:
        # Single unambiguous partial match — treat as the full command
        cmd = partial[0]
        exact = [cmd]

    if partial and not exact:
        # Multiple partial matches — show suggestions
        return _show_commands(cmd)

    # ── Execute the matched command ──
    cmd = exact[0]

    if cmd == "help":
        return _show_commands()

    if cmd == "clear":
        if not features().enabled("context_controls"):
            return _context_controls_disabled("clear")
        cleared = len(messages)
        messages.clear()
        # A cleared conversation is a new session; its accounting starts over.
        reset_session_state()
        return (f'  {GREEN}✓{RESET} Conversation cleared — {cleared} messages '
                f'dropped ({_get_model()}).')

    if cmd == "status":
        model_status = _get_model()
        window = resolve_context_window()
        cw = window.tokens
        mode = current_mode_name()
        # What the *next* request will carry, not what the session has billed.
        # The old line divided cumulative session tokens by the window, so a
        # long session read "418% of 1,000,000 used" — two different
        # quantities with the same unit.
        live_tok = (_estimate_tokens(messages) + TOOL_TOKENS + output_reserve())
        fill = (live_tok / cw * 100) if cw else 0.0
        note = '' if window.trusted else f' {YELLOW}— not measured{RESET}'
        lines = [
            f'  {BOLD}TOMAS Status{RESET}',
            f'  {DIM}{"─" * 46}{RESET}',
            f'  {CYAN}◎{RESET}  Model:     {model_status}',
            f'  {CYAN}▣{RESET}  Context:   {live_tok:,} / {cw:,} tokens '
            f'({fill:.0f}%) {DIM}·{RESET} {window.source}{note}',
            f'  {CYAN}⚙{RESET}  Mode:      {mode_color(mode)}{mode}{RESET}',
            f'  {CYAN}✉{RESET}  Messages:  {len(messages)} in history',
        ]
        try:
            from mcp_manager import read_mcp_servers
            servers = read_mcp_servers()
            enabled = sum(1 for c in servers.values() if not c.get("disabled"))
            total = len(servers)
            lines.append(f'  {CYAN}⬡{RESET}  MCP:       {enabled}/{total} servers enabled')
        except Exception:
            pass
        if mcp_manager:
            lines.append(f'  {CYAN}⚡{RESET}  MCP tools: {len(mcp_manager.tools)} loaded')
        lines.append(f'  {CYAN}⌨{RESET}  Toggle:    {DIM}⇧+Space{RESET}  ·  '
                     f'{DIM}{MODE_KEYS_HINT}{RESET}')
        # ── Token stats ──
        s = _session_tokens
        if s["calls"] > 0:
            # Billed across the whole session — this accumulates and is not a
            # fraction of anything. The window fill is on the Context line.
            lines.append(f'  {CYAN}≡{RESET}  Billed:    {s["input"]:,} in · '
                         f'{s["output"]:,} out ({s["calls"]} calls)')
        if _context_log:
            saved = sum(e.get("reclaimed_tokens", 0) for e in _context_log)
            lines.append(f'  {CYAN}⇩{RESET}  Compacted: {len(_context_log)}× '
                         f'{DIM}(~{saved:,} tokens released){RESET}')
        return '\n'.join(lines)

    if cmd == "version":
        from version import VERSION, LAST_UPDATED, git_info
        lines = [
            f'  {BOLD}TOMAS{RESET}',
            f'  {DIM}{"─" * 46}{RESET}',
            f'  {CYAN}◎{RESET}  Version:       {VERSION}',
            f'  {CYAN}▣{RESET}  Last updated:  {LAST_UPDATED}',
        ]
        # Dev-checkout-only supplement — see version.git_info's docstring for
        # why this is never the primary source.
        info = git_info()
        if info:
            commit_hash, commit_date = info
            lines.append(f'  {DIM}⎇  dev checkout — commit {commit_hash} '
                         f'({commit_date}){RESET}')
        return '\n'.join(lines)

    if cmd == "model":
        sub = parts[1].strip() if len(parts) > 1 else ""
        if not sub:
            window = resolve_context_window()
            return (f'  {CYAN}◎{RESET} {BOLD}Model:{RESET} {_get_model()} '
                    f'{DIM}·{RESET} {window.tokens:,} token context '
                    f'{DIM}({window.source}){RESET}')

        if sub.lower() in ("list", "ls"):
            import provider_manager
            active = provider_manager.get_active()
            if active is None:
                return (f'  {DIM}No provider is configured. Use{RESET} '
                        f'{CYAN}/provider list{RESET}{DIM} or agent_cli.py.{RESET}')
            current = _get_model()
            lines = [f'  {BOLD}Models on {active.name}{RESET}', f'  {DIM}{"─" * 46}{RESET}']
            try:
                models = provider_manager.available_models(active.name)
            except Exception as e:
                return f'  {RED}Could not list models: {e}{RESET}'
            # This is the raw upstream list — the same one that told the user
            # a dead model was there in the first place ("This model is
            # listed but not currently being served"). Cross-referencing the
            # observed-refusal record here means someone told "pick another
            # with /model" is not steered right back into another dead one.
            unavailable = {}
            if active.type == "zen":
                try:
                    import zen_catalog
                    unavailable = zen_catalog.unavailable()
                except Exception:
                    pass
            if not models:
                lines.append(f'  {DIM}The endpoint did not report a model list.{RESET}')
            else:
                for m in models:
                    marker = f'{GREEN}●{RESET}' if m == current else f'{DIM}○{RESET}'
                    flag = (f'  {YELLOW}(unavailable — refused within the '
                           f'last 24h){RESET}') if m in unavailable else ''
                    lines.append(f'    {marker} {m}{flag}')
            lines.append('')
            lines.append(f'  {DIM}Switch with{RESET} {CYAN}/model <name>{RESET}')
            return '\n'.join(lines)

        # ── Switch the model on the active provider, in place ──
        import provider_manager
        active = provider_manager.get_active()
        if active is None:
            return (f'  {RED}No active provider — configure one first '
                    f'({RESET}{CYAN}/provider list{RESET}{RED}).{RESET}')
        active.model = sub
        provider_manager.save(active, activate_it=True)
        # `activate()`, not a bare `apply_env("AGENT_MODEL", ...)`. The bare
        # call only ever wrote the model name — `save()` above never touches
        # `os.environ` at all — so if a *different* provider's `activate()`
        # was the last one to run in this process, `ANTHROPIC_BASE_URL` and
        # `ANTHROPIC_API_KEY` stayed pointed at it. `/model` then sent this
        # provider's model name to that stale endpoint: an Ollama cloud model
        # name reached OpenRouter and came back "not a valid model ID" — a
        # provider mismatch reported as a bad model name. `activate()` is the
        # one place that reapplies base_url/api_key/headers together with the
        # model, which is exactly what `agent_cli.py`'s TUI already does for
        # every other model switch.
        provider_manager.activate(active.name)
        # Before `_refresh_context_window`, not after: that function reads the
        # stored capabilities, and the stored capabilities still describe the
        # model we just switched away from.
        try:
            caps = provider_manager.refresh_for_model(active)
        except Exception:
            caps = active.capabilities
        reinit_client()
        cw = _refresh_context_window()
        note = (f'  {DIM}Endpoint capabilities carried over — run{RESET} '
                f'{CYAN}/provider probe{RESET} {DIM}to re-measure.{RESET}')
        if active.type == "ollama":
            note = (f'  {DIM}Tool use:{RESET} {"yes" if caps.tool_use else "no"}'
                    f'{DIM} · Vision:{RESET} {"yes" if caps.vision else "no"}')
        return (f'  {GREEN}✓{RESET} Switched model to {BOLD}{sub}{RESET} on {active.name} '
                f'({cw:,} token context).\n' + note)

    if cmd == "budget":
        return _budget_command(parts[1].strip() if len(parts) > 1 else "")

    if cmd == "mode":
        arg = parts[1].lower() if len(parts) > 1 else ""
        # One writer for the flags: /mode must go through set_mode() exactly
        # like the F5–F9 keys do. It used to write AUTO_APPROVE_LOW / YOLO_MODE
        # / BYPASS_MODE directly, so "auto" and "default" silently left YOLO or
        # BYPASS set — a safety reset that reported success while changing
        # nothing. Aliases map onto the canonical names set_mode understands;
        # unknown names fall through to the status display below.
        alias = {
            "automatic": "auto", "on": "auto", "1": "auto", "yes": "auto",
            "normal": "default", "off": "default", "0": "default", "no": "default",
            "bypass-permissions": "bypass", "nonstop": "bypass",
        }
        canonical = alias.get(arg, arg)
        if canonical in ALL_MODES:
            set_mode(canonical)
            if canonical == "auto":
                return (
                    f'  {GREEN}✓{RESET} Mode set to {BOLD}auto{RESET} — '
                    f'low-risk tools will be auto-approved.'
                )
            if canonical == "default":
                return (
                    f'  {GREEN}✓{RESET} Mode set to {BOLD}default{RESET} — '
                    f'you will confirm each tool use.'
                )
            if canonical == "strict":
                return (
                    f'  {GREEN}✓{RESET} Mode set to {BOLD}strict{RESET} — '
                    f'all tools require confirmation, risk overrides cleared.'
                )
            if canonical == "yolo":
                return (
                    f'  {BOLD}{RED}⚡ YOLO mode enabled!{RESET} {RED}All tools auto-approved.{RESET}'
                )
            return (
                f'  {BOLD}{RED}⇥ BYPASS mode enabled!{RESET} '
                f'{RED}All tools auto-approved and the turn will not stop to '
                f'ask whether to continue.{RESET}\n'
                f'  {DIM}Bounded at {MAX_AUTO_CONTINUATIONS} automatic budget '
                f'extensions ({MAX_TOOL_CALLS_PER_TURN * (MAX_AUTO_CONTINUATIONS + 1)} '
                f'tool calls per turn); loop detection still applies.{RESET}'
            )
        # No arg or unknown arg — show current mode
        mode = current_mode_name()
        lines = [
            f'  {BOLD}Current Mode{RESET}',
            f'  {DIM}{"─" * 46}{RESET}',
            f'  {CYAN}⚙{RESET}  Mode:        {BOLD}{mode}{RESET}',
            f'  {CYAN}✓{RESET}  Auto-approve: {"ON" if AUTO_APPROVE_LOW else "OFF"}',
            f'  {CYAN}⇥{RESET}  Auto-continue: '
            f'{"ON — up to %d extensions" % MAX_AUTO_CONTINUATIONS if BYPASS_MODE else "OFF — asks at the budget checkpoint"}',
            '',
            f'  {DIM}Quick keys:{RESET}',
            f'    {DIM}⇧+Space{RESET}  — toggle  ·  {DIM}Tab{RESET} cycles  ·  '
            f'{DIM}F5{RESET} auto  {DIM}F6{RESET} default  {DIM}F7{RESET} strict  '
            f'{DIM}F8{RESET} yolo  {DIM}F9{RESET} bypass',
            f'  {DIM}Slash commands:{RESET}',
            f'    /mode auto     — auto-approve low-risk tools',
            f'    /mode default  — ask before every tool',
            f'    /mode strict   — ask for everything, clear overrides',
            f'    /mode yolo     — {RED}auto-approve ALL tools (no prompts){RESET}',
            f'    /mode bypass   — {RED}yolo + never ask whether to continue{RESET}',
        ]
        return '\n'.join(lines)

    # `parts[1]` rather than a shared `arg`: this function derives the
    # argument per command, and the three below are the only ones that need
    # the raw (uncased) tail.
    if cmd == "settings":
        return _handle_settings(parts[1] if len(parts) > 1 else "")

    if cmd == "debug":
        return _handle_debug(parts[1] if len(parts) > 1 else "")

    if cmd == "export":
        return _handle_export(messages, parts[1] if len(parts) > 1 else "")

    if cmd == "config":
        return _config_menu(messages)

    if cmd == "compact":
        if not messages:
            return f'  {DIM}No conversation to compact.{RESET}'
        before = len(messages)
        # `force`: the user asked, so the automatic threshold — including
        # "never" — does not apply. Without it, choosing Never in the budget
        # page silently turned this command into a no-op that still reported
        # success.
        messages[:] = maybe_compact(messages, force=True)
        after = len(messages)
        if after == before:
            return (f'  {DIM}Nothing to gain — the conversation is already '
                    f'smaller than a summary of it would be.{RESET}')
        return f'  {GREEN}✓{RESET} Compacted ({before} → {after} messages).'

    if cmd == "zen":
        try:
            import zen_catalog
            from zen_proxy import check_status
            port = 6446
            running = check_status(port)
            # The catalogue is reported whether or not the proxy runs. Gating
            # it on the daemon was backwards twice over: the daemon is opt-in
            # (openai_adapter translates in-process), and it serves the same
            # stale list this now replaces.
            cat = zen_catalog.catalog()
            base = os.environ.get("ANTHROPIC_BASE_URL", "")
            model = os.environ.get("AGENT_MODEL", "Not set")
            active = cat.get(model)
            lines = [
                f'  {BOLD}OpenCode Zen{RESET}',
                f'  {DIM}{"─" * 46}{RESET}',
                f'  {CYAN}◈{RESET}  Endpoint: {base or "not set"}',
                f'  {CYAN}◎{RESET}  Model:    {active.label if active else model}',
                f'  {"◉" if running else "○"}  Optional local proxy: '
                f'{"running" if running else "not running"}',
                '',
                f'  {DIM}Free models — {len(cat.free())} of {len(cat.models)} '
                f'served ({cat.freshness}):{RESET}',
            ]
            for m in cat.free():
                mark = f'{GREEN}●{RESET}' if m.id == model else f'{DIM}•{RESET}'
                lines.append(f'    {mark} {m.label}')
            if not cat.free():
                lines.append(f'    {DIM}none right now{RESET}')
            lines.append('')
            # Informational only, deliberately not "switch with /model" — the
            # working set is free-tier-only for now (see
            # provider_manager.VISIBLE_PROVIDER_TYPES), so this should not
            # read as an invitation to pick one of them.
            lines.append(f'  {DIM}{len(cat.paid())} further models on this '
                         f'endpoint bill per token (not offered for now).{RESET}')
            return '\n'.join(lines)
        except Exception as e:
            return f'  {YELLOW}Zen: {e}{RESET}'

    if cmd == "provider":
        import provider_manager
        sub_raw = parts[1].strip() if len(parts) > 1 else ""
        sub = sub_raw.lower()

        if sub in ("list", "ls"):
            providers = provider_manager.visible_providers()
            if not providers:
                return f'  {DIM}No provider is configured. Use{RESET} {CYAN}agent_cli.py{RESET}{DIM} to add one.{RESET}'
            active_provider = provider_manager.get_active()
            active_name = active_provider.name if active_provider else None
            lines = [f'  {BOLD}Configured Providers{RESET}', f'  {DIM}{"─" * 46}{RESET}']
            for p in providers:
                marker = f'{GREEN}●{RESET}' if p.name == active_name else f'{DIM}○{RESET}'
                lines.append(f'    {marker} {p.name} {DIM}({p.type} · {p.model or "no model set"}){RESET}')
            lines.append('')
            lines.append(f'  {DIM}Switch with{RESET} {CYAN}/provider <name>{RESET}')
            return '\n'.join(lines)

        # ── Switch to a different configured provider ──
        if sub and sub != "probe":
            names = [p.name for p in provider_manager.visible_providers()]
            exact = [n for n in names if n.lower() == sub]
            partial = [n for n in names if sub in n.lower()]
            target = exact[0] if exact else (partial[0] if len(partial) == 1 else None)
            if target is None:
                if len(partial) > 1:
                    return (f'  {YELLOW}Multiple providers match "{sub_raw}":{RESET} '
                            f'{", ".join(partial)}')
                return (f'  {RED}No configured provider matches "{sub_raw}".{RESET}\n'
                        f'  {DIM}Run{RESET} {CYAN}/provider list{RESET} '
                        f'{DIM}to see what is configured.{RESET}')
            current = provider_manager.get_active()
            if current and target == current.name:
                return f'  {DIM}{target} is already the active provider.{RESET}'
            if not provider_manager.activate(target):
                return f'  {RED}Failed to activate "{target}".{RESET}'
            cw = _refresh_context_window()
            switched = provider_manager.get_active()
            return (f'  {GREEN}✓{RESET} Switched provider to {BOLD}{target}{RESET} — '
                    f'model {switched.model or "(unset)"} ({cw:,} token context).')

        provider = provider_manager.get_active()
        if provider is None:
            return f'  {DIM}No provider is configured. Use{RESET} {CYAN}agent_cli.py{RESET}{DIM} to add one.{RESET}'
        if sub == "probe":
            # A capability learned by degradation is sticky by design, so
            # there has to be a way to re-measure — a transient failure must
            # not disable streaming forever.
            caps = provider_manager.probe_and_persist(provider)
            head = f'  {GREEN}✓{RESET} Re-probed {CYAN}{provider.name}{RESET}'
        else:
            caps = provider.capabilities
            head = f'  {BOLD}{provider.name}{RESET} {DIM}({provider.type}){RESET}'
        when = ('never — using optimistic defaults' if not caps.probed
                else time.strftime('%Y-%m-%d %H:%M', time.localtime(caps.probed_at)))
        yn = lambda v: f'{GREEN}yes{RESET}' if v else f'{YELLOW}no{RESET}'
        # Resolved through the same path as /status and /model. Reading
        # caps.context_window straight off the provider is what let the two
        # commands report different windows for one model.
        window = resolve_context_window(provider.model, refresh=True)
        ctx_line = f'  {DIM}Context window:{RESET} {window.describe()}'
        out = [
            head,
            f'  {DIM}{"─" * 46}{RESET}',
            f'  {DIM}Model:{RESET}          {provider.model or "(unset)"}',
            f'  {DIM}Endpoint:{RESET}       {provider.base_url or "(from env)"}',
            f'  {DIM}Probed:{RESET}         {when}',
            '',
            f'  {DIM}Streaming:{RESET}      {yn(caps.streaming)}',
            f'  {DIM}Tool use:{RESET}       {yn(caps.tool_use)}',
            f'  {DIM}System prompt:{RESET}  {yn(caps.system_prompt)}',
            ctx_line,
            f'  {DIM}Tool ceiling:{RESET}   {caps.max_tools}',
        ]
        # A stale MODEL_CONTEXT_MAP is worth saying out loud: it is a hardcoded
        # table and it will rot. Silence means trusting whichever source
        # happened to answer first.
        diverged = context_window_divergence(provider.model)
        if diverged:
            known, probed = diverged
            out.append(f'  {YELLOW}⚠{RESET}  {DIM}table says {known:,}, endpoint '
                       f'says {probed:,} — MODEL_CONTEXT_MAP may be stale{RESET}')
        out += [
            '',
            f'  {DIM}Run{RESET} {CYAN}/provider probe{RESET} {DIM}to re-measure against the endpoint.{RESET}',
        ]
        return '\n'.join(out)

    if cmd in ("mcp-prompt", "mcp-prompts"):
        # Server-provided prompt templates, surfaced the same way skills are.
        if mcp_manager is None:
            return f'  {DIM}No MCP servers are connected.{RESET}'
        prompts = mcp_manager.list_prompts()
        name = parts[1] if len(parts) > 1 else ""
        if not name:
            if not prompts:
                return f'  {DIM}No connected MCP server provides prompts.{RESET}'
            lines = [f'  {BOLD}MCP Prompts{RESET}', f'  {DIM}{"─" * 46}{RESET}']
            for p in prompts:
                lines.append(f'    {CYAN}{p.get("name", "?")}{RESET} '
                             f'{DIM}[{p["server"]}] {(p.get("description") or "")[:70]}{RESET}')
            lines.append('')
            lines.append(f'  {DIM}Run one with{RESET} {CYAN}/mcp-prompt <name>{RESET}')
            return '\n'.join(lines)
        rendered = mcp_manager.get_prompt(name)
        if rendered.startswith("Error:"):
            return f'  {RED}{rendered}{RESET}'
        return rendered

    if cmd == "mcp-resources":
        return handle_read_mcp_resource({})

    if cmd == "setup":
        # The same body the trigger words would load, reached deliberately.
        # Routed through /skill's mechanism rather than duplicating it, so the
        # command and the trigger cannot describe two different procedures.
        result = cmd_skill_run("onboard")
        if result is None:
            return (f'  {RED}The onboard skill is not installed.{RESET} '
                    f'{DIM}It ships with TOMAS under skills/onboard.{RESET}')
        try:
            import onboarding
            onboarding.note_offered()   # asked for it: this session is spent
        except Exception:
            pass
        messages.append({"role": "user", "content": result})
        print(f'  {YELLOW}⚡ Setup:{RESET} {BOLD}getting to know you{RESET} '
              f'{DIM}[Esc skips any question]{RESET}')
        return None  # fall through to the agent loop

    if cmd == "skills":
        return cmd_skill_list()

    if cmd == "skill":
        skill_name = parts[1] if len(parts) > 1 else ""
        if not skill_name:
            lines = [f'  {YELLOW}Usage:{RESET} {CYAN}/skill <name>{RESET}']
            lines.append(f'  {DIM}Try {RESET}{CYAN}/skills{RESET} {DIM}to see available skills.{RESET}')
            return '\n'.join(lines)
        result = cmd_skill_run(skill_name)
        if result is None:
            return f'  {RED}Skill "{skill_name}" not found. Try{RESET} {CYAN}/skills{RESET}'
        # Inject skill content into the conversation as a user message
        messages.append({"role": "user", "content": result})
        print(f'  {YELLOW}⚡ Skill call:{RESET} {BOLD}{skill_name}{RESET} {DIM}[loaded into conversation]{RESET}')
        return None  # signals main loop to fall through to agent loop

    if cmd == "pdf-report":
        try:
            from pdf_report_skill import generate_ai_news_pdf
            output = generate_ai_news_pdf()
            return f'  {GREEN}✓{RESET} PDF report generated: {BLUE}{output}{RESET}'
        except Exception as e:
            return f'  {RED}⚠{RESET} PDF report failed: {e}'

    # ── Session commands ──
    if cmd in ("session", "sessions", "save", "load"):
        # "sessions" is an alias that defaults to the "list" subcommand.
        # parts is split with maxsplit=1, so parts[1] holds "sub [arg]".
        # Re-split to extract the subcommand and any trailing argument.
        rest = parts[1] if len(parts) > 1 else ""
        rest_parts = rest.split(maxsplit=1)
        sub = rest_parts[0].lower() if rest_parts else ""
        sid_arg = rest_parts[1].strip() if len(rest_parts) > 1 else ""
        # /sessions (alias) with no subcommand → list
        if cmd == "sessions" and not sub:
            sub = "list"
        # /save and /load are listed in SLASH_COMMANDS, offered by tab
        # completion, and were handled by nothing at all: both returned None,
        # which the REPL renders as no output whatsoever. The user typed the
        # command the help screen told them to type, saw an empty line, and
        # their session was not saved. They are the short forms of subcommands
        # that already work, so they route to them rather than growing a
        # second implementation.
        elif cmd == "save":
            sub, sid_arg = "save", ""
        elif cmd == "load":
            sub, sid_arg = "continue", (rest.strip() or "")

        if sub in ("list", "ls"):
            sessions = list_sessions(limit=20)
            if not sessions:
                return f'  {DIM}No saved sessions.{RESET}'
            lines = [f'  {BOLD}Saved Sessions ({len(sessions)}){RESET}']
            lines.append(f'  {DIM}{"─" * 60}{RESET}')
            for s in sessions:
                sid = s.get("id", "?")[:22]
                ts = s.get("timestamp_str", "?")
                proj = s.get("project", "?")
                msgs = s.get("message_count", 0)
                model = s.get("model", "?")[:25]
                summary = s.get("summary", "")[:60]
                lines.append(
                    f'  {CYAN}{sid}{RESET}  {DIM}{ts}{RESET}  '
                    f'{GREEN}{proj}{RESET}  {msgs}msgs'
                )
                if summary:
                    lines.append(f'  {DIM}  {summary}{RESET}')
                lines.append('')
            return '\n'.join(lines)

        elif sub in ("save",):
            sid = save_session(messages)
            s = _session_tokens
            return (
                f'  {GREEN}✓{RESET} Session saved: {CYAN}{sid}{RESET}\n'
                f'  {DIM}  Tokens: {s["input"]:,} in · {s["output"]:,} out'
                f' ({s["calls"]} calls){RESET}'
            )

        elif sub in ("continue", "load"):
            if not sid_arg:
                # Name the form they actually typed. Being told to use
                # `/session continue` after typing `/load` reads as though
                # `/load` was the wrong command, when it is the short form.
                usage = '/load <id>' if cmd == 'load' else '/session continue <id>'
                return (f'  {YELLOW}Usage:{RESET} {CYAN}{usage}{RESET}\n'
                        f'  {DIM}Use{RESET} {CYAN}/session list{RESET} '
                        f'{DIM}to see session IDs.{RESET}')
            sid = sid_arg
            loaded = continue_session(sid)
            if loaded is None:
                return f'  {RED}✗{RESET} Session not found: {sid}'
            # Replace current messages with the loaded ones
            messages.clear()
            reset_session_state()
            messages.extend(loaded)
            # Show the full conversation history so the user has context
            lines = [
                f'  {GREEN}✓{RESET} Continuing session {CYAN}{sid}{RESET} '
                f'({len(loaded)} messages loaded).',
                f'  {DIM}{"─" * 50}{RESET}',
            ]
            for m in loaded:
                role = m.get("role", "?")
                content = m.get("content", "")
                text = _format_block_content(content)
                if role == "user":
                    icon = f'{GREEN}◆{RESET}'
                    label = f'{GREEN}{BOLD}You{RESET}'
                elif role == "assistant":
                    icon = f'{MAGENTA}▌{RESET}'
                    label = f'{MAGENTA}{BOLD}TOMAS{RESET}'
                else:
                    icon = f'{DIM}·{RESET}'
                    label = f'{DIM}{role}{RESET}'
                display = text.replace('\n', '\n      ')
                if len(display) > 300:
                    display = display[:300] + f' {DIM}…{RESET}'
                lines.append(f'  {icon} {label}')
                lines.append(f'      {display}')
            lines.append(f'  {DIM}{"─" * 50}{RESET}')
            lines.append(f'  {DIM}Type your next message to continue the conversation.{RESET}')
            return '\n'.join(lines)

        elif sub in ("delete", "rm"):
            if not sid_arg:
                return f'  {YELLOW}Usage:{RESET} {CYAN}/session delete <id>{RESET}'
            sid = sid_arg
            if delete_session(sid):
                return f'  {GREEN}✓{RESET} Session deleted: {sid}'
            else:
                return f'  {RED}✗{RESET} Session not found: {sid}'

        elif sub in ("latest",):
            latest = get_latest_session()
            if latest is None:
                return f'  {DIM}No saved sessions.{RESET}'
            sid = latest.get("id", "?")
            ts = latest.get("timestamp_str", "?")
            proj = latest.get("project", "?")
            msgs = latest.get("message_count", 0)
            summary = latest.get("summary", "")[:80]
            lines = [
                f'  {BOLD}Latest Session{RESET}',
                f'  {DIM}{"─" * 50}{RESET}',
                f'  ID:      {CYAN}{sid}{RESET}',
                f'  When:    {ts}',
                f'  Project: {GREEN}{proj}{RESET}',
                f'  Messages: {msgs}',
            ]
            if summary:
                lines.append(f'  Summary: {DIM}{summary}{RESET}')
            lines.append('')
            lines.append(f'  {YELLOW}Tip:{RESET} Use {CYAN}/session continue {sid}{RESET} to pick up where you left off.')
            return '\n'.join(lines)

        else:
            lines = [
                f'  {BOLD}Session Management{RESET}',
                f'  {DIM}{"─" * 50}{RESET}',
                f'  {CYAN}/session list{RESET}       — List saved sessions',
                f'  {CYAN}/session save{RESET}       — Save current conversation',
                f'  {CYAN}/session continue <id>{RESET}  — Load and continue a session',
                f'  {CYAN}/session latest{RESET}     — Show latest session info',
                f'  {CYAN}/session delete <id>{RESET}  — Delete a session',
            ]
            return '\n'.join(lines)

    # ── Self-improvement commands ──
    if cmd in ("self-improve", "si"):
        sub = parts[1].lower() if len(parts) > 1 else ""

        # ── The learning store (Phase 3) ──
        # The requirement is that the user does not see the machinery. That
        # must not mean the user cannot see it: an agent that silently
        # accumulates wrong beliefs with no way to inspect them is worse than
        # one that does not learn.
        if sub in ("", "facts", "learned"):
            return _render_learned_facts()

        if sub in ("reflect", "reflection"):
            return _render_reflection_log()

        if sub in ("analyze", "scan"):
            self_improve.update_session_analysis()
            return f'  {GREEN}✓{RESET} Session analysis updated.'

        # The keyword-counting generator behind `skills`, `tips`, and
        # `patterns` was deleted in Phase 6. Point at what replaced it rather
        # than leaving three subcommands that silently show nothing.
        if sub in ("skills", "tips", "patterns"):
            return (
                f'  {DIM}The pattern/tip generator was removed — it counted keywords and\n'
                f'  filled in templates, and nothing consumed its output.{RESET}\n'
                f'  {CYAN}/self-improve facts{RESET}      — what the agent has actually learned\n'
                f'  {CYAN}/self-improve reflect{RESET}    — the reflection log\n'
                f'  {CYAN}/skills{RESET}                  — installed skills'
            )

        # Default: show status
        return self_improve.get_self_improve_status()

    if cmd == "forget":
        fact_id = parts[1].strip() if len(parts) > 1 else ""
        if not fact_id:
            return (f'  {YELLOW}Usage:{RESET} {CYAN}/forget <id>{RESET}\n'
                    f'  {DIM}Ids are shown by /si.{RESET}')
        try:
            removed = learning.forget(fact_id)
        except Exception as e:
            return f'  {RED}✗{RESET} Could not forget: {e}'
        if not removed:
            return f'  {DIM}No fact with id {fact_id}.{RESET}'
        return (f'  {GREEN}✓{RESET} Forgotten: {removed.get("fact", "")[:80]}\n'
                f'  {DIM}Tombstoned — reflection will not re-learn it.{RESET}')

    if cmd == "private":
        learning.set_enabled(not learning.is_enabled())
        if learning.is_enabled():
            return f'  {GREEN}✓{RESET} Learning re-enabled for this session.'
        return (f'  {GREEN}✓{RESET} Incognito: nothing from this session will be '
                f'reflected on or stored.')

    # ── Standing rules ──
    if cmd == "rules":
        return _handle_rules((parts[1] if len(parts) > 1 else "").strip())

    # (see _handle_rules, defined above)

    # ── Self-note commands ──
    if cmd == "note":
        # parts is split with maxsplit=1: parts[1] = "title content..."
        rest = parts[1] if len(parts) > 1 else ""
        rest_parts = rest.split(maxsplit=1)
        if len(rest_parts) < 2:
            return (
                f'  {YELLOW}Usage:{RESET} {CYAN}/note <title> <content>{RESET}\n'
                f'  {DIM}Create a self-note about something the agent learned.{RESET}'
            )
        title = rest_parts[0]
        content = rest_parts[1]
        note_id = self_notes.create_note(
            title=title,
            content=content,
            note_type="insight",
        )
        return f'  {GREEN}✓{RESET} Note created: {CYAN}{note_id}{RESET}'

    if cmd == "notes":
        notes = self_notes.list_notes(limit=20)
        if not notes:
            return f'  {DIM}No self-notes yet. Use /note <title> <content> to create one.{RESET}'
        lines = [f'  {BOLD}Self-Notes ({len(notes)}){RESET}']
        lines.append(f'  {DIM}{"─" * 50}{RESET}')
        for n in notes:
            nid = n.get("id", "?")[:22]
            title = n.get("title", "?")[:50]
            ntype = n.get("type", "insight")
            tags = n.get("tags", [])
            tag_str = f' [{", ".join(tags)}]' if tags else ""
            lines.append(f'  {DIM}[{ntype}]{RESET} {title} {DIM}{nid}{RESET}{tag_str}')
        return "\n".join(lines)

    if cmd == "exit":
        return "__exit__"

    return None


# ---------------------------------------------------------------------------
# Real-time input reader with slash suggestions
# ---------------------------------------------------------------------------

# Input history — persists across calls to read_input_with_suggestions
_input_history: list[str] = []
_history_index: int = 0


def _read_input_cross_platform(prompt: str) -> str:
    """Cross-platform input fallback for Linux/macOS.

    Uses standard input() with optional readline support for history.
    Provides basic slash-command tab completion via readline if available.
    """
    global _input_history

    # Try to enable readline for arrow-key history + tab completion
    try:
        import readline

        # Set up history file
        hist_file = Path.home() / ".tomas" / "input_history.txt"
        hist_file.parent.mkdir(parents=True, exist_ok=True)
        try:
            readline.read_history_file(str(hist_file))
        except (OSError, FileNotFoundError):
            pass

        # Tab completion for slash commands
        def _slash_completer(text: str, state: int):
            if text.startswith('/'):
                matches = [f"/{n}" for n in SLASH_COMMANDS if n.startswith(text[1:])]
                return matches[state] if state < len(matches) else None
            return None

        readline.set_completer(_slash_completer)
        readline.parse_and_bind("tab: complete")
        readline.set_history_length(100)
    except ImportError:
        hist_file = None

    try:
        result = input(prompt).strip()
    except EOFError:
        return "exit"

    # Save to history
    if result and not result.startswith('/'):
        if not _input_history or _input_history[-1] != result:
            _input_history.append(result)
            if len(_input_history) > 100:
                _input_history.pop(0)

    # Save readline history
    if hist_file:
        try:
            import readline
            readline.write_history_file(str(hist_file))
        except (OSError, ImportError):
            pass

    return result


# Past this many characters a paste is shown as a marker instead of its text.
# Small enough that a stack trace or a config block collapses, large enough
# that a sentence or a path pasted mid-thought still reads as what it is.
PASTE_COLLAPSE_CHARS = 400

#: How a paste burst is followed to its end. Two polls of 12ms — 24ms of
#: silence — used to be enough to call it over, and that is thin: a terminal
#: hands a long paste over in chunks, and the gap between them can be longer
#: than that over SSH, in a remote session, or simply behind a redraw. A
#: premature end splits the text at whatever character the queue happened to
#: run dry on.
#:
#: Raised after a prompt arrived split mid-word ("щоб во" / "на була") across
#: six fragments. That is what a premature end looks like, but it is not
#: proof: the pasted source may have carried those breaks itself, and nothing
#: in the session record can tell the two apart. Six polls costs 72ms, only
#: ever runs *after* a burst has been detected, and is never in the way of
#: typing — cheap enough to widen on a strong suspicion.
_PASTE_POLL_S = 0.012
_PASTE_IDLE_POLLS = 6


IMAGE_SUFFIXES = {".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
                  ".gif": "image/gif", ".webp": "image/webp"}

# Anthropic's per-image ceiling. Bigger than this is rejected outright, so
# it is better to say so than to send it and read the 400 back.
MAX_IMAGE_BYTES = 5 * 1024 * 1024


def grab_clipboard_image(dest: Path) -> bool:
    """Save the clipboard's image to `dest` as PNG. False if there isn't one.

    Clipboard image data never reaches the console's keystroke stream — a
    paste of an image delivers no characters at all — so it cannot be picked
    up by reading input, however that reading is done. .NET's clipboard is
    already on every Windows box, which beats taking a new imaging dependency
    for one call. It needs a single-threaded apartment, hence -STA.
    """
    if sys.platform != "win32":
        return False
    script = (
        "Add-Type -AssemblyName System.Windows.Forms,System.Drawing;"
        "$i=[System.Windows.Forms.Clipboard]::GetImage();"
        "if($i -ne $null){$i.Save('%s',"
        "[System.Drawing.Imaging.ImageFormat]::Png);'ok'}else{'none'}"
    ) % str(dest).replace("'", "''")
    try:
        out = subprocess.run(
            ["powershell", "-NoProfile", "-STA", "-Command", script],
            capture_output=True, timeout=20, text=True)
        return "ok" in (out.stdout or "") and dest.exists()
    except Exception:
        return False


def put_clipboard_text(text: str) -> bool:
    """Put `text` on the clipboard. False if it could not be done.

    `clip.exe` is on every Windows box and takes its input on stdin, so there
    is nothing to quote or escape — which matters here, because the thing being
    copied is arbitrary user text full of quotes, newlines and Cyrillic. It
    writes whatever bytes it is given, so the encoding is stated outright
    rather than left to the console code page.
    """
    if sys.platform != "win32" or not text:
        return False
    try:
        proc = subprocess.run(["clip"], input=text.encode("utf-16-le"),
                              capture_output=True, timeout=10)
        return proc.returncode == 0
    except Exception:
        return False


def _image_paths_in(text: str) -> list[Path]:
    """Existing image files named anywhere in the message.

    Dragging a file onto a terminal types its path, and the clipboard grab
    below inserts one, so path-spotting covers both without a second
    mechanism. Quotes are stripped because that is how a dragged path with a
    space arrives.
    """
    found: list[Path] = []
    pattern = r'"[^"]+"|\'[^\']+\'|\S+'
    for token in re.findall(pattern, text or ""):
        candidate = token.strip('"\'')
        if Path(candidate).suffix.lower() not in IMAGE_SUFFIXES:
            continue
        try:
            path = _resolve(candidate)
            if path.is_file() and path not in found:
                found.append(path)
        except Exception:
            continue
    return found


def build_user_content(text: str):
    """The user's turn, with any images they named attached.

    Returns the plain string when there is nothing to attach, so the ordinary
    path stays exactly as it was. Images are dropped with a printed reason
    rather than silently, and never sent to a provider whose probe says it
    cannot read them — that is a 400, not a degraded answer.
    """
    paths = _image_paths_in(text)
    if not paths:
        return text
    if not getattr(_active_capabilities(), "vision", False):
        print(f'  {YELLOW}⚠{RESET}  {_get_model()} cannot read images — '
              f'sending the message as text only.')
        return text
    blocks: list[dict] = []
    for path in paths:
        try:
            data = path.read_bytes()
            if len(data) > MAX_IMAGE_BYTES:
                print(f'  {YELLOW}⚠{RESET}  {path.name} is '
                      f'{len(data) // 1024} KB, over the {MAX_IMAGE_BYTES // 1024} KB '
                      f'limit — not attached.')
                continue
            blocks.append({
                "type": "image",
                "source": {"type": "base64",
                           "media_type": IMAGE_SUFFIXES[path.suffix.lower()],
                           "data": base64.b64encode(data).decode("ascii")},
            })
            print(f'  {DIM}🖼  attached {path.name} ({len(data) // 1024} KB){RESET}')
        except Exception as e:
            print(f'  {YELLOW}⚠{RESET}  could not attach {path.name}: {e}')
    if not blocks:
        return text
    return blocks + [{"type": "text", "text": text}]


def _shift_down() -> bool:
    """Is Shift physically held right now?

    `msvcrt.getwch` hands back a decoded character and nothing about
    modifiers, so Shift+Enter and Enter are the same `\\r` to it. Reading the
    key state settles which one it was. `GetAsyncKeyState` reports the
    physical keyboard rather than the message queue, which matters because a
    console app pumps no messages — `GetKeyState` would answer from a queue
    that never advances. Cheaper and far less invasive than replacing the
    read primitive with `ReadConsoleInput`, and it fails to False anywhere
    that is not Windows, where the caller already falls back.
    """
    try:
        import ctypes
        VK_SHIFT = 0x10
        return bool(ctypes.windll.user32.GetAsyncKeyState(VK_SHIFT) & 0x8000)
    except Exception:
        return False


def read_input_with_suggestions(prompt: str) -> str:
    """Read a line of input with live slash-command suggestions beneath the prompt.

    Uses msvcrt (Windows) for character-by-character input.  When the line
    starts with ``/``, matching ``SLASH_COMMANDS`` are shown on the line
    directly below the prompt and updated after each keystroke.

    **Interactive features**:

    * ``↑`` / ``↓``  — navigate input history (when not in command mode)
    * ``↑`` / ``↓``  — navigate through suggestion items (in command mode)
    * ``Tab``        — auto-complete to the current match or common prefix
    * ``Enter``      — accept the highlighted suggestion (if any), or submit as-is
    * ``Shift+Enter``— insert a newline instead of submitting
    * ``Esc``        — clear the input line

    The buffer is drawn across as many terminal rows as it needs, so a long
    or multi-line prompt stays readable — and selectable — instead of
    scrolling out of view. Pasted text is taken in whole rather than replayed
    as keystrokes, so line breaks inside it do not submit; anything past
    ``PASTE_COLLAPSE_CHARS`` is shown as a marker and restored on send.

    Falls back to a cross-platform input with basic slash-command completion
    when msvcrt is unavailable (Linux/macOS).
    """
    try:
        import msvcrt  # Windows-only
        import sys
    except ImportError:
        return _read_input_cross_platform(prompt)

    global _history_index   # mode switching goes through set_mode()
    base_prompt = prompt
    from version import VERSION as _VERSION

    buffer: list[str] = []
    #: Where the next character goes, as an index into `buffer`.
    #:
    #: Until this existed the buffer was append-only: every key landed at
    #: the end and backspace was the only way back, so correcting the
    #: start of a long prompt meant deleting all of it. Left/Right, Home,
    #: End and Delete all move or act on this one number, and `_place`
    #: turns it into a screen position after each redraw.
    pos = 0
    showing = False
    selected: int | None = None  # index of the currently highlighted suggestion
    drawn_rows = 1               # rows the input block occupied on the last draw
    pastes: dict[str, str] = {}  # placeholder -> the text it stands in for
    escape_armed = False         # an Esc on an empty line; a second one exits
    # Undo history for Ctrl+Z. Snapshots are taken before a *destructive* edit
    # only — Ctrl+U, Ctrl+W, a paste, a history recall. Typing is not
    # snapshotted per character: undoing one letter at a time is useless, and
    # the keystroke that actually needs taking back is the one that wiped a
    # long prompt in a single press.
    undo_stack: list[list[str]] = []

    # Reset history navigation to past-the-end
    _history_index = len(_input_history)

    # ── mode helpers ─────────────────────────────────────────────────┬─
    def _mode_badge() -> str:
        mode = current_mode_name()
        if mode == "bypass":
            # Distinct from YOLO on sight: this one also removes the stop that
            # would otherwise interrupt a long unattended run.
            return f'[{RED}{BOLD}BYPASS ⇥{RESET}]'
        if mode == "yolo":
            return f'[{RED}YOLO{RESET}]'
        color = GREEN if mode == "auto" else YELLOW
        return f'[{color}{mode}{RESET}]'

    def _build_prompt() -> str:
        badge = _mode_badge()
        idx = base_prompt.find('TOMAS')
        if idx < 0:
            return badge + ' ' + base_prompt
        # Terse on purpose: this line redraws on every keystroke, so it is
        # the one place version visibility has to earn its width. The full
        # picture (last-updated date, git info if this is a dev checkout)
        # lives in /version instead.
        after = idx + len('TOMAS')
        return (base_prompt[:idx] + badge + ' ' + base_prompt[idx:after]
                + f' {DIM}v{_VERSION}{RESET}' + base_prompt[after:])

    # Delegates to the module-level `set_mode` rather than repeating the flag
    # writes. This copy had drifted from the /mode handler's — it knew nothing
    # about read_mcp_resource, and a mode added to one was invisible to the
    # other, which is exactly how bypass would have ended up half-wired.
    _set_mode = set_mode

    # ── helpers ──────────────────────────────────────────────────────────

    def _repr() -> str:
        return ''.join(buffer)

    def _clamp() -> None:
        """Keep the insertion point inside the buffer.

        Every mutation goes through here rather than each site doing its
        own arithmetic, because the ways to get it wrong are silent: a
        `pos` past the end inserts nowhere, and a negative one inserts at
        the end.
        """
        nonlocal pos
        pos = max(0, min(pos, len(buffer)))

    def _checkpoint() -> None:
        """Remember the line before something destructive happens to it."""
        if undo_stack and undo_stack[-1] == buffer:
            return                       # nothing changed since the last one
        undo_stack.append(list(buffer))
        if len(undo_stack) > 50:
            undo_stack.pop(0)

    def _flash(message: str) -> None:
        """One-line feedback under the prompt, gone on the next redraw."""
        sys.stdout.write('\033[1B\r\033[2K'
                         f'  {DIM}{message}{RESET}'
                         '\033[1A\r')
        sys.stdout.flush()

    def _is_slash() -> bool:
        return _repr().startswith('/')

    def _cmd_filter() -> str:
        """Return the text after the leading ``/``, stripped + lowercased."""
        return _repr()[1:].strip().lower()

    def _get_matches() -> list[str]:
        cf = _cmd_filter()
        return sorted(
            name for name in SLASH_COMMANDS
            if not cf or name.startswith(cf)
        )

    def _rows() -> list[str]:
        """The physical rows prompt + buffer will occupy, as drawn.

        One column is left unused so a full row never makes the terminal
        soft-wrap on its own: every line break here is one we emitted, which
        is what makes the row count below exact and the cursor arithmetic
        safe.
        """
        from text_display import hard_wrap, term_columns
        return hard_wrap(_build_prompt() + _repr(), max(1, term_columns() - 1))

    def _place() -> None:
        """Put the terminal cursor at `pos`, from the end of the block.

        Exact rather than approximate, because `hard_wrap` breaks on
        width and not on words: wrapping the text *up to* `pos` lays out
        identically to the first part of wrapping all of it, so the last
        row of that prefix is the row the cursor belongs on and its
        display width is the column. Word wrapping would make this a
        guess, and a cursor one column out is worse than no cursor.
        """
        if pos >= len(buffer):
            return                      # the redraw already left it here
        from text_display import hard_wrap, term_columns, display_width
        head = hard_wrap(_build_prompt() + ''.join(buffer[:pos]),
                         max(1, term_columns() - 1))
        up = (drawn_rows - 1) - (len(head) - 1)
        move = (f'\033[{up}A' if up > 0 else '') + '\r'
        column = display_width(head[-1])
        if column:
            move += f'\033[{column}C'
        sys.stdout.write(move)
        sys.stdout.flush()

    def _refresh(place: bool = True):
        """Redraw prompt + buffer across as many rows as it needs.

        This used to be one row: `\\r\\033[K` clears only the row the cursor is
        on, so anything longer was scrolled horizontally behind a `…` and the
        rest was simply not on screen — you could not read back what you had
        typed, and terminal selection could not copy what it could not show.
        A buffer now carries real newlines (Shift+Enter, pasted text) as well,
        which one row cannot represent at all.

        Rewinding to the top of the block and clearing to the end of the
        screen (`\\033[J`) also wipes the suggestion row below, so the caller
        redraws it after; clearing per-row would leave that one behind.
        """
        nonlocal drawn_rows
        # One place where `pos` is guaranteed sane, rather than trusting
        # sixteen mutation sites to agree.
        _clamp()
        if drawn_rows > 1:
            sys.stdout.write(f'\033[{drawn_rows - 1}A')
        sys.stdout.write('\r\033[J')
        rows = _rows()
        sys.stdout.write('\n'.join(rows))
        drawn_rows = len(rows)
        sys.stdout.flush()
        # `_show` draws the suggestion row *below* this block and needs
        # the cursor left at the end to get there, so it places the
        # cursor itself once it is done.
        if place:
            _place()

    def _show():
        nonlocal showing, selected
        matches = _get_matches()

        # Mode status string (always shown in suggestion line)
        mode = current_mode_name()
        mode_str = (f'  {DIM}mode:{RESET} {mode_color(mode)}{mode}{RESET}  '
                    f'{DIM}⇧+Space:toggle  {MODE_KEYS_HINT}{RESET}')

        if matches:
            parts = []
            for i, name in enumerate(matches):
                if selected is not None and i == selected:
                    parts.append(f'{BOLD}{CYAN}/{name}{RESET}')
                else:
                    parts.append(f'{CYAN}/{name}{RESET}')
            hint = f'  {DIM}↑↓·Tab·Enter{RESET}'
            text = '  ' + '  '.join(parts) + hint + mode_str
        else:
            # No matches — show every command dimmed + mode info
            parts = [f'{DIM}/{n}{RESET}' for n in sorted(SLASH_COMMANDS)]
            text = '  ' + '  '.join(parts) + mode_str
            selected = None

        # The hint has to fit on exactly one row. With no matches it lists
        # every slash command, far wider than any terminal; left untruncated
        # it soft-wrapped onto a second row while the rewind below moved only
        # one, so the cursor came back a row low and the next redraw drew
        # *under* the previous one.
        #
        # The block is drawn first because `_refresh` clears to the end of the
        # screen — writing the hint before it would erase the hint. Coming
        # back, the last row is re-emitted rather than just moving up: that
        # leaves the cursor after the final character, where typing continues,
        # instead of at column 0.
        from text_display import shorten, term_columns
        _refresh(place=False)
        last_row = _rows()[-1]
        sys.stdout.write('\033[1B\r\033[2K')
        sys.stdout.write(shorten(text, max(1, term_columns() - 1)))
        sys.stdout.write('\033[1A\r' + last_row)
        sys.stdout.flush()
        _place()
        showing = True

    def _hide():
        nonlocal showing, selected
        if showing:
            showing = False
            selected = None
            _refresh()   # clearing to end of screen takes the hint row with it

    # ── paste ────────────────────────────────────────────────────────────
    #
    # A paste is not a distinct event here: the terminal replays it as
    # ordinary keystrokes, already queued. `kbhit()` is what separates it
    # from typing — a person cannot get the next character into the queue
    # before this one has been handled, so anything already waiting arrived
    # together. That is the whole detector, and it needs no bracketed-paste
    # support from the terminal.

    def _drain() -> str:
        """Take everything already queued, as one string.

        Newlines are normalised rather than acted on: inside a paste a
        Return is a line break, not "send". The short idle re-checks cover a
        paste arriving faster than this loop empties it.
        """
        chunk: list[str] = []
        prev_cr = False
        idle = 0
        while idle < _PASTE_IDLE_POLLS:
            if not msvcrt.kbhit():
                time.sleep(_PASTE_POLL_S)
                idle += 1
                continue
            idle = 0
            c = msvcrt.getwch()
            if c in ('\xe0', '\x00'):     # extended key: consume its second half
                msvcrt.getwch()
                prev_cr = False
            elif c == '\r':
                chunk.append('\n')
                prev_cr = True
            elif c == '\n':
                if not prev_cr:            # CRLF is one break, LF alone is one
                    chunk.append('\n')
                prev_cr = False
            elif c == '\t':
                chunk.append('    ')
                prev_cr = False
            elif c.isprintable():
                chunk.append(c)
                prev_cr = False
            else:
                prev_cr = False
        return ''.join(chunk)

    def _insert(text: str) -> None:
        nonlocal pos
        """Add pasted text, standing in a marker for anything long.

        A pasted file is worth sending and not worth *looking* at while you
        finish the sentence around it, so past a threshold the buffer shows
        its size instead of its content. `_expand` puts the real text back at
        send time, so the model still receives all of it.
        """
        if not text:
            return
        _checkpoint()   # a paste is one keystroke's worth of damage to undo
        if len(text) >= PASTE_COLLAPSE_CHARS:
            lines = text.count('\n') + 1
            marker = (f'[#{len(pastes) + 1} pasted {len(text)} chars'
                      f'{f", {lines} lines" if lines > 1 else ""}]')
            pastes[marker] = text
            buffer[pos:pos] = list(marker)
            pos += len(marker)
        else:
            buffer[pos:pos] = list(text)
            pos += len(text)

    def _absorb_burst(seed: str = '') -> None:
        _insert(seed + _drain())

    def _expand(text: str) -> str:
        for marker, full in pastes.items():
            text = text.replace(marker, full)
        return text

    # ── main input loop ──────────────────────────────────────────────────

    _refresh()   # draws the prompt and establishes the row count

    try:
        while True:
            # getwch returns a decoded str, so a Ukrainian, Russian,
            # accented-Latin or CJK keystroke arrives as one character
            # instead of a high byte that the old ASCII test discarded.
            ch = msvcrt.getwch()

            # Any other key disarms the exit: Esc, second thoughts, typing,
            # Esc again should not quit on the strength of the first one.
            if ch != '\x1b':
                escape_armed = False

            # ── Enter ──────────────────────────────────────────────────────
            if ch == '\r':
                # Three ways a Return can arrive, and only one of them means
                # "send". Shift+Enter is a deliberate newline. A Return in the
                # middle of a burst is a line break inside pasted text — that
                # one used to submit, so pasting three paragraphs sent three
                # half-written messages.
                #
                # A Return that lands at the *head* of a paste's next chunk,
                # with the queue momentarily empty behind it, reads as "send"
                # here and cannot be told apart from one. A time-based grace
                # window was tried and removed: it takes "paste, then press
                # Enter" — which is how the prompt is documented and tested to
                # work — away from the user to buy a case that `_drain`'s idle
                # window already covers, since that now follows a burst
                # through gaps of up to _PASTE_POLL_S * _PASTE_IDLE_POLLS.
                if _shift_down() or msvcrt.kbhit():
                    if msvcrt.kbhit():
                        _absorb_burst()
                    else:
                        buffer.insert(pos, '\n')
                        pos += 1
                    selected = None
                    _refresh()
                    _hide()
                    continue
                chosen = selected  # save before _hide() clears it
                _hide()
                sys.stdout.write('\n')
                sys.stdout.flush()
                # If a suggestion was highlighted, accept it
                if chosen is not None and _is_slash():
                    matches = _get_matches()
                    if chosen < len(matches):
                        buffer[:] = list('/' + matches[chosen])
                result = _expand(_repr())
                # Save non-empty, non-command input to history (max 100)
                if result.strip() and not result.startswith('/'):
                    if not _input_history or _input_history[-1] != result:
                        _input_history.append(result)
                        if len(_input_history) > 100:
                            _input_history.pop(0)
                return result

            # ── Ctrl+C — copy, not quit ───────────────────────────────────
            # The terminal handles Ctrl+C as copy whenever a selection
            # exists, and only forwards \x03 here when there is none. Ending
            # the session on it made the ordinary copy reflex a coin flip:
            # miss the selection and the session was over. Esc Esc quits
            # instead, and clearing the line is the useful reading of a
            # keystroke the user meant as "copy".
            if ch == '\x03':
                if buffer:
                    _checkpoint()
                    buffer.clear()
                    pos = 0
                    selected = None
                    _hide()
                    _refresh()
                continue

            # ── Arrow / function keys ─────────────────────────────────────
            if ch in ('\xe0', '\x00'):
                # One prefix branch for both extended-key families: getwch
                # reports '\xe0' for arrows and '\x00' for function keys, and
                # the second call gives the key itself.
                ext = msvcrt.getwch()

                # Ctrl+Alt+X / Alt+X — the debug view.
                #
                # An Alt'd letter arrives through the same extended-key door as
                # the function keys: a '\x00' prefix, then the key's *scan*
                # code rather than its character. 0x2D is X. Windows delivers
                # Ctrl+Alt+X here identically to Alt+X (the console API's
                # modifier bits do not survive `getwch`), so binding the scan
                # code catches both, and Alt+X alone is a harmless second way
                # in rather than a conflict — nothing else in this REPL uses it.
                if ext == '\x2d':
                    _hide()
                    print()
                    # Opens the live window, arming recording first if needed.
                    # Deliberately not a dump into this console: the point of
                    # the key is to watch traffic *while* using the session,
                    # and printing 200 lines of JSON over the prompt is the
                    # opposite of that.
                    print(_handle_debug("window"))
                    print()
                    _refresh()
                    continue

                # F5–F9 — quick mode switch
                if ext in ('\x3f', '\x40', '\x41', '\x42', '\x43'):
                    if ext == '\x3f':      # F5 — toggle auto/default
                        _set_mode("default" if AUTO_APPROVE_LOW else "auto")
                    elif ext == '\x40':    # F6 — default mode
                        _set_mode("default")
                    elif ext == '\x41':    # F7 — strict mode
                        _set_mode("strict")
                    elif ext == '\x42':    # F8 — YOLO mode
                        _set_mode("yolo")
                    else:                  # F9 — bypass mode
                        _set_mode("bypass")
                    _refresh()
                    if _is_slash():
                        _show()
                    else:
                        _hide()
                    continue

                if showing and _is_slash():
                    matches = _get_matches()
                    if not matches:
                        continue
                    if ext == 'H':  # ↑ Up arrow — navigate suggestions
                        selected = 0 if selected is None else max(0, selected - 1)
                        _show()
                    elif ext == 'P':  # ↓ Down arrow — navigate suggestions
                        selected = 0 if selected is None else min(len(matches) - 1, selected + 1)
                        _show()
                elif ext == 'H':  # ↑ Up arrow — history recall
                    if _input_history:
                        _hide()
                        _checkpoint()   # recall replaces whatever was typed
                        if _history_index > 0:
                            _history_index -= 1
                        buffer[:] = list(_input_history[_history_index])
                        pos = len(buffer)
                        _refresh()
                elif ext == 'P':  # ↓ Down arrow — history forward
                    _hide()
                    _checkpoint()
                    if _history_index < len(_input_history) - 1:
                        _history_index += 1
                        buffer[:] = list(_input_history[_history_index])
                    else:
                        _history_index = len(_input_history)
                        buffer.clear()
                    pos = len(buffer)
                    _refresh()

                # ── Cursor movement ────────────────────────────
                # These were 'silently consumed' -- the buffer was
                # append-only, so the only way to fix the start of a long
                # prompt was to delete everything after it.
                elif ext == 'K':        # ←
                    pos = max(0, pos - 1)
                    _refresh()
                elif ext == 'M':        # →
                    pos = min(len(buffer), pos + 1)
                    _refresh()
                elif ext == 'G':        # Home
                    pos = 0
                    _refresh()
                elif ext == 'O':        # End
                    pos = len(buffer)
                    _refresh()
                elif ext == 's':        # Ctrl+← -- word left
                    while pos > 0 and buffer[pos - 1].isspace():
                        pos -= 1
                    while pos > 0 and not buffer[pos - 1].isspace():
                        pos -= 1
                    _refresh()
                elif ext == 't':        # Ctrl+→ -- word right
                    while pos < len(buffer) and not buffer[pos].isspace():
                        pos += 1
                    while pos < len(buffer) and buffer[pos].isspace():
                        pos += 1
                    _refresh()
                elif ext == 'S':        # Delete -- forward, not back
                    if pos < len(buffer):
                        _checkpoint()
                        buffer.pop(pos)
                        selected = None
                        _refresh()
                        _show() if _is_slash() else _hide()
                continue

            # ── Tab — auto-complete slash commands OR cycle risk mode ──────
            if ch == '\t':
                if _is_slash():
                    cf = _cmd_filter()
                    if cf:
                        matches = sorted(
                            name for name in SLASH_COMMANDS
                            if name.startswith(cf)
                        )
                        if len(matches) == 1:
                            # Unambiguous → replace buffer with full command
                            buffer[:] = list('/' + matches[0])
                            pos = len(buffer)
                            selected = None
                        elif len(matches) > 1:
                            # Several matches → extend to the longest common prefix
                            prefix = os.path.commonprefix(matches)
                            if prefix and prefix != cf:
                                buffer[:] = list('/' + prefix)
                                pos = len(buffer)
                                selected = None
                    _refresh()
                    _show()
                else:
                    # Tab outside a slash command cycles the modes, in the one
                    # order MODE_CYCLE defines — indexing it rather than
                    # branching on the flags means a new mode joins the cycle
                    # by being added to the tuple, not by editing this if/else.
                    try:
                        nxt = (MODE_CYCLE.index(current_mode_name()) + 1) % len(MODE_CYCLE)
                    except ValueError:
                        nxt = 0          # e.g. strict, which is not in the cycle
                    _set_mode(MODE_CYCLE[nxt])
                    _refresh()
                    _hide()
                continue

            # ── Backspace ─────────────────────────────────────────────────
            if ch in ('\x08', '\x7f'):
                if buffer and pos > 0:
                    pos -= 1
                    buffer.pop(pos)
                    selected = None
                    _refresh()
                    if _is_slash():
                        _show()
                    else:
                        _hide()
                continue

            # ── Ctrl+W — delete the word before the cursor ────────────────
            # This buffer has no intra-line cursor, so backspace is the only
            # way to correct a long prompt one character at a time. Ctrl+W and
            # Ctrl+U are the two readline keys that pay off without one.
            if ch == '\x17':
                _checkpoint()
                while pos > 0 and buffer[pos - 1].isspace():
                    pos -= 1
                    buffer.pop(pos)
                while pos > 0 and not buffer[pos - 1].isspace():
                    pos -= 1
                    buffer.pop(pos)
                selected = None
                _refresh()
                _show() if _is_slash() else _hide()
                continue

            # ── Ctrl+U — clear the input line ─────────────────────────────
            if ch == '\x15':
                if buffer:
                    _checkpoint()
                    buffer.clear()
                    pos = 0
                    selected = None
                    _hide()
                    _refresh()
                continue

            # ── Ctrl+Z — undo the last destructive edit ───────────────────
            # Ctrl+U and Ctrl+W each destroy work in one keystroke and there
            # was no way back: a long prompt cleared by a mistyped Ctrl+U had
            # to be typed again from scratch. This restores the line as it was
            # before that press, repeatedly, back to the start of the turn.
            if ch == '\x1a':
                if undo_stack:
                    buffer[:] = undo_stack.pop()
                    pos = len(buffer)
                    selected = None
                    _hide()
                    _refresh()
                    if _is_slash():
                        _show()
                else:
                    _flash('nothing to undo')
                continue

            # ── Ctrl+Y — copy the current line to the clipboard ───────────
            # Ctrl+C cannot do this: the terminal intercepts it whenever a
            # selection exists and only forwards it here when there is none,
            # so binding copy to it would work about half the time. Paired
            # with ↑, this copies *any* earlier prompt — recall it, press
            # Ctrl+Y — which beats reselecting it with the mouse across a
            # block that may have scrolled.
            if ch == '\x19':
                text = _expand(_repr())
                if not text.strip():
                    _flash('nothing to copy')
                elif put_clipboard_text(text):
                    lines = text.count('\n') + 1
                    _flash(f'copied {len(text)} chars'
                           f'{f" · {lines} lines" if lines > 1 else ""} to clipboard')
                else:
                    _flash('could not reach the clipboard')
                continue

            # ── Ctrl+G — attach the image on the clipboard ─────────────────
            # Ctrl+V cannot carry this: the terminal owns that key and, with
            # an image on the clipboard, delivers no characters at all. The
            # grab writes a file and types its path, so the image travels the
            # same route as one dragged onto the window.
            if ch == '\x07':
                shot = Path(tempfile.gettempdir()) / f'tomas-clip-{int(time.time())}.png'
                if grab_clipboard_image(shot):
                    if buffer and not _repr().endswith(' '):
                        buffer.append(' ')
                    buffer.extend(str(shot))
                    pos = len(buffer)
                else:
                    _hide()
                    sys.stdout.write('\033[1B\r\033[2K'
                                     f'  {DIM}no image on the clipboard{RESET}'
                                     '\033[1A\r')
                selected = None
                _refresh()
                _hide()
                continue

            # ── Ctrl+L — clear the screen, keep what was typed ─────────────
            # Standard everywhere, and the honest way out of a terminal that
            # some other program has left dirty. `\033[3J` drops the scrollback
            # too, so the chat restarts from a genuinely clean screen.
            if ch == '\x0c':
                _hide()
                sys.stdout.write('\033[H\033[2J\033[3J')
                _refresh()
                if _is_slash():
                    _show()
                continue

            # ── Escape — clear the line, then quit ────────────────────────
            # Two presses so leaving is deliberate: the first press has
            # something to do (drop what is typed) whenever the line is not
            # empty, and only an Esc on an already-empty line arms the exit.
            if ch == '\x1b':
                if buffer:
                    buffer.clear()
                    selected = None
                    escape_armed = False
                    _hide()
                    _refresh()
                    continue
                if escape_armed:
                    _hide()
                    sys.stdout.write('\n')
                    sys.stdout.flush()
                    raise KeyboardInterrupt
                escape_armed = True
                _hide()
                _refresh()
                sys.stdout.write(f'\033[1B\r\033[2K  {DIM}press Esc again to '
                                 f'exit{RESET}\033[1A\r' + _rows()[-1])
                sys.stdout.flush()
                continue

            # ── Printable — any script ────────────────────────────────────
            # `ch.isprintable()` is true for 'a', 'П', 'ї', 'é', '日' and an
            # emoji, and false for control characters. The old test was
            # `32 <= ch[0] < 127` against a *byte*, which silently discarded
            # every Cyrillic keystroke — you could not type Ukrainian or
            # Russian into the prompt at all.
            if ch.isprintable():
                # Anything already queued behind this character arrived with
                # it. Taking the whole burst at once is also what keeps a big
                # paste from redrawing the screen once per character.
                if msvcrt.kbhit():
                    _absorb_burst(ch)
                else:
                    buffer.insert(pos, ch)
                    pos += 1
                selected = None
                _refresh()
                if _is_slash():
                    _show()
                else:
                    _hide()
                continue

            # ── Everything else (remaining control characters) ────────────
            continue

    except KeyboardInterrupt:
        raise


# ---------------------------------------------------------------------------
# Conversation history display (used when continuing a session)
# ---------------------------------------------------------------------------

def _format_block_content(content, max_len: int = 200) -> str:
    """Format a message content (str or list of blocks) into a single string."""
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        texts = []
        for block in content:
            if isinstance(block, dict):
                t = block.get("type", "")
                if t == "text":
                    texts.append(block.get("text", ""))
                elif t == "tool_use":
                    name = block.get("name", "?")
                    # ensure_ascii=False: history used to replay Cyrillic
                    # arguments as \uXXXX escapes, cut mid-escape at 80.
                    from text_display import shorten
                    inp = shorten(
                        json.dumps(block.get("input", {}), ensure_ascii=False), 80)
                    texts.append(f"[tool_use: {name}({inp})]")
                elif t == "tool_result":
                    rc = str(block.get("content", ""))[:120]
                    texts.append(f"[tool_result: {rc}]")
            elif hasattr(block, "type"):
                t = getattr(block, "type", "")
                if t == "text":
                    texts.append(getattr(block, "text", ""))
                elif t == "tool_use":
                    texts.append(f"[tool_use: {getattr(block, 'name', '?')}]")
                elif t == "tool_result":
                    texts.append("[tool_result]")
        return " ".join(texts)
    return str(content)


def _print_conversation_history(messages: list) -> None:
    """Pretty-print a full conversation history to the terminal.

    Used when continuing a session so the user can see everything that was
    said before typing a new message.
    """
    for i, m in enumerate(messages):
        role = m.get("role", "?")
        content = m.get("content", "")
        text = _format_block_content(content)

        if role == "user":
            icon = f'{GREEN}◆{RESET}'
            label = f'{GREEN}{BOLD}You{RESET}'
        elif role == "assistant":
            icon = f'{MAGENTA}▌{RESET}'
            label = f'{MAGENTA}{BOLD}TOMAS{RESET}'
        else:
            icon = f'{DIM}·{RESET}'
            label = f'{DIM}{role}{RESET}'

        # Truncate very long messages for display, but keep a reasonable preview
        display = text.replace('\n', '\n      ')
        if len(display) > 500:
            display = display[:500] + f' {DIM}…({len(text)} chars total){RESET}'

        print(f'  {icon} {label}')
        print(f'      {display}')
        print()


# ---------------------------------------------------------------------------
# REPL
# ---------------------------------------------------------------------------

def main() -> int:
    global mcp_manager, _current_context_window, CONTEXT_WINDOW, CONTINUE_SESSION_ID
    global _synthetic_replies

    # ── This session's accounting starts here, not at import time ──
    reset_session_state()

    # ── Bring the environment into line with the active provider ──
    # ~/.tomas/.env lands in os.environ at import and providers.json decides
    # who is active; nothing reconciled the two on the way in, so a session
    # could run with a *previous* provider's endpoint still live while every
    # menu showed the current one. Do it before the banner and the context
    # window, both of which read the result.
    try:
        import provider_manager
        provider_manager.sync_env_to_active()
    except Exception:
        pass    # a config we cannot read must not stop the agent from starting

    # ── Auto-start Zen proxy if needed ──
    _ensure_zen_proxy()

    # ── Fetch real context window for the current model ──
    _refresh_context_window()

    # ── Startup banner ──
    print()
    print(f'  {CYAN}{BOLD}████████╗ ██████╗ ███╗   ███╗ █████╗ ███████╗{RESET}')
    print(f'  {CYAN}{BOLD}╚══██╔══╝██╔═══██╗████╗ ████║██╔══██╗██╔════╝{RESET}')
    print(f'  {CYAN}{BOLD}   ██║   ██║   ██║██╔████╔██║███████║███████╗{RESET}')
    print(f'  {CYAN}{BOLD}   ██║   ██║   ██║██║╚██╔╝██║██╔══██║╚════██║{RESET}')
    print(f'  {CYAN}{BOLD}   ██║   ╚██████╔╝██║ ╚═╝ ██║██║  ██║███████║{RESET}')
    print(f'  {CYAN}{BOLD}   ╚═╝    ╚═════╝ ╚═╝     ╚═╝╚═╝  ╚═╝╚══════╝{RESET}')
    print()
    from version import VERSION, LAST_UPDATED
    print(f'  {DIM}Version:{RESET}  {CYAN}{VERSION}{RESET} '
          f'{DIM}(updated {LAST_UPDATED}){RESET}')
    print(f'  {DIM}Model:{RESET}    {CYAN}{_get_model()}{RESET}')
    cw_str = f'{_current_context_window:,}' if _current_context_window else '?'
    print(f'  {DIM}Context:{RESET}  {CYAN}{cw_str} tokens{RESET}')
    print(f'  {DIM}Project:{RESET}  {BLUE}{PROJECT_DIR.name}{RESET}')
    mode_name = current_mode_name()
    if mode_name == "bypass":
        mode_display = f'{RED}{BOLD}BYPASS ⇥{RESET}'
    elif mode_name == "yolo":
        mode_display = f'{RED}{BOLD}YOLO ⚡{RESET}'
    else:
        mode_display = f'{mode_color(mode_name)}{mode_name}{RESET}'
    print(f'  {DIM}Mode:{RESET}      {mode_display}')
    if mode_name == "bypass":
        print(f'  {DIM}           all tools auto-approved · never asks to '
              f'continue · up to '
              f'{MAX_TOOL_CALLS_PER_TURN * (MAX_AUTO_CONTINUATIONS + 1)} '
              f'tool calls per turn{RESET}')
    print()

    # ── Initialize MCP connections (single pass, no pre-test) ──
    # The work is in `init_mcp`; everything from here to the end of the block
    # is this front end deciding what to say about it.
    mcp_summary = init_mcp()
    if mcp_summary["error"]:
        print(f'  {RED}⚠{RESET}  MCP initialization failed: {mcp_summary["error"]}')
    try:
        disabled = set(mcp_summary["disabled"])
        connected_servers = set(mcp_summary["servers"])
        total_tools = mcp_summary["tools"]

        # One status line. Six red error lines before the user has typed
        # anything were mostly "this optional server has no credentials",
        # which is a fact, not a failure — and not one they can act on here.
        # The split into "needs credentials" and "broken" is what makes that
        # one line honest rather than merely short.
        needs_auth, broken = _classify_mcp_failures(mcp_summary["failed"])
        if connected_servers:
            line = (f'  {GREEN}✓{RESET}  {BOLD}MCP:{RESET} '
                    f'{len(connected_servers)} connected {DIM}·{RESET} {total_tools} tools')
        else:
            line = f'  {DIM}MCP:{RESET} none connected'
        extras = []
        if disabled:
            extras.append(f'{len(disabled)} disabled')
        if needs_auth:
            extras.append(f'{len(needs_auth)} need credentials')
        if broken:
            extras.append(f'{len(broken)} unavailable')
        if extras:
            line += f' {DIM}·{RESET} {DIM}' + f'{DIM} · {RESET}{DIM}'.join(extras) + RESET
        print(line)
        if needs_auth or broken:
            names = sorted(needs_auth) + sorted(broken)
            shown = ', '.join(names[:4]) + ('…' if len(names) > 4 else '')
            print(f'     {DIM}{shown} — MCP management is in the TUI '
                  f'({RESET}agent_cli.py{RESET}{DIM}); {RESET}{CYAN}/mcp-resources{RESET} '
                  f'{DIM}lists published resources{RESET}')

        # Names were resolved inside init_mcp. Which tools are *sent* is decided
        # per turn by select_tools(), against what the user is actually asking.
        if mcp_summary["renamed"]:
            print(f'  {YELLOW}⚠{RESET}  Renamed {mcp_summary["renamed"]} MCP '
                  f'tool(s) to avoid built-in name conflicts')
        if mcp_summary["dropped"]:
            print(f'  {CYAN}◈{RESET}  {mcp_summary["tools"]} MCP tools available, '
                  f'{mcp_summary["budget"]} fit per turn {DIM}(selected by '
                  f'relevance to each message; the model is told what is held '
                  f'back){RESET}')
    except Exception as e:
        print(f'  {RED}⚠{RESET}  Could not summarise MCP state: {e}')

    # ── Initialize self-improving system ──
    self_improve.init()
    init_learning()
    # One-shot: annotate sessions saved before completeness was recorded, so
    # an abandoned run is distinguishable from a finished one.
    try:
        from session_manager import backfill_completeness
        marked = backfill_completeness()
        if marked:
            print(f'  {YELLOW}⚠{RESET}  Marked {len(marked)} incomplete session(s) '
                  f'{DIM}(saved with turns that produced no reply){RESET}')
    except Exception:
        pass
    print(f'  {GREEN}✦{RESET}  Self-improving system active — learning from interactions')

    # One line, in the first five sessions only, and then never again — see
    # onboarding.py. A fresh install genuinely knows nothing about the user
    # (the fact store starts empty and reflection is a manual command), so the
    # offer is worth making; an assistant that keeps making it is one whose
    # notices stop being read at all.
    try:
        import onboarding
        onboarding.note_session_start()
        if onboarding.should_offer():
            onboarding.note_offered()
            print(f'  {MAGENTA}🧭{RESET}  {onboarding.offer_text()}')
    except Exception:
        pass

    print(f'  {DIM}─── Type {RESET}{BOLD}quit{RESET}{DIM} or {RESET}{BOLD}exit{RESET}{DIM} to leave · {RESET}{BOLD}/help{RESET}{DIM} for commands · {RESET}{BOLD}Esc Esc{RESET}{DIM} also exits ───{RESET}')
    print(f'  {DIM}    {RESET}{BOLD}Shift+Enter{RESET}{DIM} newline · paste keeps its line breaks · '
          f'{RESET}{BOLD}Ctrl+G{RESET}{DIM} attach clipboard image{RESET}')
    print(f'  {DIM}    {RESET}{BOLD}Ctrl+Y{RESET}{DIM} copy the line ({RESET}{BOLD}↑{RESET}{DIM} first for an '
          f'earlier one) · {RESET}{BOLD}Ctrl+Z{RESET}{DIM} undo a clear or paste{RESET}')
    print()

    messages: list = []

    # ── Continue previous session if requested ──
    if CONTINUE_SESSION_ID:
        try:
            from session_manager import continue_session, load_session
            loaded = continue_session(CONTINUE_SESSION_ID)
            if loaded:
                messages = loaded
                print(f'  {GREEN}✓{RESET}  Continuing session {CYAN}{CONTINUE_SESSION_ID[:16]}{RESET} ({len(loaded)} messages)')
                print(f'  {DIM}{"─" * 50}{RESET}')
                _print_conversation_history(loaded)
                print(f'  {DIM}{"─" * 50}{RESET}')
                print(f'  {DIM}Type your next message to continue the conversation.{RESET}')
                print()
        except Exception as e:
            print(f'  {YELLOW}⚠{RESET}  Could not load session: {e}')

    # ── Prefill, when asked for ──
    # Only into a genuinely fresh conversation: a continued session already
    # has the history this exists to simulate, and prepending filler to it
    # would push the user's real first turn further from the model *and*
    # re-add the filler on every subsequent continue.
    if not messages and features().enabled("prefill_context"):
        wanted = features().choice("prefill_tokens")
        primer = core_features.prefill_messages(wanted)
        if primer:
            messages.extend(primer)
            # The primer's acknowledgement is scaffolding, not a reply — see
            # `_synthetic_replies`. Without this the every-3rd-reply cap fires
            # a turn early for the whole session.
            _synthetic_replies = sum(1 for m in primer
                                     if m.get("role") == "assistant")
            share = ""
            window = CONTEXT_WINDOW or DEFAULT_CONTEXT_WINDOW
            if window:
                share = f' ({wanted / window:.0%} of the {window:,} window)'
            print(f'  {CYAN}◈{RESET}  Context prefilled: '
                  f'{CYAN}~{wanted:,} tokens{RESET}{DIM}{share} of history '
                  f'before your first message.{RESET}')
            print(f'  {DIM}Size and on/off: Settings → Prefill context.{RESET}')
            print()

    try:
        while True:
            try:
                user_input = read_input_with_suggestions(f'  {CYAN}{BOLD}TOMAS{RESET}{CYAN} »{RESET} ').strip()
            except EOFError:
                break
            if user_input.lower() in ("quit", "exit"):
                break
            if not user_input:
                continue

            # ── Record for self-improvement ──
            self_improve.record_user_message(user_input)

            # ── `#` adds a rule ──
            # The shortest path from "I want this to always happen" to a
            # stored rule. A slash command is one keystroke longer and one
            # decision harder: you have to remember whether it is /rules,
            # /remember or /note. `#` is the whole gesture, and it routes to
            # the same two stores /rules manages.
            if user_input.startswith('#') and user_input[1:].strip():
                print(f'  {MAGENTA}{BOLD}▌ TOMAS{RESET}')
                print(f'  {_rules_add(user_input[1:].strip())}')
                print()
                continue

            # ── Slash command handling ──
            if user_input.startswith('/'):
                result = handle_slash_command(user_input[1:], messages)
                if result == "__exit__":
                    break
                if result == "__continue__":
                    continue
                if result is not None:
                    # Non-None result from a processed slash command
                    print(f'  {MAGENTA}{BOLD}▌ TOMAS{RESET}')
                    print(f'  {result}')
                    print()
                    continue
                # None result: skill was loaded; fall through to agent loop
                # (skill content already injected by handle_slash_command)
            else:
                # Regular user message, plus any images it names.
                messages.append({"role": "user",
                                 "content": build_user_content(user_input)})
            # Release old tool results before considering compaction: pruning
            # is cheap and reversible (the model can re-read a file), while
            # compaction spends a model call and rewrites the whole transcript.
            # Doing it in this order means a session often no longer needs to
            # compact at all.
            # Prune, build the system prompt, then compact exactly once
            # against its real size — see _prepare_turn_context.
            messages, system_prompt = _prepare_turn_context(messages, user_input)
            # An auto-triggered skill was, until now, invisible: its body
            # went straight into the system prompt tail with nothing printed,
            # so there was no way to tell from the chat whether one had fired
            # at all. Same notice style as the explicit `/skill` path.
            # "[auto-applied]" was a promise the retrieval could not keep: it
            # is a substring match, and it fired on a question *about* a
            # document. The body is offered to the model with an instruction
            # to decide (see `skills_manager._GATE`), so the badge says what
            # actually happened — retrieved, not applied.
            for _triggered_skill in match_skills(user_input):
                print(f'  {YELLOW}⚡ Skill matched:{RESET} '
                      f'{BOLD}{_triggered_skill["name"]}{RESET} '
                      f'{DIM}[offered — the model decides if it fits]{RESET}')
            result = agent_loop(system_prompt, messages)
            # agent_loop owns both printing and transcript recording: the
            # streaming path prints tokens as they arrive, the non-streaming
            # path prints on return. Only error strings (which are returned
            # without being printed or recorded) need handling here.
            if result and result.startswith("I'm sorry"):
                print(f'  {MAGENTA}{BOLD}▌ TOMAS{RESET}')
                print(f'  {result}')
            # ── Token usage info ──
            # Behind `advanced_diagnostics`: it is per-turn accounting, and
            # accounting under every answer is the line that made the chat
            # read as a machine report rather than a conversation. Nothing is
            # lost by hiding it — `/status` asks the same question on demand,
            # and the session file records all of it either way.
            t = _last_turn_usage
            s = _session_tokens
            if s["calls"] > 0 and features().enabled("advanced_diagnostics"):
                pct = (t["input"] + t["output"]) / CONTEXT_WINDOW * 100
                elapsed = _fmt_duration(_turn_timings[-1]) if _turn_timings else "?"
                # Only shown when the provider actually reports it: a hardcoded
                # "0% cached" on an endpoint that simply does not say would read
                # as a cache that stopped working.
                cached = ""
                if s.get("cached_input"):
                    share = s["cached_input"] / max(1, s["input"]) * 100
                    cached = f'  ·  {s["cached_input"]:,} cached ({share:.0f}%)'
                # The bill nobody was shown: streamed calls re-issued
                # non-streamed to get their tool blocks. Named "duplicate"
                # rather than folded into the total because it is the number
                # that argues for removing the second call. `_stream_can_serve`
                # already serves everything recoverable, so `would_have_served`
                # is expected to read 0 now — the reason breakdown (why the
                # remaining ones fell through) is the useful part, not a
                # "recoverable" count that would always say zero.
                if s.get("duplicate_input"):
                    reasons = sorted(
                        (k[len("duplicate_reason_"):], v)
                        for k, v in s.items()
                        if k.startswith("duplicate_reason_") and v)
                    breakdown = ", ".join(f'{k}×{v}' for k, v in reasons)
                    cached += (f'  ·  {s["duplicate_input"]:,} duplicate'
                               f' ({s.get("duplicate_calls", 0)} calls'
                               + (f': {breakdown}' if breakdown else '') + ')')
                print(f'  {DIM}┄  {t["input"]:,} in  {t["output"]:,} out  ·  total: {s["input"]:,} in  {s["output"]:,} out{cached}  ·  {pct:.1f}% of {CONTEXT_WINDOW:,} ctx  ·  {elapsed}{RESET}')
                # How the turn ended, in the model's own vocabulary. This is
                # the "reached the limit" question: `max_tokens` says the
                # reply was cut off at the output ceiling, `tool_use` that it
                # stopped to run something, `end_turn` that it was simply
                # done. The core has recorded it on every turn since
                # `last_stop_reason` was added and nothing displayed it, so a
                # turn that stopped short looked identical to one that
                # finished.
                d = _last_turn_diag
                ending = d["stop_reason"] or "—"
                calls = (f'  ·  {d["tool_calls"]} tool call'
                         f'{"" if d["tool_calls"] == 1 else "s"}'
                         if d["tool_calls"] else '')
                # Only when it did not already surface as an error on screen:
                # `ErrorOccurred` prints its own detail line above, and
                # repeating it here would say the same thing twice in four
                # lines.
                why = (f'  ·  {d["error"]}'
                       if d["error"] and not d["error"].startswith("truncated")
                       else '')
                print(f'  {DIM}┄  ended: {ending}{calls}{why}{RESET}')
            # ── Self-improvement analysis, after the reply so it never adds latency ──
            try:
                self_improve.maybe_analyze_after_turn()
            except Exception:
                pass
            print()
    except KeyboardInterrupt:
        print()
        print(f'  {DIM}Bye!{RESET}')
    finally:
        # ── Save session on exit ──
        if messages:
            try:
                sid = save_session(
                    messages,
                    model=_get_model(),
                    token_usage=dict(_session_tokens),
                    session_id=CONTINUE_SESSION_ID,
                )
                print(f'  {DIM}💾 Session saved: {sid}{RESET}')
            except Exception as e:
                print(f'  {DIM}⚠  Session save skipped: {e}{RESET}')
            CONTINUE_SESSION_ID = None
            # Reflection calls the model over the whole transcript, which can
            # take several seconds on a long session. Shown as a spinner so
            # the wait isn't mistaken for a hang right after "Session saved"
            # already printed above.
            thinking = Thinking(label="reflecting on this session")
            thinking.start()
            try:
                reflect_on_session_end(messages)
            finally:
                thinking.stop()
        # ── Clean up MCP connections ──
        if mcp_manager:
            mcp_manager.disconnect_all()
        # ── Release local models this session loaded ──
        _release_local_models()
    return 0

if __name__ == "__main__":
    sys.exit(main())