"""pptx2txt.py — convert any .pptx to a TXT file while preserving as much
formatting detail as possible.

Strategy
--------
Uses ``python-pptx`` for the structural walk (slides, shapes, text frames,
tables, notes) and reaches into the underlying ``lxml`` tree for the
formatting attributes that ``python-pptx`` does not expose
(underline, strikethrough, character color).

What it preserves
-----------------
* Slide boundaries — ``===== Slide N (<layout>) =====`` markers
* Slide titles — extracted from the title placeholder, labeled
* Hidden slides — ``(hidden)`` flag in the slide header
* Body shapes — text from every shape with a text frame
* Tables — rendered as tab-separated rows
* Speaker notes — emitted as a ``----- Notes -----`` block per slide
* Inline formatting — bold / italic / underline / mono with inline markers
* Font name + size — optional per-line annotation (``--annotate-fonts``)
* Embedded images — optionally extracted (``--extract-images``)
* Grouped shapes — traversed recursively
* Full structured data — optional sidecar ``.json`` (``--json``)

Limitations
-----------
* Charts and SmartArt show up as ``[chart: ...]`` / ``[smart-art: ...]``
  placeholders; their internal data is not exported.
* Animations, transitions, and slide-master overrides are not preserved
  beyond the per-slide text they show.
* Text inside placeholders that are not yet filled in is silently skipped.

Usage
-----
    python pptx2txt.py file.pptx
    python pptx2txt.py file.pptx -o out.txt
    python pptx2txt.py file.pptx --annotate-fonts
    python pptx2txt.py file.pptx --extract-images img/
    python pptx2txt.py file.pptx --json
"""
from __future__ import annotations

import argparse
import json
import re
import sys
from pathlib import Path
from typing import Any

try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    from pptx.util import Emu
except ImportError:
    sys.stderr.write(
        "ERROR: python-pptx is required. Install it with:\n"
        "    python -m pip install python-pptx\n"
    )
    sys.exit(1)

try:
    from lxml import etree  # noqa: F401  (ensure lxml present)
except ImportError:
    sys.stderr.write("ERROR: lxml is required (python-pptx depends on it).\n")
    sys.exit(1)


# XML namespaces used in the OOXML spec
_A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
_P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"


# -- Inline formatting helpers (mirrors docx2txt.py) -------------------------

def _bool(elem) -> bool:
    if elem is None:
        return False
    val = elem.get("val")
    if val is None:
        return True
    return val.lower() not in ("false", "0", "off")


def run_formatting(run) -> dict[str, Any]:
    """Return a dict of formatting flags and properties for a run."""
    rPr = run._r.find(f"{{{_A_NS}}}rPr")
    flags: dict[str, Any] = {
        "bold": bool(run.font.bold) if run.font.bold is not None else False,
        "italic": bool(run.font.italic) if run.font.italic is not None else False,
        "underline": False,
        "strike": False,
        "color": None,
        "font": None,
        "size_pt": None,
    }

    if run.font.underline is True or (isinstance(run.font.underline, str)
                                      and run.font.underline):
        flags["underline"] = True
    elif rPr is not None:
        u = rPr.find(f"{{{_A_NS}}}u")
        flags["underline"] = _bool(u)

    if rPr is not None:
        strike = rPr.find(f"{{{_A_NS}}}strike")
        if strike is not None and strike.get("noStrike") != "1":
            flags["strike"] = True
        solidFill = rPr.find(f"{{{_A_NS}}}solidFill")
        if solidFill is not None:
            srgb = solidFill.find(f"{{{_A_NS}}}srgbClr")
            if srgb is not None and srgb.get("val") not in (None, "auto"):
                flags["color"] = "#" + srgb.get("val").upper()

    if run.font.name:
        flags["font"] = run.font.name
    if run.font.size is not None:
        flags["size_pt"] = float(run.font.size.pt)

    return flags


def is_mono(flags: dict[str, Any]) -> bool:
    name = (flags.get("font") or "").lower()
    return any(token in name for token in (
        "courier", "consolas", "monaco", "menlo", "roboto mono", "fira",
        "source code", "cascadia", "jetbrains", "consola",
    ))


def wrap_inline(text: str, flags: dict[str, Any]) -> str:
    if not text:
        return text
    if is_mono(flags):
        if "\n" in text:
            return text
        return f"`{text}`"
    if flags.get("bold") and flags.get("italic"):
        text = f"***{text}***"
    elif flags.get("bold"):
        text = f"**{text}**"
    elif flags.get("italic"):
        text = f"*{text}*"
    if flags.get("underline"):
        text = f"_{text}_"
    if flags.get("strike"):
        text = f"~~{text}~~"
    return text


def format_run_label(flags: dict[str, Any]) -> str:
    parts = []
    if flags.get("font"):
        parts.append(flags["font"])
    if flags.get("size_pt") is not None:
        parts.append(f"{flags['size_pt']:.1f}pt")
    style = []
    if flags.get("bold"):
        style.append("B")
    if flags.get("italic"):
        style.append("I")
    if flags.get("underline"):
        style.append("U")
    if flags.get("strike"):
        style.append("S")
    if is_mono(flags):
        style.append("M")
    if style:
        parts.append("/".join(style))
    return " ".join(parts) if parts else "?"


# -- Paragraph / bullet handling ---------------------------------------------

def bullet_marker(paragraph) -> str:
    """Return a bullet character (or numbered prefix) for a paragraph, or ''."""
    pPr = paragraph._pPr
    if pPr is None:
        return ""
    # python-pptx surfaces buChar / buAutoNum via paragraph attributes
    try:
        if paragraph.level is not None and paragraph.level > 0:
            indent = "    " * paragraph.level
        else:
            indent = ""
    except Exception:
        indent = ""

    # buChar: literal bullet glyph (•, –, etc.)
    buChar = pPr.find(f"{{{_A_NS}}}buChar")
    if buChar is not None and buChar.get("char"):
        return indent + buChar.get("char") + " "

    # buAutoNum: numbered/bulleted
    buAuto = pPr.find(f"{{{_A_NS}}}buAutoNum")
    if buAuto is not None:
        kind = buAuto.get("type", "arabicPeriod")
        # We can't know the actual sequence number without walking the
        # slide's lists, so emit a generic marker.
        prefix = {
            "arabicPeriod": "#.",
            "arabicPlain": "#",
            "romanLcPeriod": "i.",
            "romanUcPeriod": "I.",
            "alphaLcPeriod": "a.",
            "alphaUcPeriod": "A.",
        }.get(kind, "#.")
        return indent + prefix + " "

    # buNone or no bullet at all -> not a list item
    return ""


def render_paragraph(paragraph, *, annotate: bool) -> str:
    pieces: list[str] = []
    for run in paragraph.runs:
        if not run.text:
            continue
        flags = run_formatting(run)
        pieces.append(wrap_inline(run.text, flags))
    line = "".join(pieces).rstrip()
    if not line.strip():
        return ""

    marker = bullet_marker(paragraph)
    if marker:
        line = marker + line
    if annotate:
        for run in paragraph.runs:
            if run.text.strip():
                line = f"[{format_run_label(run_formatting(run))}] {line}"
                break
    return line


# -- Shape / slide traversal -------------------------------------------------

def shape_placeholder_role(shape) -> str | None:
    """Return 'title' / 'body' / 'ctrTitle' / 'subTitle' for placeholders."""
    if not shape.is_placeholder:
        return None
    fmt = shape.placeholder_format
    if fmt is None:
        return None
    # idx 0 is the title; type carries the semantic role
    try:
        if fmt.idx == 0 or fmt.type in (13, 14, 15):  # 13=TITLE, 14=CTR_TITLE, ...
            return "title"
    except Exception:
        pass
    return "body"


def shape_position(shape) -> str:
    """Short `(x, y, w x h)` label in points for a shape."""
    try:
        left = Emu(shape.left).pt if shape.left is not None else 0
        top = Emu(shape.top).pt if shape.top is not None else 0
        width = Emu(shape.width).pt if shape.width is not None else 0
        height = Emu(shape.height).pt if shape.height is not None else 0
        return f"(x={left:.0f} y={top:.0f} {width:.0f}x{height:.0f} pt)"
    except Exception:
        return ""


def render_shape(shape, *, annotate: bool, depth: int = 0) -> list[str]:
    """Render any shape to a list of text lines. Recurses into groups."""
    indent = "  " * depth
    out: list[str] = []

    # Group: recurse
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        out.append(f"{indent}[group] {shape.name}")
        for child in shape.shapes:
            out.extend(render_shape(child, annotate=annotate, depth=depth + 1))
        return out

    # Placeholder without a text frame (image placeholder etc.)
    if not shape.has_text_frame and not shape.has_table:
        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
            out.append(f"{indent}[image] {shape.name} {shape_position(shape)}")
        elif shape.shape_type == MSO_SHAPE_TYPE.CHART:
            out.append(f"{indent}[chart] {shape.name} {shape_position(shape)}")
        elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
            out.append(f"{indent}[shape] {shape.name} {shape_position(shape)}")
        else:
            out.append(f"{indent}[{shape.shape_type}] {shape.name} {shape_position(shape)}")
        return out

    role = shape_placeholder_role(shape)

    if shape.has_table:
        out.append(f"{indent}[table] {shape.name}")
        out.extend(render_table(shape.table, indent=indent + "  "))
        return out

    if shape.has_text_frame:
        for p in shape.text_frame.paragraphs:
            rendered = render_paragraph(p, annotate=annotate)
            if rendered:
                prefix = ""
                if role == "title" and p is shape.text_frame.paragraphs[0]:
                    prefix = "TITLE: "
                out.append(indent + prefix + rendered)
        return out

    return out


def render_table(table, *, indent: str = "") -> list[str]:
    """Render a table as tab-separated rows under a [table] header.

    The caller (``render_shape``) is responsible for adding the shape name;
    this function only emits the rows + closing tag.
    """
    rows: list[list[str]] = []
    for row in table.rows:
        cells: list[str] = []
        for cell in row.cells:
            cell_parts: list[str] = []
            for p in cell.text_frame.paragraphs:
                rendered = render_paragraph(p, annotate=False)
                if rendered:
                    cell_parts.append(rendered)
            cells.append("\n".join(cell_parts).replace("\t", "    "))
        rows.append(cells)
    if not rows:
        return []
    width = max(len(r) for r in rows)
    out: list[str] = []
    for r in rows:
        r = r + [""] * (width - len(r))
        out.append(indent + "\t".join(r))
    out.append(f"{indent}[/table]")
    return out


def render_notes(slide) -> list[str]:
    """Render the speaker notes of a slide (if any)."""
    if not slide.has_notes_slide:
        return []
    notes = slide.notes_slide.notes_text_frame
    parts: list[str] = []
    for p in notes.paragraphs:
        rendered = render_paragraph(p, annotate=False)
        if rendered:
            parts.append(rendered)
    if not parts:
        return []
    return ["----- Notes -----"] + parts + ["----- /Notes -----"]


# -- Main extraction ---------------------------------------------------------

def slide_layout_name(slide) -> str:
    try:
        return slide.slide_layout.name or "(layout)"
    except Exception:
        return "(layout)"


def extract(prs: Presentation, *, annotate: bool) -> list[str]:
    out: list[str] = []
    out.append("===== Presentation =====")
    if prs.core_properties.title:
        out.append(f"Title: {prs.core_properties.title}")
    if prs.core_properties.author:
        out.append(f"Author: {prs.core_properties.author}")
    n = len(prs.slides)
    out.append(f"Slides: {n}")
    out.append("")

    for i, slide in enumerate(prs.slides, start=1):
        layout = slide_layout_name(slide)
        hidden = ""
        # python-pptx >=0.6.21 exposes slide.show (ShowPlaceholders enum)
        # but the underlying element is the most reliable check.
        try:
            show_attr = slide._element.get("show", "1")
            if show_attr == "0":
                hidden = " (hidden)"
        except Exception:
            pass

        out.append(f"===== Slide {i} — {layout}{hidden} =====")

        for shape in slide.shapes:
            out.extend(render_shape(shape, annotate=annotate))

        # Speaker notes
        out.extend(render_notes(slide))
        out.append("")  # blank between slides

    return out


# -- Image extraction --------------------------------------------------------

def extract_images(prs: Presentation, out_dir: Path) -> int:
    out_dir.mkdir(parents=True, exist_ok=True)
    saved = 0
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                try:
                    blob = shape.image.blob
                    ext = shape.image.ext or "png"
                except Exception:
                    continue
                path = out_dir / f"slide{slide_idx:03d}_img{saved + 1:03d}.{ext}"
                path.write_bytes(blob)
                saved += 1
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                for child in shape.shapes:
                    if child.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        try:
                            blob = child.image.blob
                            ext = child.image.ext or "png"
                        except Exception:
                            continue
                        path = out_dir / f"slide{slide_idx:03d}_img{saved + 1:03d}.{ext}"
                        path.write_bytes(blob)
                        saved += 1
    return saved


# -- Strip-formatting mode ---------------------------------------------------

_STRIP_PATTERNS = [
    (re.compile(r"\*\*(.+?)\*\*"), r"\1"),
    (re.compile(r"\*(.+?)\*"), r"\1"),
    (re.compile(r"_(.+?)_"), r"\1"),
    (re.compile(r"~~(.+?)~~"), r"\1"),
    (re.compile(r"`(.+?)`"), r"\1"),
]


def strip_formatting(text: str) -> str:
    for pat, repl in _STRIP_PATTERNS:
        text = pat.sub(repl, text)
    return text


# -- CLI ---------------------------------------------------------------------

def build_argparser() -> argparse.ArgumentParser:
    ap = argparse.ArgumentParser(
        description="Convert a .pptx to a formatting-preserving TXT file.",
        formatter_class=argparse.RawDescriptionHelpFormatter,
    )
    ap.add_argument("pptx", help="Input .pptx file")
    ap.add_argument("-o", "--output", help="Output TXT path (default: <pptx>.txt)")
    ap.add_argument("--no-format", action="store_true",
                    help="Strip **bold**/*italic*/`mono` markers, plain text only")
    ap.add_argument("--annotate-fonts", action="store_true",
                    help="Prefix each paragraph with its dominant font + size")
    ap.add_argument("--extract-images", metavar="DIR",
                    help="Extract embedded images into DIR")
    ap.add_argument("--json", action="store_true",
                    help="Also write a sidecar .json with full structured data")
    return ap


def main(argv: list[str] | None = None) -> int:
    args = build_argparser().parse_args(argv)
    in_path = Path(args.pptx)
    if not in_path.is_file():
        sys.stderr.write(f"ERROR: file not found: {in_path}\n")
        return 1

    out_path = Path(args.output) if args.output else in_path.with_suffix(".txt")
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(in_path)
    lines = extract(prs, annotate=args.annotate_fonts)
    text = "\n".join(lines)
    if args.no_format:
        text = strip_formatting(text)

    out_path.write_text(text, encoding="utf-8")
    n_slides = len(prs.slides)
    n_shapes = sum(len(s.shapes) for s in prs.slides)
    print(
        f"Wrote {out_path}  "
        f"({n_slides} slide{'s' if n_slides != 1 else ''}, "
        f"{n_shapes} shape{'s' if n_shapes != 1 else ''})"
    )

    if args.extract_images:
        n = extract_images(prs, Path(args.extract_images))
        print(f"Extracted {n} image{'s' if n != 1 else ''} -> {args.extract_images}")

    if args.json:
        json_path = out_path.with_suffix(".json")
        structured = {
            "source": str(in_path),
            "title": prs.core_properties.title,
            "author": prs.core_properties.author,
            "slides": [
                {
                    "layout": s.slide_layout.name,
                    "shapes": [
                        {"name": sh.name, "type": str(sh.shape_type),
                         "text": sh.text_frame.text if sh.has_text_frame else None}
                        for sh in s.shapes
                    ],
                    "notes": s.notes_slide.notes_text_frame.text if s.has_notes_slide else None,
                }
                for s in prs.slides
            ],
        }
        json_path.write_text(json.dumps(structured, ensure_ascii=False, indent=2),
                             encoding="utf-8")
        print(f"Wrote {json_path}  (full structured data)")

    return 0


if __name__ == "__main__":
    sys.exit(main())
